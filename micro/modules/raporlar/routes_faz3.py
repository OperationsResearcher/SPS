"""Faz 3 — Premium dosya ürünleri (stratejik yıllık, yatırımcı sunum, ESG, audit, bireysel karne batch)."""
from __future__ import annotations

import re as _re
from collections import defaultdict
from datetime import datetime, timedelta, timezone, date as _date

from flask import render_template, jsonify, request, current_app, send_file
from flask_login import login_required, current_user
from app.utils.decorators import require_module
from sqlalchemy import func, and_, or_, text, select

from platform_core import app_bp
from app.models import db


def _hk_safe_name() -> str:
    """Aktif kurumun dosya-güvenli adı (rapor dosya adlarında). Fonksiyon-kapsamı
    sorununu önler — eskiden yerel `_safe_filename` değişkeni başka fonksiyonlarda
    NameError veriyordu."""
    from app.models.core import Tenant
    tid = getattr(current_user, "tenant_id", None)
    t = db.session.get(Tenant, tid) if tid else None
    return _re.sub(r'[^\w\-]', '_', (t.name if t and t.name else "Kurum"))[:50]
from app.models.core import User, Strategy, SubStrategy, Tenant
from app.models.process import (
    Process, ProcessKpi, KpiData, IndividualPerformanceIndicator,
    ProcessActivity, ProcessSubStrategyLink, process_leaders,
)
from app.models.k_vektor import KVektorStrategyWeight
from app.models.plan_year import PlanYear, KpiYearConfig
from app.models.initiative import Initiative
from app.services.plan_year_service import get_active_plan_year_for_user, list_plan_years
from app.services.score_engine_service import compute_process_scores_internal

from .helpers import _tid_or_none, MUDA_MAX_PROCESSES, _ai_text
from flask_babel import gettext as _

# ═══════════════════════════════════════════════════════════════════════════
# FAZ 3 — PREMIUM ÜRÜNLER (indirilebilir PDF/PPTX/Excel/ZIP)
# ═══════════════════════════════════════════════════════════════════════════

def _pdf_helpers():
    """Reportlab PDF üretici yardımcılar."""
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Image,
    )
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    return {
        "SimpleDocTemplate": SimpleDocTemplate, "Paragraph": Paragraph,
        "Spacer": Spacer, "PageBreak": PageBreak, "Table": Table,
        "TableStyle": TableStyle, "A4": A4, "cm": cm, "colors": colors,
        "styles": getSampleStyleSheet(),
        "ParagraphStyle": ParagraphStyle,
        "TA_CENTER": TA_CENTER, "TA_JUSTIFY": TA_JUSTIFY,
    }


# ─── ST-16: Stratejik Yıllık Kitap PDF ─────────────────────────────────────

@app_bp.route("/reports/strategic-annual")
@login_required
@require_module("raporlar")
def raporlar_stratejik_yillik():
    return render_template("platform/raporlar/stratejik_yillik.html")


@app_bp.route("/reports/api/strategic-annual/preview")
@login_required
@require_module("raporlar")
def raporlar_api_stratejik_yillik_preview():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    tenant = db.session.get(Tenant, tid)
    active_py = get_active_plan_year_for_user(current_user)

    sections = [
        {"no": 1, "title": _("Kapak"), "content": _("Kurum logo + yıl + başlık")},
        {"no": 2, "title": _("Önsöz"), "content": _("AI özet — 1 sayfa giriş")},
        {"no": 3, "title": _("Yönetici Özeti"), "content": _("Yıl özetinde 3 cümle + 5 kritik metrik")},
        {"no": 4, "title": _("Kurumsal Kimlik"), "content": _("Vizyon, misyon, değerler, etik")},
        {"no": 5, "title": _("Stratejik Direkler"), "content": _("Ana stratejiler + K-Vektör ağırlık")},
        {"no": 6, "title": _("Süreç Mükemmelliği"), "content": _("Süreç sağlığı + CMMI olgunluk")},
        {"no": 7, "title": _("Initiative Portföyü"), "content": _("Yıl içi girişimler + bütçe")},
        {"no": 8, "title": _("Performans Göstergeleri"), "content": _("PG hedef-gerçek özet")},
        {"no": 9, "title": _("Risk Yönetimi"), "content": _("Risk heatmap + top 10")},
        {"no": 10, "title": _("Sürdürülebilirlik (ESG)"), "content": _("Carbon + sosyal + yönetişim")},
        {"no": 11, "title": _("İK ve Yetenek"), "content": _("Çalışan + departman + bireysel başarı")},
        {"no": 12, "title": _("Yıllar Arası Karşılaştırma"), "content": _("Önceki yıllarla evrim")},
        {"no": 13, "title": _("Önümüzdeki Dönem"), "content": _("Yol haritası + AI önerileri")},
        {"no": 14, "title": _("Ekler"), "content": _("Veri kaynağı, metodoloji, terimler")},
    ]
    return jsonify({"success": True, "preview": {
        "tenant_name": tenant.name if tenant else "—",
        "plan_year": active_py.year if active_py else None,
        "sections": sections,
        "estimated_pages": 35,
        "format": _("PDF (reportlab)"),
    }})


@app_bp.route("/reports/api/strategic-annual/generate", methods=["GET", "POST"])
@login_required
@require_module("raporlar")
def raporlar_api_stratejik_yillik_generate():
    from flask import send_file
    import io
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    h = _pdf_helpers()
    tenant = db.session.get(Tenant, tid)
    active_py = get_active_plan_year_for_user(current_user)
    year_label = active_py.year if active_py else _date.today().year
    from html import escape as _he
    import re as _re
    _raw_name = tenant.name if tenant else "Kurum"
    tname = _he(_raw_name)  # ReportLab Paragraph markup injection koruması
    _safe_filename = _re.sub(r'[^\w\-]', '_', _raw_name)[:50]

    # Veri toplama
    strat_q = Strategy.query.filter_by(tenant_id=tid, is_active=True)
    if active_py:
        strat_q = strat_q.filter_by(plan_year_id=active_py.id)
    strategies = strat_q.order_by(Strategy.code).all()
    weights = {w.strategy_id: w.weight_raw or 0 for w in
               KVektorStrategyWeight.query.filter_by(tenant_id=tid).all()}
    total_w = sum(weights.values()) or 1

    proc_count = Process.query.filter_by(tenant_id=tid, is_active=True,
                                          plan_year_id=active_py.id if active_py else None).count()
    kpi_count = ProcessKpi.query.join(Process).filter(
        Process.tenant_id == tid, Process.plan_year_id == (active_py.id if active_py else None),
        ProcessKpi.is_active.is_(True),
    ).count()
    initiatives = Initiative.query.filter_by(tenant_id=tid, is_active=True).all()
    user_count = User.query.filter_by(tenant_id=tid, is_active=True).count()
    meas_count = db.session.query(func.count(KpiData.id)).join(ProcessKpi).join(
        Process).filter(Process.tenant_id == tid).scalar() or 0

    # AI önsöz
    onsoz = _ai_text(
        prompt=f"{tname} {year_label} Stratejik Yıllık Kitabı için 4-5 cümlelik bir önsöz yaz. "
               f"Veri: {len(strategies)} strateji, {proc_count} süreç, {kpi_count} PG, "
               f"{len(initiatives)} initiative, {user_count} çalışan, {meas_count} ölçüm. "
               "Resmî, ilham verici ton kullan.",
        fallback=(f"{tname} olarak {year_label} yılında stratejik planlama yolculuğumuzda "
                  f"{len(strategies)} ana direk üzerinden {proc_count} süreç ve {kpi_count} "
                  f"performans göstergesi ile yönettik. {user_count} çalışanımızın katkısıyla "
                  f"{meas_count:,} ölçüm verisi sisteme aktarıldı. Bu yıllık, "
                  "yıl içinde elde ettiğimiz başarıları, karşılaştığımız sorunları ve "
                  "önümüzdeki dönem için yol haritamızı bir arada sunmaktadır."),
        tid=tid, endpoint="ai_stratejik_yillik_onsoz", max_tokens=350,
    )

    # PDF üret
    buf = io.BytesIO()
    doc = h["SimpleDocTemplate"](
        buf, pagesize=h["A4"],
        leftMargin=2*h["cm"], rightMargin=2*h["cm"],
        topMargin=2.5*h["cm"], bottomMargin=2.5*h["cm"],
        title=f"{tname} {year_label} Stratejik Yıllık", author="Kokpitim",
    )
    styles = h["styles"]
    title_style = h["ParagraphStyle"](
        "TitleBig", parent=styles["Title"], fontSize=32, leading=40,
        textColor=h["colors"].HexColor("#0f172a"), spaceAfter=20, alignment=h["TA_CENTER"],
    )
    h1 = h["ParagraphStyle"]("H1", parent=styles["Heading1"], fontSize=20,
        textColor=h["colors"].HexColor("#4f46e5"), spaceAfter=16, spaceBefore=12)
    h2 = h["ParagraphStyle"]("H2", parent=styles["Heading2"], fontSize=14,
        textColor=h["colors"].HexColor("#1e293b"), spaceAfter=10, spaceBefore=10)
    body = h["ParagraphStyle"]("Body", parent=styles["BodyText"], fontSize=11,
        leading=16, textColor=h["colors"].HexColor("#0f172a"),
        spaceAfter=8, alignment=h["TA_JUSTIFY"])
    small = h["ParagraphStyle"]("Small", parent=body, fontSize=9,
        textColor=h["colors"].HexColor("#64748b"))

    P = h["Paragraph"]
    elems = []

    # Kapak
    elems.append(h["Spacer"](1, 6*h["cm"]))
    elems.append(P(tname, title_style))
    elems.append(P(f"{year_label}", title_style))
    elems.append(h["Spacer"](1, 1*h["cm"]))
    elems.append(P(_("STRATEJİK YILLIK"), h["ParagraphStyle"](
        "Sub", parent=body, fontSize=16, alignment=h["TA_CENTER"],
        textColor=h["colors"].HexColor("#64748b"), spaceAfter=8)))
    elems.append(h["Spacer"](1, 3*h["cm"]))
    elems.append(P(f"{_('Hazırlık tarihi')}: {_date.today().strftime('%d %B %Y')}",
                    h["ParagraphStyle"]("date", parent=small, alignment=h["TA_CENTER"])))
    elems.append(h["PageBreak"]())

    # Önsöz
    elems.append(P(_("Önsöz"), h1))
    elems.append(P(onsoz, body))
    elems.append(h["PageBreak"]())

    # Yönetici özeti
    elems.append(P(_("Yönetici Özeti"), h1))
    elems.append(P(f"<b>{year_label}</b> {_('yılı için kurum geneli görünümü:')}", body))
    summary_rows = [
        [_("Strateji sayısı"), str(len(strategies))],
        [_("Aktif süreç"), str(proc_count)],
        [_("Performans göstergesi"), str(kpi_count)],
        [_("Stratejik girişim"), str(len(initiatives))],
        [_("Aktif çalışan"), str(user_count)],
        [_("Toplam ölçüm"), f"{meas_count:,}"],
    ]
    tbl = h["Table"](summary_rows, colWidths=[8*h["cm"], 5*h["cm"]])
    tbl.setStyle(h["TableStyle"]([
        ("BACKGROUND", (0, 0), (0, -1), h["colors"].HexColor("#f8fafc")),
        ("TEXTCOLOR", (1, 0), (1, -1), h["colors"].HexColor("#4f46e5")),
        ("FONTNAME", (1, 0), (1, -1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 11),
        ("BOX", (0, 0), (-1, -1), 1, h["colors"].HexColor("#e2e8f0")),
        ("INNERGRID", (0, 0), (-1, -1), 0.5, h["colors"].HexColor("#e2e8f0")),
        ("PADDING", (0, 0), (-1, -1), 10),
    ]))
    elems.append(tbl)
    elems.append(h["PageBreak"]())

    # Kurumsal Kimlik
    elems.append(P(_("Kurumsal Kimlik"), h1))
    if tenant:
        if tenant.purpose:
            elems.append(P(_("Amaç"), h2)); elems.append(P(tenant.purpose, body))
        if tenant.vision:
            elems.append(P(_("Vizyon"), h2)); elems.append(P(tenant.vision, body))
        if tenant.core_values:
            elems.append(P(_("Değerler"), h2)); elems.append(P(tenant.core_values, body))
        if tenant.code_of_ethics:
            elems.append(P(_("Etik Kurallar"), h2)); elems.append(P(tenant.code_of_ethics, body))
    else:
        elems.append(P(_("Kurumsal kimlik bilgisi tanımlanmamış."), small))
    elems.append(h["PageBreak"]())

    # Stratejik Direkler
    elems.append(P(f"{_('Stratejik Direkler')} — {len(strategies)} {_('Ana Strateji')}", h1))
    if strategies:
        rows = [[_("Kod"), _("Ad"), _("K-Vektör %")]]
        for s in strategies:
            w_pct = round(weights.get(s.id, 0) / total_w * 100, 1)
            rows.append([s.code or "—", (s.title or "")[:60], f"%{w_pct}"])
        tbl = h["Table"](rows, colWidths=[2*h["cm"], 11*h["cm"], 3*h["cm"]])
        tbl.setStyle(h["TableStyle"]([
            ("BACKGROUND", (0, 0), (-1, 0), h["colors"].HexColor("#4f46e5")),
            ("TEXTCOLOR", (0, 0), (-1, 0), h["colors"].white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 10),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1),
             [h["colors"].white, h["colors"].HexColor("#f8fafc")]),
            ("BOX", (0, 0), (-1, -1), 1, h["colors"].HexColor("#e2e8f0")),
            ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#e2e8f0")),
            ("PADDING", (0, 0), (-1, -1), 8),
        ]))
        elems.append(tbl)
    else:
        elems.append(P(_("Strateji tanımlanmamış."), small))
    elems.append(h["PageBreak"]())

    # Initiative Portföyü
    elems.append(P(f"{_('Initiative Portföyü')} — {len(initiatives)} {_('Girişim')}", h1))
    if initiatives:
        rows = [[_("Kod"), _("Ad"), _("Durum"), _("İlerleme"), _("Bütçe (₺)")]]
        for i in initiatives[:30]:
            rows.append([
                i.code or "—", (i.name or "")[:40],
                i.status or "—", f"%{i.progress_pct or 0}",
                f"{float(i.budget_total or 0):,.0f}",
            ])
        tbl = h["Table"](rows, colWidths=[3*h["cm"], 6*h["cm"], 2.5*h["cm"], 2*h["cm"], 3*h["cm"]])
        tbl.setStyle(h["TableStyle"]([
            ("BACKGROUND", (0, 0), (-1, 0), h["colors"].HexColor("#16a34a")),
            ("TEXTCOLOR", (0, 0), (-1, 0), h["colors"].white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1),
             [h["colors"].white, h["colors"].HexColor("#f8fafc")]),
            ("BOX", (0, 0), (-1, -1), 1, h["colors"].HexColor("#e2e8f0")),
            ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#e2e8f0")),
            ("PADDING", (0, 0), (-1, -1), 6),
        ]))
        elems.append(tbl)
    else:
        elems.append(P(_("Initiative tanımlanmamış."), small))
    elems.append(h["PageBreak"]())

    # Önümüzdeki dönem
    elems.append(P(_("Önümüzdeki Dönem"), h1))
    next_year_text = _ai_text(
        prompt=f"{tname} için {year_label+1} yılı önerilen 3-4 stratejik öncelik yaz. "
               "Mevcut yıldaki verilere bakarak (özellikle düşük performanslı alanlar).",
        fallback=(f"{year_label+1} yılında stratejik plan döngüsünün doğal devamı olarak: "
                  "(1) mevcut stratejilerin değerlendirilmesi, "
                  "(2) yeni initiative tahsisleri, "
                  "(3) operasyonel mükemmellik odakları, "
                  "(4) ESG metriklerinin yatırımcı standartlarına çekilmesi öngörülmektedir."),
        tid=tid, endpoint="ai_stratejik_yillik_next", max_tokens=350,
    )
    elems.append(P(next_year_text, body))
    elems.append(h["PageBreak"]())

    # Kapanış
    elems.append(P(_("Teşekkürler"), h1))
    elems.append(P(f"{_('Bu yıllık,')} {tname} {_('ekibinin yıl boyunca verdiği emeğin yansımasıdır. Kokpitim platformu üzerinde sistemli olarak toplanan veriler ile üretilmiştir.')}", body))
    elems.append(h["Spacer"](1, 3*h["cm"]))
    elems.append(P(f"<b>{tname}</b> · {year_label} {_('Stratejik Yıllık')}",
                    h["ParagraphStyle"]("FootCenter", parent=body, alignment=h["TA_CENTER"])))
    elems.append(P(_("Kokpitim — Kurumsal Performans Yönetim Platformu"),
                    h["ParagraphStyle"]("FootSmall", parent=small, alignment=h["TA_CENTER"])))

    doc.build(elems)
    buf.seek(0)

    filename = f"{_safe_filename}_{year_label}_stratejik_yillik.pdf"
    return send_file(buf, mimetype="application/pdf",
                     as_attachment=True, download_name=filename)


# ─── AI-17: Yatırımcı Sunum PPTX ───────────────────────────────────────────

@app_bp.route("/reports/investor-presentation")
@login_required
@require_module("raporlar")
def raporlar_yatirimci_sunum():
    return render_template("platform/raporlar/yatirimci_sunum.html")


@app_bp.route("/reports/api/investor-presentation/preview")
@login_required
@require_module("raporlar")
def raporlar_api_yatirimci_sunum_preview():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    tenant = db.session.get(Tenant, tid)
    active_py = get_active_plan_year_for_user(current_user)

    slides = [
        {"no": 1, "title": _("Kapak"), "content": f"{tenant.name if tenant else 'Kurum'} — {_('Yatırımcı Sunumu')}"},
        {"no": 2, "title": _("Yönetici Özeti"), "content": _("AI 4-cümle özet")},
        {"no": 3, "title": _("Vizyon & Misyon"), "content": _("Stratejik kimlik")},
        {"no": 4, "title": _("Pazardaki Yerimiz"), "content": _("Sektör + rakip benchmark")},
        {"no": 5, "title": _("İş Modeli"), "content": _("Gelir kalemleri + büyüme stratejisi")},
        {"no": 6, "title": _("Stratejik Direkler"), "content": f"{Strategy.query.filter_by(tenant_id=tid, is_active=True).count()} {_('ana strateji')}"},
        {"no": 7, "title": _("Operasyonel Mükemmellik"), "content": _("Süreç olgunluk + verimlilik")},
        {"no": 8, "title": _("Finansal Performans"), "content": _("Bütçe + initiative ROI")},
        {"no": 9, "title": _("İnsan Kaynağı"), "content": _("Çalışan + yetenek geliştirme")},
        {"no": 10, "title": _("ESG & Sürdürülebilirlik"), "content": _("Çevre + sosyal + yönetişim")},
        {"no": 11, "title": _("Risk Yönetimi"), "content": _("Risk register + mitigation")},
        {"no": 12, "title": _("Initiative Portföyü"), "content": f"{Initiative.query.filter_by(tenant_id=tid, is_active=True).count()} {_('girişim')}"},
        {"no": 13, "title": _("Geçmiş Performans (7 yıl)"), "content": _("Yıllar arası trend")},
        {"no": 14, "title": _("3 Yıllık Projeksiyon"), "content": _("Senaryo bazlı tahmin")},
        {"no": 15, "title": _("Yatırım Talebi"), "content": _("Bütçe + use of funds")},
        {"no": 16, "title": _("Yönetim Ekibi"), "content": _("Lider ekip")},
        {"no": 17, "title": _("Soru-Cevap"), "content": _("Hazır cevaplar")},
        {"no": 18, "title": _("Teşekkürler"), "content": _("İletişim")},
    ]
    return jsonify({"success": True, "preview": {
        "tenant_name": tenant.name if tenant else "—",
        "plan_year": active_py.year if active_py else None,
        "slides": slides,
        "format": _("PPTX (python-pptx, 16:9)"),
        "slide_count": len(slides),
    }})


@app_bp.route("/reports/api/investor-presentation/generate", methods=["GET", "POST"])
@login_required
@require_module("raporlar")
def raporlar_api_yatirimci_sunum_generate():
    from flask import send_file
    import io
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return jsonify({"success": False, "message": "python-pptx yok"}), 500

    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    tenant = db.session.get(Tenant, tid)
    active_py = get_active_plan_year_for_user(current_user)
    year_label = str(active_py.year) if active_py else str(_date.today().year)
    tname = tenant.name if tenant else "Kurum"

    strategies = Strategy.query.filter_by(tenant_id=tid, is_active=True,
        plan_year_id=active_py.id if active_py else None).order_by(Strategy.code).all()
    inits = Initiative.query.filter_by(tenant_id=tid, is_active=True).all()
    proc_count = Process.query.filter_by(tenant_id=tid, is_active=True).count()
    user_count = User.query.filter_by(tenant_id=tid, is_active=True).count()
    meas_count = db.session.query(func.count(KpiData.id)).join(ProcessKpi).join(
        Process).filter(Process.tenant_id == tid).scalar() or 0
    total_budget = sum(float(i.budget_total or 0) for i in inits)

    # AI özet
    investor_summary = _ai_text(
        prompt=f"{tname} için yatırımcılara yönelik 4-cümle özet yaz. "
               f"Strateji: {len(strategies)}, Süreç: {proc_count}, "
               f"Çalışan: {user_count}, Initiative bütçe: {total_budget:,.0f}₺. "
               "Güven verici, sayısal, profesyonel ton.",
        fallback=(f"{tname}, {len(strategies)} stratejik direk üzerinde "
                  f"{proc_count} süreç ve {user_count} çalışanın katkısıyla "
                  f"sürdürülebilir büyüme yolculuğundadır. {len(inits)} aktif stratejik girişim "
                  f"ile {total_budget:,.0f}₺'lik yatırım portföyü yönetilmekte, "
                  f"{meas_count:,} veri noktası ile karar verme süreçleri desteklenmektedir. "
                  "Bu sunum, mevcut performansımızı ve yatırımcı önerimizi sunmaktadır."),
        tid=tid, endpoint="ai_yatirimci_sunum", max_tokens=400,
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(0x4f, 0x46, 0xe5)
    DARK = RGBColor(0x0f, 0x17, 0x2a)
    GRAY = RGBColor(0x64, 0x74, 0x8b)
    GREEN = RGBColor(0x10, 0xb9, 0x81)

    def add_slide_title(slide, txt, color):
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        p = tx.text_frame.paragraphs[0]
        p.text = txt
        p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = color

    def add_body(slide, text, bullets=None, font=14):
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
        bf = body.text_frame; bf.word_wrap = True
        if text:
            p = bf.paragraphs[0]; p.text = text
            p.font.size = Pt(font); p.font.color.rgb = DARK
        if bullets:
            for b in bullets:
                bp = bf.add_paragraph(); bp.text = "• " + str(b)
                bp.font.size = Pt(font - 2); bp.font.color.rgb = GRAY
                bp.space_before = Pt(4)

    def add_footer(slide, no):
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.3))
        ft.text_frame.text = f"{tname} · {year_label} · Slayt {no}"
        ft.text_frame.paragraphs[0].font.size = Pt(10)
        ft.text_frame.paragraphs[0].font.color.rgb = GRAY

    # Slaytları üret
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, tname, DARK)
    add_body(s, f"\n\n{year_label} {_('Yatırımcı Sunumu')}\n\n{_date.today().strftime('%d %B %Y')}", font=20)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, _("Yönetici Özeti"), BLUE)
    add_body(s, investor_summary); add_footer(s, 2)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, _("Vizyon"), BLUE)
    add_body(s, (tenant.vision if tenant else None) or _("Vizyon tanımlanmamış.")); add_footer(s, 3)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, _("Misyon ve Değerler"), BLUE)
    bullets = []
    if tenant and tenant.purpose: bullets.append(_("Amaç: ") + (tenant.purpose[:200]))
    if tenant and tenant.core_values: bullets.append(_("Değerler: ") + (tenant.core_values[:200]))
    add_body(s, None, bullets); add_footer(s, 4)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, _("Stratejik Direkler"), BLUE)
    add_body(s, f"{len(strategies)} {_('ana stratejik direk:')}",
        bullets=[f"{st.code} — {st.title}" for st in strategies[:8]])
    add_footer(s, 5)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, _("Operasyonel Performans"), BLUE)
    add_body(s, _("Yıllık operasyonel göstergeler:"),
        bullets=[f"{proc_count} {_('aktif süreç')}", f"{meas_count:,} {_('KPI ölçümü')}",
                 f"{user_count} {_('aktif çalışan')}",
                 _("Süreçler CMMI olgunluk seviyesinde izlenmektedir")])
    add_footer(s, 6)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, _("Finansal Görünüm"), BLUE)
    add_body(s, _("Stratejik girişim portföyü:"),
        bullets=[f"{_('Toplam initiative bütçesi:')} ₺{total_budget:,.0f}",
                 f"{_('Aktif initiative:')} {len(inits)}",
                 f"{_('Ortalama bütçe/initiative:')} ₺{total_budget/max(len(inits),1):,.0f}",
                 _("EVM (PV/EV/AC) modeliyle takip")])
    add_footer(s, 7)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, _("Initiative Portföyü"), BLUE)
    add_body(s, f"{_('En öne çıkan')} {min(8, len(inits))} {_('stratejik girişim:')}",
        bullets=[f"{i.code or 'INI'} — {i.name} (%{i.progress_pct or 0} {_('ilerleme, durum:')} {i.status})"
                 for i in sorted(inits, key=lambda x: -(x.progress_pct or 0))[:8]])
    add_footer(s, 8)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, _("İnsan Sermayesi"), BLUE)
    add_body(s, f"{user_count} {_('aktif çalışan ile yetenek odaklı büyüme.')}",
        bullets=[_("Bireysel performans takibi"), _("Süreç sahipliği ve liderlik geliştirme"),
                 _("2FA güvenlik standartları")])
    add_footer(s, 9)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, _("ESG ve Sürdürülebilirlik"), GREEN)
    add_body(s, _("Çevre, sosyal ve yönetişim taahhütleri:"),
        bullets=[_("Scope 1+2+3 emisyon takibi"), _("SDG katkı haritası"),
                 _("İş güvenliği ve çeşitlilik metrikleri"), _("Bağımsız yatırımcı raporlaması")])
    add_footer(s, 10)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, _("Risk Yönetimi"), BLUE)
    add_body(s, _("Stratejik risk register ile proaktif yönetim:"),
        bullets=[_("5×5 risk heatmap (olasılık × etki)"),
                 _("Sahip atamalı mitigation planları"),
                 _("Çeyreklik risk değerlendirme"), _("AI Erken Uyarı sistemi")])
    add_footer(s, 11)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, _("Yatırım Talebi"), BLUE)
    add_body(s, _("Önerilen yatırım kullanım planı:"),
        bullets=[_("Ürün geliştirme ve Ar-Ge (~%40)"),
                 _("Operasyonel kapasite (~%30)"),
                 _("Pazarlama ve büyüme (~%20)"),
                 _("Yedek (~%10)")])
    add_footer(s, 12)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, _("Teşekkürler"), DARK)
    add_body(s, f"\n\n{tname}\n\n{_('İletişim: investor@example.com')}\n\n{_('Detaylı analiz: Kokpitim platformu')}")
    add_footer(s, 13)

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    filename = f"{_hk_safe_name()}_{year_label}_yatirimci_sunum.pptx"
    return send_file(buf,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True, download_name=filename)


# ─── ES-08: GRI/CDP/TCFD ESG Rapor PDF ─────────────────────────────────────

@app_bp.route("/reports/esg-report")
@login_required
@require_module("raporlar")
def raporlar_esg_rapor():
    return render_template("platform/raporlar/esg_rapor.html")


@app_bp.route("/reports/api/esg-report/generate", methods=["GET", "POST"])
@login_required
@require_module("raporlar")
def raporlar_api_esg_rapor_generate():
    from flask import send_file
    from app.models.esg import EsgMetric, EsgMetricValue
    import io

    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    tenant = db.session.get(Tenant, tid)
    tname = tenant.name if tenant else "Kurum"
    year = _date.today().year

    h = _pdf_helpers()
    buf = io.BytesIO()
    doc = h["SimpleDocTemplate"](buf, pagesize=h["A4"],
        leftMargin=2*h["cm"], rightMargin=2*h["cm"], topMargin=2*h["cm"], bottomMargin=2*h["cm"],
        title=f"{tname} ESG Yıllık Raporu {year}")
    styles = h["styles"]
    h1 = h["ParagraphStyle"]("H1", parent=styles["Heading1"], fontSize=18,
        textColor=h["colors"].HexColor("#16a34a"), spaceAfter=14)
    h2 = h["ParagraphStyle"]("H2", parent=styles["Heading2"], fontSize=13,
        textColor=h["colors"].HexColor("#0f172a"), spaceAfter=8)
    body = h["ParagraphStyle"]("Body", parent=styles["BodyText"], fontSize=10.5,
        leading=15, textColor=h["colors"].HexColor("#0f172a"),
        spaceAfter=6, alignment=h["TA_JUSTIFY"])

    P = h["Paragraph"]; elems = []

    # Kapak
    elems.append(P(f"<font color='#16a34a'>{tname}</font>",
                    h["ParagraphStyle"]("Cover", parent=body, fontSize=24, alignment=h["TA_CENTER"])))
    elems.append(h["Spacer"](1, 0.3*h["cm"]))
    elems.append(P(f"{_('ESG YILLIK RAPOR')} — {year}",
                    h["ParagraphStyle"]("Cover2", parent=body, fontSize=18, alignment=h["TA_CENTER"],
                    textColor=h["colors"].HexColor("#475569"))))
    elems.append(h["Spacer"](1, 0.5*h["cm"]))
    elems.append(P(_("GRI · CDP · TCFD Uyumlu"),
                    h["ParagraphStyle"]("Cover3", parent=body, fontSize=11, alignment=h["TA_CENTER"],
                    textColor=h["colors"].HexColor("#94a3b8"))))
    elems.append(h["PageBreak"]())

    # GRI 102 — Genel Bilgiler
    elems.append(P(_("GRI 102 — Organizasyonel Profil"), h1))
    elems.append(P(f"<b>{_('Kurum:')}</b> {tname}", body))
    if tenant:
        if tenant.sector: elems.append(P(f"<b>{_('Sektör:')}</b> {tenant.sector}", body))
        if tenant.employee_count: elems.append(P(f"<b>{_('Çalışan sayısı:')}</b> {tenant.employee_count}", body))
        if tenant.vision: elems.append(P(f"<b>{_('Vizyon:')}</b> {tenant.vision[:300]}", body))
    elems.append(h["Spacer"](1, 0.5*h["cm"]))

    # GRI 305 — Emisyonlar
    elems.append(P(_("GRI 305 — Emisyonlar"), h1))
    e_metrics = EsgMetric.query.filter_by(tenant_id=tid, is_active=True).filter(
        EsgMetric.category == "E").all()
    if e_metrics:
        rows = [[_("Kod"), _("Metrik"), _("Scope"), _("Birim")]]
        for m in e_metrics:
            rows.append([m.code or "—", m.name or "—", m.scope or "—", m.unit or "—"])
        tbl = h["Table"](rows, colWidths=[2.5*h["cm"], 7*h["cm"], 3*h["cm"], 3*h["cm"]])
        tbl.setStyle(h["TableStyle"]([
            ("BACKGROUND", (0, 0), (-1, 0), h["colors"].HexColor("#16a34a")),
            ("TEXTCOLOR", (0, 0), (-1, 0), h["colors"].white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1),
             [h["colors"].white, h["colors"].HexColor("#f0fdf4")]),
            ("BOX", (0, 0), (-1, -1), 1, h["colors"].HexColor("#e2e8f0")),
            ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#e2e8f0")),
            ("PADDING", (0, 0), (-1, -1), 6),
        ]))
        elems.append(tbl)

        # En son yıl toplam
        by_year_scope = defaultdict(lambda: defaultdict(float))
        for m in e_metrics:
            for v in m.values:
                if v.value is not None and m.scope:
                    by_year_scope[v.year][m.scope] += v.value
        years_sorted = sorted(by_year_scope.keys())
        if years_sorted:
            latest = years_sorted[-1]
            elems.append(h["Spacer"](1, 0.5*h["cm"]))
            elems.append(P(f"<b>{latest} {_('Yılı Toplam Emisyonlar:')}</b>", body))
            for scope in ("scope1", "scope2", "scope3"):
                v = by_year_scope[latest].get(scope, 0)
                if v > 0:
                    elems.append(P(f"  {scope.upper()}: {v:,.2f} tCO₂e", body))
    else:
        elems.append(P(_("Henüz çevre metriği tanımlanmamış."), body))

    elems.append(h["PageBreak"]())

    # GRI 403 — Sosyal
    elems.append(P(_("GRI 403 — Sosyal Göstergeler (LTIFR, vb.)"), h1))
    s_metrics = EsgMetric.query.filter_by(tenant_id=tid, is_active=True, category="S").all()
    if s_metrics:
        for m in s_metrics:
            elems.append(P(f"<b>{m.name}</b>", body))
            elems.append(P(f"{_('Kod:')} {m.code or '—'} · {_('Birim:')} {m.unit or '—'}", body))
            if m.description: elems.append(P(m.description, body))
    else:
        elems.append(P(_("Henüz sosyal metriği tanımlanmamış."), body))

    elems.append(h["Spacer"](1, 1*h["cm"]))

    # GRI 102-22 — Yönetişim
    elems.append(P(_("GRI 102-22 — Yönetişim"), h1))
    g_metrics = EsgMetric.query.filter_by(tenant_id=tid, is_active=True, category="G").all()
    if g_metrics:
        for m in g_metrics:
            elems.append(P(f"<b>{m.name}:</b> {m.description or '—'}", body))
    else:
        elems.append(P(_("Henüz yönetişim metriği tanımlanmamış."), body))

    elems.append(h["PageBreak"]())

    # TCFD — İklim Risk Açıklaması (TASK-233: statik metin yerine kurum verisi)
    elems.append(P(_("TCFD — İklim Risk Açıklaması"), h1))
    elems.append(P(_("Yönetişim, Strateji, Risk Yönetimi, Metrik ve Hedefler dört "
                   "boyutta iklim ile ilgili açıklamalar:"), body))

    # 1. Yönetişim — G kategorisi metrik varlığı
    if g_metrics:
        elems.append(P(_("<b>1. Yönetişim:</b> Kurumda %(n)s yönetişim (G) metriği tanımlı "
                       "ve izlenmektedir; iklim riskleri bu çerçevede yönetim seviyesinde "
                       "değerlendirilmektedir.") % {"n": len(g_metrics)}, body))
    else:
        elems.append(P(_("<b>1. Yönetişim:</b> Henüz yönetişim (G) metriği tanımlanmamıştır; "
                       "iklim risklerinin yönetim kurulu seviyesinde izlenmesi için metrik "
                       "tanımlanması önerilir."), body))

    # 2. Strateji — hedef/baz değeri atanmış çevre metrikleri
    e_targeted = [m for m in (e_metrics or [])
                  if m.target_value is not None or m.baseline_value is not None]
    if e_targeted:
        elems.append(P(_("<b>2. Strateji:</b> %(n)s çevre metriği hedef/baz yılı ile "
                       "izlenmektedir:") % {"n": len(e_targeted)}, body))
        for m in e_targeted[:8]:
            parts = []
            if m.baseline_value is not None and m.baseline_year:
                parts.append(_("baz %(y)s: %(v)s") % {"y": m.baseline_year, "v": m.baseline_value})
            if m.target_value is not None:
                parts.append(_("hedef: %(v)s %(u)s") % {"v": m.target_value, "u": m.unit or ""})
            elems.append(P(f"  • {m.name} ({', '.join(parts)})", body))
    else:
        elems.append(P(_("<b>2. Strateji:</b> Çevre metriklerine henüz hedef/baz değer "
                       "atanmamıştır; kısa/orta/uzun vadeli stratejik etki analizi için "
                       "hedef tanımlanması önerilir."), body))

    # 3. Risk yönetimi — risk envanterindeki iklim/çevre başlıklı kayıtlar
    try:
        from app.models.k_radar_domain import RiskHeatmapItem
        _all_risks = RiskHeatmapItem.query.filter_by(tenant_id=tid, is_active=True).all()
        _kw = ("iklim", "çevre", "karbon", "enerji", "emisyon", "su ", "atık")
        climate_risks = [r for r in _all_risks
                         if any(k in (r.title or "").lower() for k in _kw)]
    except Exception:
        _all_risks, climate_risks = [], []
    if climate_risks:
        elems.append(P(_("<b>3. Risk yönetimi:</b> Risk envanterinde %(c)s iklim/çevre "
                       "ilişkili kayıt izlenmektedir (toplam %(t)s aktif risk):")
                       % {"c": len(climate_risks), "t": len(_all_risks)}, body))
        for r in sorted(climate_risks,
                        key=lambda x: -(x.rpn or x.probability * x.impact))[:5]:
            elems.append(P(f"  • {r.title} (RPN: {r.rpn or r.probability * r.impact})", body))
    else:
        elems.append(P(_("<b>3. Risk yönetimi:</b> Risk envanterinde iklim/çevre etiketli "
                       "kayıt bulunmamaktadır (toplam %(t)s aktif risk). İklim risklerinin "
                       "risk register'a entegrasyonu önerilir.") % {"t": len(_all_risks)}, body))

    # 4. Metrik ve hedefler — Scope 1/2/3 son yıl değerleri (E bölümündeki hesapla aynı)
    _scope_line = ""
    try:
        _scope_vals = by_year_scope.get(latest, {})
        _scope_line = " · ".join(
            f"{s.upper()}: {v:,.2f} tCO₂e" for s, v in _scope_vals.items() if v > 0)
    except NameError:
        pass  # e_metrics boşsa by_year_scope tanımsız — veri yok mesajına düşer
    if _scope_line:
        elems.append(P(_("<b>4. Metrik ve hedefler:</b> Son ölçüm yılı emisyonları — ")
                       + _scope_line, body))
    else:
        elems.append(P(_("<b>4. Metrik ve hedefler:</b> Scope 1/2/3 emisyon ölçümü henüz "
                       "kayıtlı değildir; ölçüm ve Net Zero hedefi tanımlanması önerilir."), body))

    # Kapanış
    elems.append(h["Spacer"](1, 1*h["cm"]))
    elems.append(P(_("Bu rapor Kokpitim platformu üzerinden, GRI/CDP/TCFD standart "
                   "yapısına uygun olarak otomatik üretilmiştir. Detay veri ve "
                   "yöntem için ek ekleri talep edebilirsiniz."),
                   h["ParagraphStyle"]("Foot", parent=body, fontSize=9,
                       textColor=h["colors"].HexColor("#64748b"))))

    doc.build(elems); buf.seek(0)
    filename = f"{_hk_safe_name()}_ESG_Raporu_{year}.pdf"
    return send_file(buf, mimetype="application/pdf",
                     as_attachment=True, download_name=filename)


# ─── RK-14: Audit Çıktı Paketi PDF ─────────────────────────────────────────

@app_bp.route("/reports/audit-package")
@login_required
@require_module("raporlar")
def raporlar_audit_paketi():
    return render_template("platform/raporlar/audit_paketi.html")


@app_bp.route("/reports/api/audit-package/generate", methods=["GET", "POST"])
@login_required
@require_module("raporlar")
def raporlar_api_audit_paketi_generate():
    from flask import send_file
    from app.models.audit import AuditLog
    import io

    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    tenant = db.session.get(Tenant, tid)
    tname = tenant.name if tenant else "Kurum"
    year = _date.today().year

    h = _pdf_helpers()
    buf = io.BytesIO()
    doc = h["SimpleDocTemplate"](buf, pagesize=h["A4"],
        leftMargin=2*h["cm"], rightMargin=2*h["cm"], topMargin=2*h["cm"], bottomMargin=2*h["cm"],
        title=f"{tname} Audit Çıktı Paketi {year}")
    styles = h["styles"]
    h1 = h["ParagraphStyle"]("H1", parent=styles["Heading1"], fontSize=18,
        textColor=h["colors"].HexColor("#dc2626"), spaceAfter=14)
    body = h["ParagraphStyle"]("Body", parent=styles["BodyText"], fontSize=10.5,
        leading=14, textColor=h["colors"].HexColor("#0f172a"),
        spaceAfter=6, alignment=h["TA_JUSTIFY"])
    small = h["ParagraphStyle"]("Sm", parent=body, fontSize=9,
        textColor=h["colors"].HexColor("#64748b"))

    P = h["Paragraph"]; elems = []

    # Kapak
    elems.append(P(f"<b>{tname}</b>",
                    h["ParagraphStyle"]("C1", parent=body, fontSize=24, alignment=h["TA_CENTER"])))
    elems.append(P(f"{_('AUDIT ÇIKTI PAKETİ')} — {year}",
                    h["ParagraphStyle"]("C2", parent=body, fontSize=16, alignment=h["TA_CENTER"],
                    textColor=h["colors"].HexColor("#dc2626"))))
    elems.append(h["Spacer"](1, 1*h["cm"]))
    elems.append(P(_("Bu doküman, üçüncü parti denetçilere kurum içi stratejik plan, "
                   "performans, risk ve kullanıcı aktivite verisinin özet sunumudur."),
                   h["ParagraphStyle"]("CI", parent=body, alignment=h["TA_CENTER"])))
    elems.append(h["PageBreak"]())

    # 1. Tenant özeti
    elems.append(P(_("1. Kurum Bilgileri"), h1))
    info_rows = [[_("Kurum"), tname]]
    if tenant:
        if tenant.short_name: info_rows.append([_("Kısa Ad"), tenant.short_name])
        if tenant.sector: info_rows.append([_("Sektör"), tenant.sector])
        if tenant.employee_count: info_rows.append([_("Çalışan"), str(tenant.employee_count)])
        info_rows.append([_("Plan Year Aktif"), _("Evet") if tenant.plan_year_enabled else _("Hayır")])
        info_rows.append([_("Tenant ID"), str(tid)])
    tbl = h["Table"](info_rows, colWidths=[5*h["cm"], 11*h["cm"]])
    tbl.setStyle(h["TableStyle"]([
        ("BACKGROUND", (0, 0), (0, -1), h["colors"].HexColor("#f8fafc")),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("BOX", (0, 0), (-1, -1), 0.5, h["colors"].HexColor("#cbd5e1")),
        ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#e2e8f0")),
        ("PADDING", (0, 0), (-1, -1), 8),
    ]))
    elems.append(tbl)
    elems.append(h["PageBreak"]())

    # 2. Stratejik plan özet
    elems.append(P(_("2. Stratejik Plan Yapısı"), h1))
    strat = Strategy.query.filter_by(tenant_id=tid, is_active=True).count()
    sub = db.session.query(func.count(SubStrategy.id)).join(Strategy).filter(
        Strategy.tenant_id == tid, SubStrategy.is_active.is_(True)).scalar() or 0
    proc = Process.query.filter_by(tenant_id=tid, is_active=True).count()
    kpi = ProcessKpi.query.join(Process).filter(
        Process.tenant_id == tid, ProcessKpi.is_active.is_(True)).count()
    sp_rows = [[_("Ana Strateji"), str(strat)], [_("Alt Strateji"), str(sub)],
               [_("Süreç"), str(proc)], [_("Performans Göstergesi"), str(kpi)]]
    tbl = h["Table"](sp_rows, colWidths=[8*h["cm"], 4*h["cm"]])
    tbl.setStyle(h["TableStyle"]([
        ("BACKGROUND", (0, 0), (0, -1), h["colors"].HexColor("#f8fafc")),
        ("FONTSIZE", (0, 0), (-1, -1), 11),
        ("BOX", (0, 0), (-1, -1), 0.5, h["colors"].HexColor("#cbd5e1")),
        ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#e2e8f0")),
        ("PADDING", (0, 0), (-1, -1), 8),
    ]))
    elems.append(tbl)
    elems.append(h["PageBreak"]())

    # 3. Audit log son 90 gün
    elems.append(P(_("3. Audit Log — Son 90 Gün"), h1))
    last_90 = datetime.now(timezone.utc) - timedelta(days=90)
    audit_logs = AuditLog.query.filter(
        AuditLog.tenant_id == tid, AuditLog.created_at >= last_90,
    ).order_by(AuditLog.created_at.desc()).limit(50).all()

    # Action dağılımı
    action_count = defaultdict(int)
    for a in audit_logs:
        action_count[a.action] += 1
    if action_count:
        elems.append(P(f"<b>{_('Aksiyon Dağılımı (top 50):')}</b>", body))
        for act, cnt in sorted(action_count.items(), key=lambda x: -x[1]):
            elems.append(P(f"  {act}: {cnt}", small))
        elems.append(h["Spacer"](1, 0.5*h["cm"]))

    # En son 20 audit log
    if audit_logs:
        elems.append(P(f"<b>{_('Son 20 İşlem:')}</b>", body))
        rows = [[_("Tarih"), _("Kullanıcı"), _("Aksiyon"), _("Kaynak")]]
        for a in audit_logs[:20]:
            rows.append([
                a.created_at.strftime("%d.%m %H:%M") if a.created_at else "—",
                (a.username or "—")[:25],
                a.action or "—",
                (a.resource_type or "—")[:20],
            ])
        tbl = h["Table"](rows, colWidths=[3*h["cm"], 4.5*h["cm"], 4*h["cm"], 4*h["cm"]])
        tbl.setStyle(h["TableStyle"]([
            ("BACKGROUND", (0, 0), (-1, 0), h["colors"].HexColor("#dc2626")),
            ("TEXTCOLOR", (0, 0), (-1, 0), h["colors"].white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1),
             [h["colors"].white, h["colors"].HexColor("#fef2f2")]),
            ("PADDING", (0, 0), (-1, -1), 5),
            ("BOX", (0, 0), (-1, -1), 0.5, h["colors"].HexColor("#fecaca")),
            ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#fecaca")),
        ]))
        elems.append(tbl)
    else:
        elems.append(P(_("Son 90 günde audit log kaydı yok."), small))
    elems.append(h["PageBreak"]())

    # 4. Kullanıcı + 2FA
    elems.append(P(_("4. Kullanıcı Güvenlik Özeti"), h1))
    total_u = User.query.filter_by(tenant_id=tid, is_active=True).count()
    totp_u = User.query.filter_by(tenant_id=tid, is_active=True, totp_enabled=True).count()
    elems.append(P(f"<b>{_('Toplam aktif kullanıcı:')}</b> {total_u}", body))
    elems.append(P(f"<b>{_('2FA etkin:')}</b> {totp_u} (%{round(totp_u/max(total_u,1)*100,1)})", body))
    elems.append(P(f"<b>{_('2FA pasif:')}</b> {total_u - totp_u}", body))

    # Footer
    elems.append(h["Spacer"](1, 2*h["cm"]))
    elems.append(P(f"{_('Bu rapor')} {_date.today().strftime('%d.%m.%Y')} {_('tarihinde Kokpitim platformundan otomatik üretilmiştir.')}",
                   h["ParagraphStyle"]("Foot", parent=small, alignment=h["TA_CENTER"])))

    doc.build(elems); buf.seek(0)
    filename = f"{_hk_safe_name()}_Audit_Paketi_{year}.pdf"
    return send_file(buf, mimetype="application/pdf",
                     as_attachment=True, download_name=filename)


# ─── HR-01: Bireysel Karne PDF Batch (ZIP) ─────────────────────────────────

@app_bp.route("/reports/individual-scorecard-batch")
@login_required
@require_module("raporlar")
def raporlar_bireysel_karne_batch():
    return render_template("platform/raporlar/bireysel_karne_batch.html")


@app_bp.route("/reports/api/individual-scorecard-batch/preview")
@login_required
@require_module("raporlar")
def raporlar_api_bireysel_karne_batch_preview():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    # PG'si olan kullanıcı sayısı
    users_with_pg = db.session.query(func.count(func.distinct(IndividualPerformanceIndicator.user_id))).join(
        User, User.id == IndividualPerformanceIndicator.user_id
    ).filter(
        User.tenant_id == tid, User.is_active.is_(True),
        IndividualPerformanceIndicator.is_active.is_(True),
    ).scalar() or 0
    total = User.query.filter_by(tenant_id=tid, is_active=True).count()
    return jsonify({"success": True, "preview": {
        "total_users": total, "users_with_pg": users_with_pg,
        "format": _("ZIP (PDF her kullanıcı için)"),
        "estimated_size_mb": round(users_with_pg * 0.05, 2),
    }})


@app_bp.route("/reports/api/individual-scorecard-batch/generate", methods=["GET", "POST"])
@login_required
@require_module("raporlar")
def raporlar_api_bireysel_karne_batch_generate():
    from flask import send_file
    import io, zipfile

    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    tenant = db.session.get(Tenant, tid)
    tname = tenant.name if tenant else "Kurum"

    # PG'si olan kullanıcılar
    user_ids_with_pg = [r[0] for r in db.session.query(
        func.distinct(IndividualPerformanceIndicator.user_id)
    ).join(User, User.id == IndividualPerformanceIndicator.user_id).filter(
        User.tenant_id == tid, User.is_active.is_(True),
        IndividualPerformanceIndicator.is_active.is_(True),
    ).all()]

    if not user_ids_with_pg:
        return jsonify({"success": False, "message": _("Bireysel PG'si olan kullanıcı yok.")}), 400

    h = _pdf_helpers()
    P = h["Paragraph"]
    styles = h["styles"]
    h1 = h["ParagraphStyle"]("H1", parent=styles["Heading1"], fontSize=18,
        textColor=h["colors"].HexColor("#4f46e5"), spaceAfter=12)
    body = h["ParagraphStyle"]("Body", parent=styles["BodyText"], fontSize=10.5,
        leading=14, textColor=h["colors"].HexColor("#0f172a"), spaceAfter=4)

    # Toplu olarak tüm kullanıcı + PG'leri tek sorguda topla (N+1 önlemi)
    _uids = user_ids_with_pg[:100]
    _users_map = {u.id: u for u in User.query.filter(User.id.in_(_uids)).all()} if _uids else {}
    _pgs_by_uid = defaultdict(list)
    if _uids:
        for p in IndividualPerformanceIndicator.query.filter(
            IndividualPerformanceIndicator.user_id.in_(_uids),
            IndividualPerformanceIndicator.is_active.is_(True),
        ).all():
            _pgs_by_uid[p.user_id].append(p)

    # ZIP container
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for uid in _uids:
            u = _users_map.get(uid)
            if not u:
                continue
            uname = f"{u.first_name or ''} {u.last_name or ''}".strip() or u.email
            pgs = _pgs_by_uid.get(uid, [])
            aligned = sum(1 for p in pgs if p.source_process_id or p.source_process_kpi_id)
            alignment_pct = round(aligned / max(len(pgs), 1) * 100, 1)

            # Tek PDF üret
            pdf_buf = io.BytesIO()
            doc = h["SimpleDocTemplate"](pdf_buf, pagesize=h["A4"],
                leftMargin=2*h["cm"], rightMargin=2*h["cm"], topMargin=2*h["cm"], bottomMargin=2*h["cm"],
                title=f"{uname} Bireysel Karne")
            elems = []
            elems.append(P(f"<b>{tname}</b>",
                h["ParagraphStyle"]("Top", parent=body, fontSize=10,
                    textColor=h["colors"].HexColor("#64748b"))))
            elems.append(P(_("BİREYSEL PERFORMANS KARNESİ"), h1))
            elems.append(P(f"<b>{uname}</b>", body))
            elems.append(P(f"{_('E-posta:')} {u.email}", body))
            if u.department: elems.append(P(f"{_('Departman:')} {u.department}", body))
            elems.append(h["Spacer"](1, 0.5*h["cm"]))
            elems.append(P(f"<b>{_('Özet:')}</b> {len(pgs)} {_('bireysel PG ·')} "
                f"{aligned} {_('kurum stratejisine bağlı')} (%{alignment_pct} {_('hizalama')})", body))
            elems.append(h["Spacer"](1, 0.5*h["cm"]))

            if pgs:
                rows = [[_("Kod"), _("PG Adı"), _("Hedef"), _("Birim"), _("Tip")]]
                for p in pgs[:30]:
                    rows.append([
                        (p.code or "—")[:15],
                        (p.name or "")[:40],
                        (p.target_value or "—")[:10],
                        (p.unit or "—")[:10],
                        _("Kurumdan") if p.source_process_id else _("Bireysel"),
                    ])
                tbl = h["Table"](rows, colWidths=[2*h["cm"], 7*h["cm"], 2.5*h["cm"], 2*h["cm"], 3*h["cm"]])
                tbl.setStyle(h["TableStyle"]([
                    ("BACKGROUND", (0, 0), (-1, 0), h["colors"].HexColor("#4f46e5")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), h["colors"].white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1),
                     [h["colors"].white, h["colors"].HexColor("#f8fafc")]),
                    ("PADDING", (0, 0), (-1, -1), 5),
                    ("BOX", (0, 0), (-1, -1), 0.5, h["colors"].HexColor("#e2e8f0")),
                    ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#e2e8f0")),
                ]))
                elems.append(tbl)
            elems.append(h["Spacer"](1, 1.5*h["cm"]))
            elems.append(P(_("Bu karne Kokpitim platformu üzerinden otomatik üretilmiştir."),
                h["ParagraphStyle"]("F", parent=body, fontSize=9,
                    textColor=h["colors"].HexColor("#94a3b8"))))

            doc.build(elems)
            pdf_buf.seek(0)
            safe_name = uname.replace(" ", "_").replace("/", "_")[:40]
            zf.writestr(f"{safe_name}_karne.pdf", pdf_buf.read())

    zip_buf.seek(0)
    filename = f"{_hk_safe_name()}_bireysel_karne_batch.zip"
    return send_file(zip_buf, mimetype="application/zip",
                     as_attachment=True, download_name=filename)


