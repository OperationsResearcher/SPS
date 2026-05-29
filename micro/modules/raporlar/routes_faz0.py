"""Faz 0 — Rapor 1-10 (veri kalitesi, K-vektör, sankey, departman, liderlik, initiative bubble, sabah özeti, evrim filmi, ai sunum)."""
from __future__ import annotations

from collections import defaultdict
from datetime import datetime, timedelta, date as _date

from flask import render_template, jsonify, request, current_app, send_file
from flask_login import login_required, current_user
from sqlalchemy import func, and_, or_, text, select
from sqlalchemy.orm import joinedload, selectinload

from platform_core import app_bp
from app.models import db
from app.models.core import User, Strategy, SubStrategy, Tenant
from app.models.process import (
    Process, ProcessKpi, KpiData, IndividualPerformanceIndicator,
    ProcessActivity, ProcessSubStrategyLink, process_leaders,
)
from app.models.k_vektor import KVektorStrategyWeight
from app.models.plan_year import PlanYear, KpiYearConfig
from app.models.initiative import Initiative
from app.services.plan_year_service import get_active_plan_year_for_user, list_plan_years
from app.services.score_engine_service import compute_process_scores_internal

from .helpers import _tid_or_none, MUDA_MAX_PROCESSES

# ─── Rapor 1: Veri Kalitesi ──────────────────────────────────────────────────

@app_bp.route("/raporlar/veri-kalitesi")
@login_required
def raporlar_veri_kalitesi():
    """Veri Kalitesi Raporu — PG doluluk, eksik alanlar, son giriş tarihleri."""
    from app.services.plan_year_service import list_plan_years
    years = []
    active_year = None
    if current_user.tenant_id:
        years = [py.year for py in list_plan_years(current_user.tenant_id)]
        apy = get_active_plan_year_for_user(current_user)
        active_year = apy.year if apy else None
    return render_template("platform/raporlar/veri_kalitesi.html",
                           plan_years=years, active_year=active_year)


@app_bp.route("/raporlar/api/veri-kalitesi")
@login_required
def raporlar_api_veri_kalitesi():
    """Veri kalitesi metriklerini JSON döner."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    year_param = request.args.get("year", type=int)
    if year_param:
        from app.services.plan_year_service import get_plan_year as _gpy
        active_py = _gpy(tid, year_param)
    else:
        active_py = get_active_plan_year_for_user(current_user)
    py_filter = (Process.plan_year_id == active_py.id) if active_py else True

    # Tüm aktif PG'ler (aktif plan yılındaki süreçlere bağlı) — N+1 önlemek için
    # process'i eager yükle (sonradan k.process.name erişiliyor)
    kpis_q = ProcessKpi.query.options(joinedload(ProcessKpi.process)).join(Process).filter(
        Process.tenant_id == tid,
        Process.is_active.is_(True),
        ProcessKpi.is_active.is_(True),
        py_filter if active_py else True,
    )
    all_kpis = kpis_q.all()
    total_kpi = len(all_kpis)

    if total_kpi == 0:
        return jsonify({
            "success": True,
            "summary": {"total_kpi": 0, "score": 0},
            "categories": {"kritik": [], "orta": [], "iyi": []},
            "department_breakdown": [],
            "by_process": [],
        })

    # Her PG için son ölçüm tarihi + ölçüm sayısı
    kpi_id_to_last = {}
    kpi_id_to_count = {}
    rows = db.session.query(
        KpiData.process_kpi_id,
        func.max(KpiData.data_date),
        func.count(KpiData.id),
    ).filter(
        KpiData.process_kpi_id.in_(select(ProcessKpi.id).where(ProcessKpi.id.in_([k.id for k in all_kpis]))),
        KpiData.is_active.is_(True),
    ).group_by(KpiData.process_kpi_id).all()
    for kid, last_date, cnt in rows:
        kpi_id_to_last[kid] = last_date
        kpi_id_to_count[kid] = cnt

    # Kategorize: kritik / orta / iyi
    today = datetime.utcnow().date()
    kritik, orta, iyi = [], [], []

    for k in all_kpis:
        last = kpi_id_to_last.get(k.id)
        cnt = kpi_id_to_count.get(k.id, 0)
        issues = []
        if not last:
            issues.append("hiç ölçüm yok")
        elif (today - last).days > 180:
            issues.append(f"son ölçüm {(today - last).days} gün önce")
        if not k.target_value:
            issues.append("hedef tanımsız")
        if not k.unit:
            issues.append("birim tanımsız")
        if not k.basari_puani_araliklari:
            issues.append("başarı puan aralığı tanımsız")

        item = {
            "id": k.id,
            "process_id": k.process_id,
            "code": k.code or f"PG-{k.id}",
            "name": k.name,
            "process_name": k.process.name if k.process else "?",
            "last_data_date": last.isoformat() if last else None,
            "data_count": cnt,
            "issues": issues,
        }
        if len(issues) >= 3 or (not last):
            kritik.append(item)
        elif issues:
            orta.append(item)
        else:
            iyi.append(item)

    # Süreç bazlı doluluk
    by_process = defaultdict(lambda: {"total": 0, "filled": 0, "name": ""})
    for k in all_kpis:
        pname = k.process.name if k.process else "?"
        by_process[k.process_id]["name"] = pname
        by_process[k.process_id]["total"] += 1
        if kpi_id_to_last.get(k.id):
            by_process[k.process_id]["filled"] += 1
    by_process_list = sorted(
        [{"process_id": pid, **v, "fill_rate": round(v["filled"] / v["total"] * 100, 1) if v["total"] else 0}
         for pid, v in by_process.items()],
        key=lambda x: x["fill_rate"]
    )

    # Genel skor
    score = round(
        (len(iyi) * 100 + len(orta) * 60 + len(kritik) * 20) / (total_kpi * 100) * 100, 1
    ) if total_kpi else 0

    return jsonify({
        "success": True,
        "summary": {
            "total_kpi": total_kpi,
            "score": score,
            "kritik_count": len(kritik),
            "orta_count": len(orta),
            "iyi_count": len(iyi),
            "plan_year": active_py.year if active_py else None,
        },
        "categories": {
            "kritik": sorted(kritik, key=lambda x: -len(x["issues"]))[:20],
            "orta": orta[:20],
            "iyi_count_only": len(iyi),
        },
        "by_process": by_process_list,
    })


# ─── Rapor 2: K-Vektör Çarpıklık ──────────────────────────────────────────────

@app_bp.route("/raporlar/k-vektor-carpiklik")
@login_required
def raporlar_kv_carpiklik():
    """K-Vektör Çarpıklık — strateji ağırlığı × skor uyumsuzluğu."""
    return render_template("platform/raporlar/kv_carpiklik.html")


@app_bp.route("/raporlar/api/k-vektor-carpiklik")
@login_required
def raporlar_api_kv_carpiklik():
    """Her ana stratejinin K-Vektör ağırlığı vs gerçek skor uyumsuzluğu."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    active_py = get_active_plan_year_for_user(current_user)
    py_id = active_py.id if active_py else None

    # Stratejiler + ağırlık
    strat_q = Strategy.query.options(selectinload(Strategy.sub_strategies)).filter_by(tenant_id=tid, is_active=True)
    if py_id:
        strat_q = strat_q.filter(or_(Strategy.plan_year_id == py_id, Strategy.plan_year_id.is_(None)))
    strategies = strat_q.order_by(Strategy.code).all()

    weights = {w.strategy_id: w.weight_raw for w in
               KVektorStrategyWeight.query.filter_by(tenant_id=tid).all()}

    # Skor motorundan strateji skorları
    try:
        today = datetime.utcnow().date()
        score_year = active_py.year if active_py else today.year
        proc_scores, _ = compute_process_scores_internal(
            tid, score_year, today, persist_pg_scores=False, plan_year=active_py
        )
    except Exception as e:
        current_app.logger.warning(f"[raporlar/kv-carpiklik] skor hesaplanamadı: {e}")
        proc_scores = {}

    # Her strateji için: alt-strateji aracılığıyla bağlı süreçlerin ortalama skoru
    rows = []
    total_weight = sum(weights.values()) if weights else 0
    for s in strategies:
        sub_ids = [ss.id for ss in s.sub_strategies if getattr(ss, "is_active", True)]
        # bu alt-stratejilere bağlı PG'ler -> hangi süreçlerden
        related_proc_ids = set()
        if sub_ids:
            related_proc_ids = set(r[0] for r in db.session.query(ProcessKpi.process_id).filter(
                ProcessKpi.sub_strategy_id.in_(sub_ids),
                ProcessKpi.is_active.is_(True),
            ).all())
        # bu süreçlerin skor ortalaması
        proc_score_vals = [proc_scores.get(pid) for pid in related_proc_ids
                           if proc_scores.get(pid) is not None]
        avg_score = round(sum(proc_score_vals) / len(proc_score_vals), 1) if proc_score_vals else None

        weight_raw = weights.get(s.id, 0) or 0
        weight_pct = round((weight_raw / total_weight * 100), 1) if total_weight else 0

        # Çarpıklık: yüksek ağırlık + düşük skor = sorunlu
        skew = None
        skew_label = ""
        if avg_score is not None and weight_pct > 0:
            # 0-100 skalada: weight_pct ağırlık, avg_score performans
            skew = round(weight_pct - avg_score, 1)
            if skew > 10:
                skew_label = "Ağır ama düşük performans"
            elif skew < -10:
                skew_label = "Hafif ama yüksek performans"
            else:
                skew_label = "Dengeli"

        rows.append({
            "code": s.code,
            "title": s.title,
            "weight_pct": weight_pct,
            "avg_score": avg_score,
            "skew": skew,
            "skew_label": skew_label,
            "related_process_count": len(related_proc_ids),
        })

    # Genel sağlık: kaç strateji "dengeli"
    balanced = sum(1 for r in rows if r["skew_label"] == "Dengeli")
    unbalanced = sum(1 for r in rows if r["skew_label"] in ("Ağır ama düşük performans", "Hafif ama yüksek performans"))

    return jsonify({
        "success": True,
        "summary": {
            "total_strategies": len(rows),
            "balanced": balanced,
            "unbalanced": unbalanced,
            "plan_year": active_py.year if active_py else None,
        },
        "strategies": rows,
    })


# ─── Rapor 3: Stratejik Hizalama Sankey ──────────────────────────────────────

@app_bp.route("/raporlar/hizalama-sankey")
@login_required
def raporlar_hizalama_sankey():
    return render_template("platform/raporlar/hizalama_sankey.html")


@app_bp.route("/raporlar/api/hizalama-sankey")
@login_required
def raporlar_api_hizalama_sankey():
    """Vizyon → Strateji → Alt Strateji → Süreç akış verisi."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    active_py = get_active_plan_year_for_user(current_user)
    py_id = active_py.id if active_py else None

    # Stratejiler + ağırlık
    weights = {w.strategy_id: w.weight_raw or 0 for w in
               KVektorStrategyWeight.query.filter_by(tenant_id=tid).all()}
    total_w = sum(weights.values()) or 1

    strat_q = Strategy.query.options(selectinload(Strategy.sub_strategies)).filter_by(tenant_id=tid, is_active=True)
    if py_id:
        strat_q = strat_q.filter(or_(Strategy.plan_year_id == py_id, Strategy.plan_year_id.is_(None)))
    strategies = strat_q.order_by(Strategy.code).all()

    nodes = [{"id": "vision", "label": "Vizyon", "level": 0, "color": "#0f172a"}]
    links = []
    node_index = {"vision": 0}

    for s in strategies:
        sid = f"s-{s.id}"
        node_index[sid] = len(nodes)
        nodes.append({
            "id": sid,
            "label": f"{s.code or ''} {s.title}".strip()[:40],
            "level": 1,
            "color": "#4f46e5",
        })
        w = weights.get(s.id, 0)
        links.append({"source": 0, "target": node_index[sid], "value": round(w / total_w * 100, 1)})

        # Alt stratejiler
        for ss in s.sub_strategies:
            if not getattr(ss, "is_active", True):
                continue
            ssid = f"ss-{ss.id}"
            node_index[ssid] = len(nodes)
            nodes.append({
                "id": ssid,
                "label": f"{ss.code or ''} {ss.title}".strip()[:40],
                "level": 2,
                "color": "#8b5cf6",
            })
            links.append({"source": node_index[sid], "target": node_index[ssid], "value": 1})

            # Bağlı süreçler (process_sub_strategy_links)
            for pssl in ss.process_sub_strategy_links:
                p = pssl.process
                if not p or not p.is_active:
                    continue
                if p.tenant_id != tid:
                    continue
                if py_id and p.plan_year_id and p.plan_year_id != py_id:
                    continue
                pid = f"p-{p.id}"
                if pid not in node_index:
                    node_index[pid] = len(nodes)
                    nodes.append({
                        "id": pid,
                        "label": f"{p.code or ''} {p.name}".strip()[:40],
                        "level": 3,
                        "color": "#10b981",
                    })
                contrib = pssl.contribution_pct or 1
                links.append({"source": node_index[ssid], "target": node_index[pid], "value": round(contrib, 1)})

    return jsonify({
        "success": True,
        "summary": {
            "vision_nodes": 1,
            "strategy_nodes": len(strategies),
            "sub_strategy_nodes": sum(1 for n in nodes if n["level"] == 2),
            "process_nodes": sum(1 for n in nodes if n["level"] == 3),
            "total_links": len(links),
            "plan_year": active_py.year if active_py else None,
        },
        "nodes": nodes,
        "links": links,
    })


# ─── Rapor 4: PG Hedef Revizyon Sıklığı ──────────────────────────────────────

@app_bp.route("/raporlar/hedef-revizyon")
@login_required
def raporlar_hedef_revizyon():
    return render_template("platform/raporlar/hedef_revizyon.html")


@app_bp.route("/raporlar/api/hedef-revizyon")
@login_required
def raporlar_api_hedef_revizyon():
    """Yıl bazlı hedef override sayımı (kpi_year_configs vs ProcessKpi.target_value)."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    # Tüm plan yıllarını + kpi_year_configs sayımı (override yapılmış PG'ler)
    py_list = list_plan_years(tid)
    py_data = []
    revised_kpis = []  # tek tek revize edilen PG detayı

    for py in py_list:
        configs = KpiYearConfig.query.filter_by(plan_year_id=py.id).all()
        revised_count = 0
        for cfg in configs:
            kpi = cfg.process_kpi
            if not kpi:
                continue
            # Override aktif sayılır: target_value/weight/period vs farklıysa
            differences = []
            if cfg.target_value and cfg.target_value != (kpi.target_value or ""):
                differences.append("hedef")
            if cfg.weight is not None and cfg.weight != (kpi.weight or 0):
                differences.append("ağırlık")
            if cfg.period and cfg.period != (kpi.period or ""):
                differences.append("periyot")
            if differences:
                revised_count += 1
                revised_kpis.append({
                    "year": py.year,
                    "kpi_code": kpi.code or f"PG-{kpi.id}",
                    "kpi_name": kpi.name,
                    "diff_fields": differences,
                    "base_target": kpi.target_value,
                    "year_target": cfg.target_value,
                })

        py_data.append({
            "year": py.year,
            "status": py.status,
            "config_count": len(configs),
            "revised_count": revised_count,
        })

    return jsonify({
        "success": True,
        "summary": {
            "total_years": len(py_list),
            "total_revisions": len(revised_kpis),
            "years_with_revisions": sum(1 for x in py_data if x["revised_count"] > 0),
        },
        "by_year": py_data,
        "revised_kpis": revised_kpis[:50],
    })


# ─── Rapor 5: Departman Performans Skoru ─────────────────────────────────────

@app_bp.route("/raporlar/departman-performans")
@login_required
def raporlar_departman_performans():
    return render_template("platform/raporlar/departman_performans.html")


@app_bp.route("/raporlar/api/departman-performans")
@login_required
def raporlar_api_departman_performans():
    """Departman bazlı kullanıcı + bireysel PG + atanan görev sayısı."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    rows = db.session.query(
        User.department,
        func.count(func.distinct(User.id)).label("user_count"),
    ).filter(
        User.tenant_id == tid,
        User.is_active.is_(True),
    ).group_by(User.department).all()

    # Departman başına bireysel PG sayısı + başarı medyanı
    dept_data = []
    for dept, user_count in rows:
        users_in_dept = db.session.query(User.id).filter(
            User.tenant_id == tid, User.is_active.is_(True),
            User.department == dept,
        ).subquery()

        pg_count = db.session.query(func.count(IndividualPerformanceIndicator.id)).filter(
            IndividualPerformanceIndicator.user_id.in_(users_in_dept),
            IndividualPerformanceIndicator.is_active.is_(True),
        ).scalar() or 0

        # Önemli PG sayısı (is_important)
        important_pg = db.session.query(func.count(IndividualPerformanceIndicator.id)).filter(
            IndividualPerformanceIndicator.user_id.in_(users_in_dept),
            IndividualPerformanceIndicator.is_active.is_(True),
            IndividualPerformanceIndicator.is_important.is_(True),
        ).scalar() or 0

        # Departman içindeki kullanıcıların 2FA oranı (security proxy)
        totp_enabled = db.session.query(func.count(User.id)).filter(
            User.tenant_id == tid, User.is_active.is_(True),
            User.department == dept, User.totp_enabled.is_(True),
        ).scalar() or 0

        dept_data.append({
            "department": dept or "(belirsiz)",
            "user_count": user_count,
            "pg_count": pg_count,
            "pg_per_user": round(pg_count / user_count, 2) if user_count else 0,
            "important_pg": important_pg,
            "totp_rate": round(totp_enabled / user_count * 100, 1) if user_count else 0,
        })

    dept_data.sort(key=lambda x: -x["user_count"])

    return jsonify({
        "success": True,
        "summary": {
            "total_departments": len(dept_data),
            "total_users": sum(d["user_count"] for d in dept_data),
            "total_pgs": sum(d["pg_count"] for d in dept_data),
        },
        "departments": dept_data,
    })


# ─── Rapor 6: Yönetici Liderlik Skoru ────────────────────────────────────────

@app_bp.route("/raporlar/yonetici-liderlik")
@login_required
def raporlar_yonetici_liderlik():
    return render_template("platform/raporlar/yonetici_liderlik.html")


@app_bp.route("/raporlar/api/yonetici-liderlik")
@login_required
def raporlar_api_yonetici_liderlik():
    """Süreç liderlerinin liderliğindeki süreçlerin ortalama performans skoru."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    active_py = get_active_plan_year_for_user(current_user)
    py_id = active_py.id if active_py else None

    # Süreç liderleri (M:N tablosu)
    rows = db.session.query(
        process_leaders.c.user_id,
        process_leaders.c.process_id,
    ).join(Process, Process.id == process_leaders.c.process_id).filter(
        Process.tenant_id == tid,
        Process.is_active.is_(True),
    ).all()

    if not rows:
        # Fallback: hiç process_leaders yoksa user.role==yonetici olanları göster
        managers = User.query.filter_by(
            tenant_id=tid, is_active=True,
        ).join(User.role).filter(
            (User.role.has(name="yonetici")) | (User.role.has(name="tenant_admin"))
        ).all()
        return jsonify({
            "success": True,
            "summary": {
                "total_leaders": len(managers),
                "has_leader_data": False,
                "note": "Süreç-lider eşleştirmesi bulunamadı. Rol bazlı yönetici listesi gösterildi.",
            },
            "leaders": [{
                "user_id": u.id,
                "name": f"{u.first_name or ''} {u.last_name or ''}".strip() or u.email,
                "email": u.email,
                "department": u.department,
                "led_process_count": 0,
                "avg_process_score": None,
            } for u in managers[:25]],
        })

    # Skor motorundan süreç skorları
    try:
        today = _date.today()
        score_year = active_py.year if active_py else today.year
        proc_scores, _ = compute_process_scores_internal(
            tid, score_year, today, persist_pg_scores=False, plan_year=active_py
        )
    except Exception:
        proc_scores = {}

    # Lider → bağlı süreçler agregasyonu
    leader_map = defaultdict(list)
    for uid, pid in rows:
        leader_map[uid].append(pid)

    leaders = []
    for uid, pids in leader_map.items():
        u = User.query.get(uid)
        if not u or not u.is_active:
            continue
        scores = [proc_scores.get(pid) for pid in pids if proc_scores.get(pid) is not None]
        avg = round(sum(scores) / len(scores), 1) if scores else None
        leaders.append({
            "user_id": uid,
            "name": f"{u.first_name or ''} {u.last_name or ''}".strip() or u.email,
            "email": u.email,
            "department": u.department,
            "led_process_count": len(pids),
            "avg_process_score": avg,
        })
    leaders.sort(key=lambda x: -(x["avg_process_score"] or 0))

    return jsonify({
        "success": True,
        "summary": {
            "total_leaders": len(leaders),
            "has_leader_data": True,
            "with_score": sum(1 for l in leaders if l["avg_process_score"] is not None),
            "avg_score_overall": round(
                sum(l["avg_process_score"] for l in leaders if l["avg_process_score"]) /
                max(sum(1 for l in leaders if l["avg_process_score"]), 1), 1
            ),
        },
        "leaders": leaders,
    })


# ─── Rapor 7: Initiative Portföy Bubble ──────────────────────────────────────

@app_bp.route("/raporlar/initiative-bubble")
@login_required
def raporlar_initiative_bubble():
    return render_template("platform/raporlar/initiative_bubble.html")


@app_bp.route("/raporlar/api/initiative-bubble")
@login_required
def raporlar_api_initiative_bubble():
    """Initiative portföyü: bütçe × ilerleme × öncelik."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    inits = Initiative.query.filter_by(tenant_id=tid, is_active=True).all()
    PRIO_COLOR = {
        "critical": "#dc2626", "high": "#f59e0b",
        "medium": "#0ea5e9", "low": "#94a3b8",
    }
    STATUS_OPACITY = {
        "completed": 0.55, "cancelled": 0.25, "on_hold": 0.4,
        "in_progress": 1.0, "planned": 0.7,
    }
    bubbles = []
    for i in inits:
        budget = float(i.budget_total or 0)
        spent = float(i.budget_spent or 0)
        bubbles.append({
            "id": i.id,
            "code": i.code or f"INI-{i.id}",
            "name": i.name,
            "status": i.status,
            "priority": i.priority,
            "color": PRIO_COLOR.get(i.priority, "#64748b"),
            "opacity": STATUS_OPACITY.get(i.status, 0.7),
            "budget_total": budget,
            "budget_spent": spent,
            "budget_usage_pct": round(spent / budget * 100, 1) if budget else 0,
            "progress_pct": i.progress_pct or 0,
            "start_year": i.start_year,
            "end_year": i.end_year,
            "duration": (i.end_year or 0) - (i.start_year or 0) + 1,
        })

    # Özet
    by_status = defaultdict(int)
    by_priority = defaultdict(int)
    for b in bubbles:
        by_status[b["status"]] += 1
        by_priority[b["priority"]] += 1

    return jsonify({
        "success": True,
        "summary": {
            "total": len(bubbles),
            "total_budget": sum(b["budget_total"] for b in bubbles),
            "total_spent": sum(b["budget_spent"] for b in bubbles),
            "by_status": dict(by_status),
            "by_priority": dict(by_priority),
            "avg_progress": round(sum(b["progress_pct"] for b in bubbles) / max(len(bubbles), 1), 1),
        },
        "bubbles": bubbles,
    })


# ─── Rapor 8: Operasyonel Sabah Özeti+ ──────────────────────────────────────

@app_bp.route("/raporlar/sabah-ozeti")
@login_required
def raporlar_sabah_ozeti():
    return render_template("platform/raporlar/sabah_ozeti.html")


@app_bp.route("/raporlar/api/sabah-ozeti")
@login_required
def raporlar_api_sabah_ozeti():
    """Bugünün ve son 7 günün operasyonel özeti."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    today = _date.today()
    week_ago = today - timedelta(days=7)
    week_ahead = today + timedelta(days=7)

    # Tenant süreçleri
    proc_ids_subq = db.session.query(Process.id).filter(
        Process.tenant_id == tid, Process.is_active.is_(True),
    ).subquery()

    # Faaliyet sayımları
    overdue = ProcessActivity.query.filter(
        ProcessActivity.process_id.in_(proc_ids_subq),
        ProcessActivity.is_active.is_(True),
        ProcessActivity.end_at < datetime.utcnow(),
        ProcessActivity.status != "Tamamlandı",
    ).count()

    today_due = ProcessActivity.query.filter(
        ProcessActivity.process_id.in_(proc_ids_subq),
        ProcessActivity.is_active.is_(True),
        func.date(ProcessActivity.end_at) == today,
    ).count()

    upcoming = ProcessActivity.query.filter(
        ProcessActivity.process_id.in_(proc_ids_subq),
        ProcessActivity.is_active.is_(True),
        ProcessActivity.end_at.between(today, week_ahead),
    ).count()

    completed_week = ProcessActivity.query.filter(
        ProcessActivity.process_id.in_(proc_ids_subq),
        ProcessActivity.status == "Tamamlandı",
        ProcessActivity.completed_at.between(week_ago, today + timedelta(days=1)),
    ).count()

    # KPI ölçüm akışı
    kpi_ids_subq = db.session.query(ProcessKpi.id).filter(
        ProcessKpi.process_id.in_(proc_ids_subq),
        ProcessKpi.is_active.is_(True),
    ).subquery()

    measurements_today = KpiData.query.filter(
        KpiData.process_kpi_id.in_(kpi_ids_subq),
        KpiData.is_active.is_(True),
        func.date(KpiData.created_at) == today,
    ).count()

    measurements_week = KpiData.query.filter(
        KpiData.process_kpi_id.in_(kpi_ids_subq),
        KpiData.is_active.is_(True),
        KpiData.created_at >= week_ago,
    ).count()

    # En son 5 ölçüm girişi (recent activity)
    recent_data = db.session.query(
        KpiData.id, KpiData.created_at, KpiData.actual_value,
        ProcessKpi.name.label("kpi_name"),
        ProcessKpi.code.label("kpi_code"),
        User.first_name, User.last_name, User.email,
    ).join(ProcessKpi, ProcessKpi.id == KpiData.process_kpi_id).join(
        Process, Process.id == ProcessKpi.process_id
    ).join(User, User.id == KpiData.user_id, isouter=True).filter(
        Process.tenant_id == tid,
        KpiData.is_active.is_(True),
    ).order_by(KpiData.created_at.desc()).limit(8).all()

    recent_entries = [{
        "kpi_code": r.kpi_code,
        "kpi_name": r.kpi_name,
        "value": r.actual_value,
        "when": r.created_at.isoformat() if r.created_at else None,
        "who": (f"{r.first_name or ''} {r.last_name or ''}").strip() or r.email or "—",
    } for r in recent_data]

    return jsonify({
        "success": True,
        "today": today.isoformat(),
        "metrics": {
            "overdue_activities": overdue,
            "today_due_activities": today_due,
            "upcoming_7d_activities": upcoming,
            "completed_7d_activities": completed_week,
            "measurements_today": measurements_today,
            "measurements_7d": measurements_week,
        },
        "recent_entries": recent_entries,
    })


# ─── Rapor 9: Yıllar Arası Evrim Filmi ──────────────────────────────────────

@app_bp.route("/raporlar/evrim-filmi")
@login_required
def raporlar_evrim_filmi():
    return render_template("platform/raporlar/evrim_filmi.html")


@app_bp.route("/raporlar/api/evrim-filmi")
@login_required
def raporlar_api_evrim_filmi():
    """Yıllar boyunca strateji ağacının evrimi — her yıl için snapshot."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    py_list = list_plan_years(tid)
    if not py_list:
        return jsonify({"success": True, "years": [], "summary": {"total_years": 0}})

    snapshots = []
    for py in sorted(py_list, key=lambda x: x.year):
        # Bu plan yılındaki tüm strateji-alt-süreç-PG
        strategies = Strategy.query.options(selectinload(Strategy.sub_strategies)).filter_by(
            tenant_id=tid, plan_year_id=py.id, is_active=True
        ).order_by(Strategy.code).all()

        nodes = []
        for s in strategies:
            sub_list = []
            for ss in s.sub_strategies:
                if not getattr(ss, "is_active", True):
                    continue
                # Bu alt-stratejiye bağlı süreçler (process_sub_strategy_links üzerinden)
                proc_codes = []
                for pssl in ss.process_sub_strategy_links:
                    p = pssl.process
                    if p and p.is_active and (not p.plan_year_id or p.plan_year_id == py.id):
                        proc_codes.append(p.code or f"P-{p.id}")
                sub_list.append({
                    "code": ss.code or "",
                    "title": ss.title or "",
                    "process_codes": list(set(proc_codes)),
                })
            nodes.append({
                "code": s.code or "",
                "title": s.title or "",
                "id": s.id,
                "sub_strategies": sub_list,
            })

        # Yılın özet sayıları
        proc_count = Process.query.filter_by(
            tenant_id=tid, plan_year_id=py.id, is_active=True
        ).count()
        kpi_count = ProcessKpi.query.join(Process).filter(
            Process.tenant_id == tid,
            Process.plan_year_id == py.id,
            ProcessKpi.is_active.is_(True),
            Process.is_active.is_(True),
        ).count()
        _pk_subq = (
            select(ProcessKpi.id)
            .join(Process, ProcessKpi.process_id == Process.id)
            .where(Process.tenant_id == tid, Process.plan_year_id == py.id)
        )
        meas_count = db.session.query(func.count(KpiData.id)).filter(
            KpiData.year == py.year,
            KpiData.process_kpi_id.in_(_pk_subq),
        ).scalar() or 0

        snapshots.append({
            "year": py.year,
            "status": py.status,
            "strategies": nodes,
            "stats": {
                "strategy_count": len(strategies),
                "sub_count": sum(len(s["sub_strategies"]) for s in nodes),
                "process_count": proc_count,
                "kpi_count": kpi_count,
                "measurement_count": meas_count,
            },
        })

    return jsonify({
        "success": True,
        "summary": {
            "total_years": len(snapshots),
            "year_range": f"{snapshots[0]['year']}–{snapshots[-1]['year']}" if snapshots else "—",
        },
        "years": snapshots,
    })


# ─── Rapor 10: AI Yıl Sonu Sunum Üretici ────────────────────────────────────

@app_bp.route("/raporlar/ai-sunum")
@login_required
def raporlar_ai_sunum():
    return render_template("platform/raporlar/ai_sunum.html")


@app_bp.route("/raporlar/api/ai-sunum/preview")
@login_required
def raporlar_api_ai_sunum_preview():
    """Sunumun veri özetini ve üretilecek slayt başlıklarını döner (preview)."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    active_py = get_active_plan_year_for_user(current_user)
    tenant = db.session.get(Tenant, tid)

    strat_count = Strategy.query.filter_by(
        tenant_id=tid, is_active=True,
        plan_year_id=active_py.id if active_py else None,
    ).count() if active_py else Strategy.query.filter_by(tenant_id=tid, is_active=True).count()

    proc_count = Process.query.filter_by(
        tenant_id=tid, is_active=True,
        plan_year_id=active_py.id if active_py else None,
    ).count() if active_py else Process.query.filter_by(tenant_id=tid, is_active=True).count()

    init_count = Initiative.query.filter_by(tenant_id=tid, is_active=True).count()
    user_count = User.query.filter_by(tenant_id=tid, is_active=True).count()
    meas_count = db.session.query(func.count(KpiData.id)).join(
        ProcessKpi
    ).join(Process).filter(Process.tenant_id == tid).scalar() or 0

    slides = [
        {"no": 1, "title": "Kapak", "content": f"{tenant.name if tenant else 'Kurum'} — {active_py.year if active_py else 'Yıllık'} Stratejik Plan Sunumu"},
        {"no": 2, "title": "Yönetici Özeti", "content": "AI tarafından üretilen 3-cümle özet"},
        {"no": 3, "title": "Vizyon ve Misyon", "content": "Kurumsal kimlik alanları"},
        {"no": 4, "title": "Stratejik Yıl Özeti", "content": f"{strat_count} strateji · {proc_count} süreç · {init_count} initiative"},
        {"no": 5, "title": "Strateji Ağacı", "content": "Ana stratejiler + ağırlık dağılımı"},
        {"no": 6, "title": "Süreç Sağlık Heatmap", "content": f"{proc_count} süreç performans özeti"},
        {"no": 7, "title": "K-Vektör Analizi", "content": "Ağırlık × performans çarpıklık tespiti"},
        {"no": 8, "title": "Initiative Portföyü", "content": f"{init_count} girişimin bütçe + ilerleme durumu"},
        {"no": 9, "title": "Yıllar Arası Değişim", "content": "Önceki yılla karşılaştırma — strateji, süreç, PG farkları"},
        {"no": 10, "title": "İnsan Kaynakları", "content": f"{user_count} çalışan + departman dağılımı"},
        {"no": 11, "title": "Risk Haritası", "content": "Top 10 risk + olasılık × etki matrisi"},
        {"no": 12, "title": "ESG & Sürdürülebilirlik", "content": "E/S/G metrikleri + SDG katkı"},
        {"no": 13, "title": "Veri Kalitesi", "content": f"{meas_count:,} toplam ölçüm · doluluk oranları"},
        {"no": 14, "title": "Önümüzdeki Dönem Yol Haritası", "content": "Yeni initiative önerileri + öncelikler"},
        {"no": 15, "title": "Kapanış", "content": "Teşekkür + Q&A"},
    ]

    return jsonify({
        "success": True,
        "preview": {
            "tenant_name": tenant.name if tenant else "—",
            "plan_year": active_py.year if active_py else None,
            "data_points": {
                "strategies": strat_count,
                "processes": proc_count,
                "initiatives": init_count,
                "users": user_count,
                "kpi_measurements": meas_count,
            },
            "slides": slides,
            "estimated_size_kb": 280,
            "ai_tokens_estimate": 4500,
        },
    })


@app_bp.route("/raporlar/api/ai-sunum/generate", methods=["GET", "POST"])
@login_required
def raporlar_api_ai_sunum_generate():
    """Gerçek PowerPoint dosyasını üretir + indirme URL'i döner."""
    from flask import send_file
    import io
    import os
    from datetime import datetime as _dt

    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return jsonify({
            "success": False,
            "message": "python-pptx kurulu değil. Kurulum: pip install python-pptx",
        }), 500

    active_py = get_active_plan_year_for_user(current_user)
    tenant = db.session.get(Tenant, tid)

    # Veri özetleri
    tenant_name = tenant.name if tenant else "Kurum"
    year_label = str(active_py.year) if active_py else _dt.utcnow().year

    strategies = []
    if active_py:
        strat_q = Strategy.query.filter_by(
            tenant_id=tid, is_active=True, plan_year_id=active_py.id,
        ).order_by(Strategy.code).all()
        weights = {w.strategy_id: w.weight_raw for w in
                   KVektorStrategyWeight.query.filter_by(tenant_id=tid).all()}
        total_w = sum(weights.values()) or 1
        for s in strat_q:
            strategies.append({
                "code": s.code, "title": s.title,
                "weight_pct": round((weights.get(s.id, 0) or 0) / total_w * 100, 1),
            })

    proc_count = Process.query.filter_by(
        tenant_id=tid, is_active=True,
        plan_year_id=active_py.id if active_py else None,
    ).count()
    kpi_count = ProcessKpi.query.join(Process).filter(
        Process.tenant_id == tid,
        Process.plan_year_id == (active_py.id if active_py else None),
        ProcessKpi.is_active.is_(True),
    ).count()
    initiatives = Initiative.query.filter_by(tenant_id=tid, is_active=True).all()
    user_count = User.query.filter_by(tenant_id=tid, is_active=True).count()
    meas_count = db.session.query(func.count(KpiData.id)).join(ProcessKpi).join(
        Process
    ).filter(Process.tenant_id == tid).scalar() or 0

    # AI özeti — LLM gateway ile (yoksa fallback şablon)
    ai_summary = (
        f"{tenant_name} {year_label} dönemini {len(strategies)} ana strateji, "
        f"{proc_count} aktif süreç ve {kpi_count} performans göstergesi ile yönetti. "
        f"Yıl içinde {len(initiatives)} stratejik girişim portföyünde takip edildi ve "
        f"{user_count} çalışan tarafından toplam {meas_count:,} ölçüm verisi sisteme girildi. "
        "Önümüzdeki dönemde stratejik öncelikler ve yapısal iyileştirmeler için yol haritası bu sunumda yer almaktadır."
    )
    try:
        from app.services.llm_gateway import call_llm  # opsiyonel
        prompt = (
            f"Türkçe, 4 cümlelik yönetici özeti yaz. Kurum: {tenant_name}, Yıl: {year_label}. "
            f"Veri: {len(strategies)} strateji, {proc_count} süreç, {kpi_count} PG, "
            f"{len(initiatives)} initiative, {user_count} çalışan, {meas_count} ölçüm. "
            "Resmî ama özlü ton kullan."
        )
        result = call_llm(prompt=prompt, tenant_id=tid, user_id=current_user.id,
                          endpoint="ai_yil_sonu_sunum", max_output_tokens=400)
        if result and result.get("text"):
            ai_summary = result["text"].strip()
    except Exception as e:
        current_app.logger.info(f"[ai-sunum] LLM yok/başarısız ({e}), fallback metin kullanıldı")

    # PowerPoint oluştur
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(0x4f, 0x46, 0xe5)
    DARK = RGBColor(0x0f, 0x17, 0x2a)
    GRAY = RGBColor(0x64, 0x74, 0x8b)
    LIGHT_BG = RGBColor(0xf8, 0xfa, 0xfc)

    def add_title_slide(prs, title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Title
        tx = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = DARK
        # Subtitle
        sub = slide.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(12), Inches(1))
        sub.text_frame.text = subtitle
        sub.text_frame.paragraphs[0].font.size = Pt(20)
        sub.text_frame.paragraphs[0].font.color.rgb = GRAY
        return slide

    def add_content_slide(prs, slide_no, title, body_text, body_bullets=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Header bar
        from pptx.shapes.autoshape import Shape  # noqa
        # Title
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BLUE
        # Body
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
        bf = body.text_frame
        bf.word_wrap = True
        if body_text:
            p0 = bf.paragraphs[0]
            p0.text = body_text
            p0.font.size = Pt(16)
            p0.font.color.rgb = DARK
        if body_bullets:
            for bul in body_bullets:
                bp = bf.add_paragraph()
                bp.text = "• " + str(bul)
                bp.font.size = Pt(14)
                bp.font.color.rgb = GRAY
                bp.space_before = Pt(6)
        # Footer
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.3))
        ft.text_frame.text = f"{tenant_name} · {year_label} · Slayt {slide_no}"
        ft.text_frame.paragraphs[0].font.size = Pt(10)
        ft.text_frame.paragraphs[0].font.color.rgb = GRAY
        return slide

    # Slaytları üret
    add_title_slide(prs,
        f"{tenant_name}",
        f"{year_label} Stratejik Plan Sunumu · {_dt.utcnow().strftime('%d %B %Y')}")

    add_content_slide(prs, 2, "Yönetici Özeti", ai_summary)

    add_content_slide(prs, 3, f"Stratejik Yıl Özeti — {year_label}",
        "Tüm sistem aşağıdaki ana boyutlarda izlendi:",
        body_bullets=[
            f"{len(strategies)} ana strateji",
            f"{proc_count} aktif süreç",
            f"{kpi_count} performans göstergesi",
            f"{len(initiatives)} stratejik girişim (initiative)",
            f"{user_count} aktif çalışan",
            f"{meas_count:,} toplam KPI ölçümü",
        ])

    add_content_slide(prs, 4, "Ana Stratejiler ve K-Vektör Ağırlıkları",
        None, body_bullets=[
            f"{s['code']} — {s['title']} (Ağırlık: %{s['weight_pct']})"
            for s in strategies[:12]
        ] or ["Aktif plan yılında strateji bulunamadı."])

    add_content_slide(prs, 5, "Initiative Portföyü",
        f"Aktif ve tamamlanmış toplam {len(initiatives)} girişim.",
        body_bullets=[
            f"{i.code or 'INI-'+str(i.id)} — {i.name} (durum: {i.status}, ilerleme: %{i.progress_pct or 0})"
            for i in initiatives[:10]
        ])

    add_content_slide(prs, 6, "Veri Akışı ve Ölçüm Hacmi",
        f"{tenant_name} sistemine girilen toplam {meas_count:,} ölçüm verisi.",
        body_bullets=[
            f"Aktif süreç: {proc_count}",
            f"Aktif PG: {kpi_count}",
            f"Ortalama ölçüm/PG: {round(meas_count / max(kpi_count, 1), 1)}",
            "Veri girişi düzenli sürdürüldü — kalite skoru raporu detay sağlar.",
        ])

    add_content_slide(prs, 7, "Önümüzdeki Dönem — Öneriler",
        "Bu sunumun arka planındaki Kokpitim verisinden çıkan öneriler:",
        body_bullets=[
            "Stratejik öncelikleri K-Vektör çarpıklık raporundan gözden geçir",
            "Yıllar arası karşılaştırma sayfasından evrim haritasını incele",
            "Veri kalitesi raporundaki eksik PG'leri planla",
            "Initiative portföyünde geri kalanlara odak",
        ])

    add_content_slide(prs, 8, "Teşekkürler",
        f"{tenant_name} ekibinin katkılarıyla hazırlandı.",
        body_bullets=["Sorular ve detay için: kokpitim.com"])

    # Bellek üzerinden gönder
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = f"{tenant_name.replace(' ', '_')}_{year_label}_yil_sonu_sunum.pptx"
    return send_file(
        buf,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=filename,
    )


