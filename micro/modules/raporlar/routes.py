"""Raporlar modülü — yeni rapor önerileri için staging.

Sol menüden "Raporlar" → /raporlar landing → kartlar → tek tek raporlar.
docs/rapor/27mayisraporu.md'deki önerilerden seçilenler burada inşa edilir.
"""
from __future__ import annotations

from collections import defaultdict
from datetime import datetime, timedelta, date as _date

from flask import render_template, jsonify, request, current_app
from flask_login import login_required, current_user
from sqlalchemy import func, and_, or_, text

from platform_core import app_bp
from app.models import db
from app.models.core import User, Strategy, SubStrategy, Tenant
from app.models.process import (
    Process, ProcessKpi, KpiData, IndividualPerformanceIndicator,
    ProcessActivity, ProcessSubStrategyLink, process_leaders,
)
from app.models.k_vektor import KVektorStrategyWeight
from app.models.plan_year import PlanYear, KpiYearConfig
from app.models.initiative import Initiative
from app.services.plan_year_service import get_active_plan_year_for_user, list_plan_years
from app.services.score_engine_service import compute_process_scores_internal


# ─── /raporlar — Landing kart-grid ────────────────────────────────────────────

@app_bp.route("/raporlar")
@login_required
def raporlar_index():
    """Yeni rapor önerileri için landing — kart grid."""
    return render_template("platform/raporlar/index.html")


# ─── Rapor 1: Veri Kalitesi ──────────────────────────────────────────────────

@app_bp.route("/raporlar/veri-kalitesi")
@login_required
def raporlar_veri_kalitesi():
    """Veri Kalitesi Raporu — PG doluluk, eksik alanlar, son giriş tarihleri."""
    return render_template("platform/raporlar/veri_kalitesi.html")


@app_bp.route("/raporlar/api/veri-kalitesi")
@login_required
def raporlar_api_veri_kalitesi():
    """Veri kalitesi metriklerini JSON döner."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    active_py = get_active_plan_year_for_user(current_user)
    py_filter = (Process.plan_year_id == active_py.id) if active_py else True

    # Tüm aktif PG'ler (aktif plan yılındaki süreçlere bağlı)
    kpis_q = ProcessKpi.query.join(Process).filter(
        Process.tenant_id == tid,
        Process.is_active.is_(True),
        ProcessKpi.is_active.is_(True),
        py_filter if active_py else True,
    )
    all_kpis = kpis_q.all()
    total_kpi = len(all_kpis)

    if total_kpi == 0:
        return jsonify({
            "success": True,
            "summary": {"total_kpi": 0, "score": 0},
            "categories": {"kritik": [], "orta": [], "iyi": []},
            "department_breakdown": [],
            "by_process": [],
        })

    # Her PG için son ölçüm tarihi + ölçüm sayısı
    kpi_id_to_last = {}
    kpi_id_to_count = {}
    rows = db.session.query(
        KpiData.process_kpi_id,
        func.max(KpiData.data_date),
        func.count(KpiData.id),
    ).filter(
        KpiData.process_kpi_id.in_([k.id for k in all_kpis]),
        KpiData.is_active.is_(True),
    ).group_by(KpiData.process_kpi_id).all()
    for kid, last_date, cnt in rows:
        kpi_id_to_last[kid] = last_date
        kpi_id_to_count[kid] = cnt

    # Kategorize: kritik / orta / iyi
    today = datetime.utcnow().date()
    kritik, orta, iyi = [], [], []

    for k in all_kpis:
        last = kpi_id_to_last.get(k.id)
        cnt = kpi_id_to_count.get(k.id, 0)
        issues = []
        if not last:
            issues.append("hiç ölçüm yok")
        elif (today - last).days > 180:
            issues.append(f"son ölçüm {(today - last).days} gün önce")
        if not k.target_value:
            issues.append("hedef tanımsız")
        if not k.unit:
            issues.append("birim tanımsız")
        if not k.basari_puani_araliklari:
            issues.append("başarı puan aralığı tanımsız")

        item = {
            "id": k.id,
            "code": k.code or f"PG-{k.id}",
            "name": k.name,
            "process_name": k.process.name if k.process else "?",
            "last_data_date": last.isoformat() if last else None,
            "data_count": cnt,
            "issues": issues,
        }
        if len(issues) >= 3 or (not last):
            kritik.append(item)
        elif issues:
            orta.append(item)
        else:
            iyi.append(item)

    # Süreç bazlı doluluk
    by_process = defaultdict(lambda: {"total": 0, "filled": 0, "name": ""})
    for k in all_kpis:
        pname = k.process.name if k.process else "?"
        by_process[k.process_id]["name"] = pname
        by_process[k.process_id]["total"] += 1
        if kpi_id_to_last.get(k.id):
            by_process[k.process_id]["filled"] += 1
    by_process_list = sorted(
        [{"process_id": pid, **v, "fill_rate": round(v["filled"] / v["total"] * 100, 1) if v["total"] else 0}
         for pid, v in by_process.items()],
        key=lambda x: x["fill_rate"]
    )

    # Genel skor
    score = round(
        (len(iyi) * 100 + len(orta) * 60 + len(kritik) * 20) / (total_kpi * 100) * 100, 1
    ) if total_kpi else 0

    return jsonify({
        "success": True,
        "summary": {
            "total_kpi": total_kpi,
            "score": score,
            "kritik_count": len(kritik),
            "orta_count": len(orta),
            "iyi_count": len(iyi),
            "plan_year": active_py.year if active_py else None,
        },
        "categories": {
            "kritik": sorted(kritik, key=lambda x: -len(x["issues"]))[:20],
            "orta": orta[:20],
            "iyi_count_only": len(iyi),
        },
        "by_process": by_process_list,
    })


# ─── Rapor 2: K-Vektör Çarpıklık ──────────────────────────────────────────────

@app_bp.route("/raporlar/k-vektor-carpiklik")
@login_required
def raporlar_kv_carpiklik():
    """K-Vektör Çarpıklık — strateji ağırlığı × skor uyumsuzluğu."""
    return render_template("platform/raporlar/kv_carpiklik.html")


@app_bp.route("/raporlar/api/k-vektor-carpiklik")
@login_required
def raporlar_api_kv_carpiklik():
    """Her ana stratejinin K-Vektör ağırlığı vs gerçek skor uyumsuzluğu."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    active_py = get_active_plan_year_for_user(current_user)
    py_id = active_py.id if active_py else None

    # Stratejiler + ağırlık
    strat_q = Strategy.query.filter_by(tenant_id=tid, is_active=True)
    if py_id:
        strat_q = strat_q.filter(or_(Strategy.plan_year_id == py_id, Strategy.plan_year_id.is_(None)))
    strategies = strat_q.order_by(Strategy.code).all()

    weights = {w.strategy_id: w.weight_raw for w in
               KVektorStrategyWeight.query.filter_by(tenant_id=tid).all()}

    # Skor motorundan strateji skorları
    try:
        today = datetime.utcnow().date()
        score_year = active_py.year if active_py else today.year
        proc_scores, _ = compute_process_scores_internal(
            tid, score_year, today, persist_pg_scores=False, plan_year=active_py
        )
    except Exception as e:
        current_app.logger.warning(f"[raporlar/kv-carpiklik] skor hesaplanamadı: {e}")
        proc_scores = {}

    # Her strateji için: alt-strateji aracılığıyla bağlı süreçlerin ortalama skoru
    rows = []
    total_weight = sum(weights.values()) if weights else 0
    for s in strategies:
        sub_ids = [ss.id for ss in s.sub_strategies if getattr(ss, "is_active", True)]
        # bu alt-stratejilere bağlı PG'ler -> hangi süreçlerden
        related_proc_ids = set()
        if sub_ids:
            related_proc_ids = set(r[0] for r in db.session.query(ProcessKpi.process_id).filter(
                ProcessKpi.sub_strategy_id.in_(sub_ids),
                ProcessKpi.is_active.is_(True),
            ).all())
        # bu süreçlerin skor ortalaması
        proc_score_vals = [proc_scores.get(pid) for pid in related_proc_ids
                           if proc_scores.get(pid) is not None]
        avg_score = round(sum(proc_score_vals) / len(proc_score_vals), 1) if proc_score_vals else None

        weight_raw = weights.get(s.id, 0) or 0
        weight_pct = round((weight_raw / total_weight * 100), 1) if total_weight else 0

        # Çarpıklık: yüksek ağırlık + düşük skor = sorunlu
        skew = None
        skew_label = ""
        if avg_score is not None and weight_pct > 0:
            # 0-100 skalada: weight_pct ağırlık, avg_score performans
            skew = round(weight_pct - avg_score, 1)
            if skew > 10:
                skew_label = "Ağır ama düşük performans"
            elif skew < -10:
                skew_label = "Hafif ama yüksek performans"
            else:
                skew_label = "Dengeli"

        rows.append({
            "code": s.code,
            "title": s.title,
            "weight_pct": weight_pct,
            "avg_score": avg_score,
            "skew": skew,
            "skew_label": skew_label,
            "related_process_count": len(related_proc_ids),
        })

    # Genel sağlık: kaç strateji "dengeli"
    balanced = sum(1 for r in rows if r["skew_label"] == "Dengeli")
    unbalanced = sum(1 for r in rows if r["skew_label"] in ("Ağır ama düşük performans", "Hafif ama yüksek performans"))

    return jsonify({
        "success": True,
        "summary": {
            "total_strategies": len(rows),
            "balanced": balanced,
            "unbalanced": unbalanced,
            "plan_year": active_py.year if active_py else None,
        },
        "strategies": rows,
    })


# ─── Rapor 3: Stratejik Hizalama Sankey ──────────────────────────────────────

@app_bp.route("/raporlar/hizalama-sankey")
@login_required
def raporlar_hizalama_sankey():
    return render_template("platform/raporlar/hizalama_sankey.html")


@app_bp.route("/raporlar/api/hizalama-sankey")
@login_required
def raporlar_api_hizalama_sankey():
    """Vizyon → Strateji → Alt Strateji → Süreç akış verisi."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    active_py = get_active_plan_year_for_user(current_user)
    py_id = active_py.id if active_py else None

    # Stratejiler + ağırlık
    weights = {w.strategy_id: w.weight_raw or 0 for w in
               KVektorStrategyWeight.query.filter_by(tenant_id=tid).all()}
    total_w = sum(weights.values()) or 1

    strat_q = Strategy.query.filter_by(tenant_id=tid, is_active=True)
    if py_id:
        strat_q = strat_q.filter(or_(Strategy.plan_year_id == py_id, Strategy.plan_year_id.is_(None)))
    strategies = strat_q.order_by(Strategy.code).all()

    nodes = [{"id": "vision", "label": "Vizyon", "level": 0, "color": "#0f172a"}]
    links = []
    node_index = {"vision": 0}

    for s in strategies:
        sid = f"s-{s.id}"
        node_index[sid] = len(nodes)
        nodes.append({
            "id": sid,
            "label": f"{s.code or ''} {s.title}".strip()[:40],
            "level": 1,
            "color": "#4f46e5",
        })
        w = weights.get(s.id, 0)
        links.append({"source": 0, "target": node_index[sid], "value": round(w / total_w * 100, 1)})

        # Alt stratejiler
        for ss in s.sub_strategies:
            if not getattr(ss, "is_active", True):
                continue
            ssid = f"ss-{ss.id}"
            node_index[ssid] = len(nodes)
            nodes.append({
                "id": ssid,
                "label": f"{ss.code or ''} {ss.title}".strip()[:40],
                "level": 2,
                "color": "#8b5cf6",
            })
            links.append({"source": node_index[sid], "target": node_index[ssid], "value": 1})

            # Bağlı süreçler (process_sub_strategy_links)
            for pssl in ss.process_sub_strategy_links:
                p = pssl.process
                if not p or not p.is_active:
                    continue
                if p.tenant_id != tid:
                    continue
                if py_id and p.plan_year_id and p.plan_year_id != py_id:
                    continue
                pid = f"p-{p.id}"
                if pid not in node_index:
                    node_index[pid] = len(nodes)
                    nodes.append({
                        "id": pid,
                        "label": f"{p.code or ''} {p.name}".strip()[:40],
                        "level": 3,
                        "color": "#10b981",
                    })
                contrib = pssl.contribution_pct or 1
                links.append({"source": node_index[ssid], "target": node_index[pid], "value": round(contrib, 1)})

    return jsonify({
        "success": True,
        "summary": {
            "vision_nodes": 1,
            "strategy_nodes": len(strategies),
            "sub_strategy_nodes": sum(1 for n in nodes if n["level"] == 2),
            "process_nodes": sum(1 for n in nodes if n["level"] == 3),
            "total_links": len(links),
            "plan_year": active_py.year if active_py else None,
        },
        "nodes": nodes,
        "links": links,
    })


# ─── Rapor 4: PG Hedef Revizyon Sıklığı ──────────────────────────────────────

@app_bp.route("/raporlar/hedef-revizyon")
@login_required
def raporlar_hedef_revizyon():
    return render_template("platform/raporlar/hedef_revizyon.html")


@app_bp.route("/raporlar/api/hedef-revizyon")
@login_required
def raporlar_api_hedef_revizyon():
    """Yıl bazlı hedef override sayımı (kpi_year_configs vs ProcessKpi.target_value)."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    # Tüm plan yıllarını + kpi_year_configs sayımı (override yapılmış PG'ler)
    py_list = list_plan_years(tid)
    py_data = []
    revised_kpis = []  # tek tek revize edilen PG detayı

    for py in py_list:
        configs = KpiYearConfig.query.filter_by(plan_year_id=py.id).all()
        revised_count = 0
        for cfg in configs:
            kpi = cfg.process_kpi
            if not kpi:
                continue
            # Override aktif sayılır: target_value/weight/period vs farklıysa
            differences = []
            if cfg.target_value and cfg.target_value != (kpi.target_value or ""):
                differences.append("hedef")
            if cfg.weight is not None and cfg.weight != (kpi.weight or 0):
                differences.append("ağırlık")
            if cfg.period and cfg.period != (kpi.period or ""):
                differences.append("periyot")
            if differences:
                revised_count += 1
                revised_kpis.append({
                    "year": py.year,
                    "kpi_code": kpi.code or f"PG-{kpi.id}",
                    "kpi_name": kpi.name,
                    "diff_fields": differences,
                    "base_target": kpi.target_value,
                    "year_target": cfg.target_value,
                })

        py_data.append({
            "year": py.year,
            "status": py.status,
            "config_count": len(configs),
            "revised_count": revised_count,
        })

    return jsonify({
        "success": True,
        "summary": {
            "total_years": len(py_list),
            "total_revisions": len(revised_kpis),
            "years_with_revisions": sum(1 for x in py_data if x["revised_count"] > 0),
        },
        "by_year": py_data,
        "revised_kpis": revised_kpis[:50],
    })


# ─── Rapor 5: Departman Performans Skoru ─────────────────────────────────────

@app_bp.route("/raporlar/departman-performans")
@login_required
def raporlar_departman_performans():
    return render_template("platform/raporlar/departman_performans.html")


@app_bp.route("/raporlar/api/departman-performans")
@login_required
def raporlar_api_departman_performans():
    """Departman bazlı kullanıcı + bireysel PG + atanan görev sayısı."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    rows = db.session.query(
        User.department,
        func.count(func.distinct(User.id)).label("user_count"),
    ).filter(
        User.tenant_id == tid,
        User.is_active.is_(True),
    ).group_by(User.department).all()

    # Departman başına bireysel PG sayısı + başarı medyanı
    dept_data = []
    for dept, user_count in rows:
        users_in_dept = db.session.query(User.id).filter(
            User.tenant_id == tid, User.is_active.is_(True),
            User.department == dept,
        ).subquery()

        pg_count = db.session.query(func.count(IndividualPerformanceIndicator.id)).filter(
            IndividualPerformanceIndicator.user_id.in_(users_in_dept),
            IndividualPerformanceIndicator.is_active.is_(True),
        ).scalar() or 0

        # Önemli PG sayısı (is_important)
        important_pg = db.session.query(func.count(IndividualPerformanceIndicator.id)).filter(
            IndividualPerformanceIndicator.user_id.in_(users_in_dept),
            IndividualPerformanceIndicator.is_active.is_(True),
            IndividualPerformanceIndicator.is_important.is_(True),
        ).scalar() or 0

        # Departman içindeki kullanıcıların 2FA oranı (security proxy)
        totp_enabled = db.session.query(func.count(User.id)).filter(
            User.tenant_id == tid, User.is_active.is_(True),
            User.department == dept, User.totp_enabled.is_(True),
        ).scalar() or 0

        dept_data.append({
            "department": dept or "(belirsiz)",
            "user_count": user_count,
            "pg_count": pg_count,
            "pg_per_user": round(pg_count / user_count, 2) if user_count else 0,
            "important_pg": important_pg,
            "totp_rate": round(totp_enabled / user_count * 100, 1) if user_count else 0,
        })

    dept_data.sort(key=lambda x: -x["user_count"])

    return jsonify({
        "success": True,
        "summary": {
            "total_departments": len(dept_data),
            "total_users": sum(d["user_count"] for d in dept_data),
            "total_pgs": sum(d["pg_count"] for d in dept_data),
        },
        "departments": dept_data,
    })


# ─── Rapor 6: Yönetici Liderlik Skoru ────────────────────────────────────────

@app_bp.route("/raporlar/yonetici-liderlik")
@login_required
def raporlar_yonetici_liderlik():
    return render_template("platform/raporlar/yonetici_liderlik.html")


@app_bp.route("/raporlar/api/yonetici-liderlik")
@login_required
def raporlar_api_yonetici_liderlik():
    """Süreç liderlerinin liderliğindeki süreçlerin ortalama performans skoru."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    active_py = get_active_plan_year_for_user(current_user)
    py_id = active_py.id if active_py else None

    # Süreç liderleri (M:N tablosu)
    rows = db.session.query(
        process_leaders.c.user_id,
        process_leaders.c.process_id,
    ).join(Process, Process.id == process_leaders.c.process_id).filter(
        Process.tenant_id == tid,
        Process.is_active.is_(True),
    ).all()

    if not rows:
        # Fallback: hiç process_leaders yoksa user.role==yonetici olanları göster
        managers = User.query.filter_by(
            tenant_id=tid, is_active=True,
        ).join(User.role).filter(
            (User.role.has(name="yonetici")) | (User.role.has(name="tenant_admin"))
        ).all()
        return jsonify({
            "success": True,
            "summary": {
                "total_leaders": len(managers),
                "has_leader_data": False,
                "note": "Süreç-lider eşleştirmesi bulunamadı. Rol bazlı yönetici listesi gösterildi.",
            },
            "leaders": [{
                "user_id": u.id,
                "name": f"{u.first_name or ''} {u.last_name or ''}".strip() or u.email,
                "email": u.email,
                "department": u.department,
                "led_process_count": 0,
                "avg_process_score": None,
            } for u in managers[:25]],
        })

    # Skor motorundan süreç skorları
    try:
        today = _date.today()
        score_year = active_py.year if active_py else today.year
        proc_scores, _ = compute_process_scores_internal(
            tid, score_year, today, persist_pg_scores=False, plan_year=active_py
        )
    except Exception:
        proc_scores = {}

    # Lider → bağlı süreçler agregasyonu
    leader_map = defaultdict(list)
    for uid, pid in rows:
        leader_map[uid].append(pid)

    leaders = []
    for uid, pids in leader_map.items():
        u = User.query.get(uid)
        if not u or not u.is_active:
            continue
        scores = [proc_scores.get(pid) for pid in pids if proc_scores.get(pid) is not None]
        avg = round(sum(scores) / len(scores), 1) if scores else None
        leaders.append({
            "user_id": uid,
            "name": f"{u.first_name or ''} {u.last_name or ''}".strip() or u.email,
            "email": u.email,
            "department": u.department,
            "led_process_count": len(pids),
            "avg_process_score": avg,
        })
    leaders.sort(key=lambda x: -(x["avg_process_score"] or 0))

    return jsonify({
        "success": True,
        "summary": {
            "total_leaders": len(leaders),
            "has_leader_data": True,
            "with_score": sum(1 for l in leaders if l["avg_process_score"] is not None),
            "avg_score_overall": round(
                sum(l["avg_process_score"] for l in leaders if l["avg_process_score"]) /
                max(sum(1 for l in leaders if l["avg_process_score"]), 1), 1
            ),
        },
        "leaders": leaders,
    })


# ─── Rapor 7: Initiative Portföy Bubble ──────────────────────────────────────

@app_bp.route("/raporlar/initiative-bubble")
@login_required
def raporlar_initiative_bubble():
    return render_template("platform/raporlar/initiative_bubble.html")


@app_bp.route("/raporlar/api/initiative-bubble")
@login_required
def raporlar_api_initiative_bubble():
    """Initiative portföyü: bütçe × ilerleme × öncelik."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    inits = Initiative.query.filter_by(tenant_id=tid, is_active=True).all()
    PRIO_COLOR = {
        "critical": "#dc2626", "high": "#f59e0b",
        "medium": "#0ea5e9", "low": "#94a3b8",
    }
    STATUS_OPACITY = {
        "completed": 0.55, "cancelled": 0.25, "on_hold": 0.4,
        "in_progress": 1.0, "planned": 0.7,
    }
    bubbles = []
    for i in inits:
        budget = float(i.budget_total or 0)
        spent = float(i.budget_spent or 0)
        bubbles.append({
            "id": i.id,
            "code": i.code or f"INI-{i.id}",
            "name": i.name,
            "status": i.status,
            "priority": i.priority,
            "color": PRIO_COLOR.get(i.priority, "#64748b"),
            "opacity": STATUS_OPACITY.get(i.status, 0.7),
            "budget_total": budget,
            "budget_spent": spent,
            "budget_usage_pct": round(spent / budget * 100, 1) if budget else 0,
            "progress_pct": i.progress_pct or 0,
            "start_year": i.start_year,
            "end_year": i.end_year,
            "duration": (i.end_year or 0) - (i.start_year or 0) + 1,
        })

    # Özet
    by_status = defaultdict(int)
    by_priority = defaultdict(int)
    for b in bubbles:
        by_status[b["status"]] += 1
        by_priority[b["priority"]] += 1

    return jsonify({
        "success": True,
        "summary": {
            "total": len(bubbles),
            "total_budget": sum(b["budget_total"] for b in bubbles),
            "total_spent": sum(b["budget_spent"] for b in bubbles),
            "by_status": dict(by_status),
            "by_priority": dict(by_priority),
            "avg_progress": round(sum(b["progress_pct"] for b in bubbles) / max(len(bubbles), 1), 1),
        },
        "bubbles": bubbles,
    })


# ─── Rapor 8: Operasyonel Sabah Özeti+ ──────────────────────────────────────

@app_bp.route("/raporlar/sabah-ozeti")
@login_required
def raporlar_sabah_ozeti():
    return render_template("platform/raporlar/sabah_ozeti.html")


@app_bp.route("/raporlar/api/sabah-ozeti")
@login_required
def raporlar_api_sabah_ozeti():
    """Bugünün ve son 7 günün operasyonel özeti."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    today = _date.today()
    week_ago = today - timedelta(days=7)
    week_ahead = today + timedelta(days=7)

    # Tenant süreçleri
    proc_ids_subq = db.session.query(Process.id).filter(
        Process.tenant_id == tid, Process.is_active.is_(True),
    ).subquery()

    # Faaliyet sayımları
    overdue = ProcessActivity.query.filter(
        ProcessActivity.process_id.in_(proc_ids_subq),
        ProcessActivity.is_active.is_(True),
        ProcessActivity.end_at < datetime.utcnow(),
        ProcessActivity.status != "Tamamlandı",
    ).count()

    today_due = ProcessActivity.query.filter(
        ProcessActivity.process_id.in_(proc_ids_subq),
        ProcessActivity.is_active.is_(True),
        func.date(ProcessActivity.end_at) == today,
    ).count()

    upcoming = ProcessActivity.query.filter(
        ProcessActivity.process_id.in_(proc_ids_subq),
        ProcessActivity.is_active.is_(True),
        ProcessActivity.end_at.between(today, week_ahead),
    ).count()

    completed_week = ProcessActivity.query.filter(
        ProcessActivity.process_id.in_(proc_ids_subq),
        ProcessActivity.status == "Tamamlandı",
        ProcessActivity.completed_at.between(week_ago, today + timedelta(days=1)),
    ).count()

    # KPI ölçüm akışı
    kpi_ids_subq = db.session.query(ProcessKpi.id).filter(
        ProcessKpi.process_id.in_(proc_ids_subq),
        ProcessKpi.is_active.is_(True),
    ).subquery()

    measurements_today = KpiData.query.filter(
        KpiData.process_kpi_id.in_(kpi_ids_subq),
        KpiData.is_active.is_(True),
        func.date(KpiData.created_at) == today,
    ).count()

    measurements_week = KpiData.query.filter(
        KpiData.process_kpi_id.in_(kpi_ids_subq),
        KpiData.is_active.is_(True),
        KpiData.created_at >= week_ago,
    ).count()

    # En son 5 ölçüm girişi (recent activity)
    recent_data = db.session.query(
        KpiData.id, KpiData.created_at, KpiData.actual_value,
        ProcessKpi.name.label("kpi_name"),
        ProcessKpi.code.label("kpi_code"),
        User.first_name, User.last_name, User.email,
    ).join(ProcessKpi, ProcessKpi.id == KpiData.process_kpi_id).join(
        Process, Process.id == ProcessKpi.process_id
    ).join(User, User.id == KpiData.user_id, isouter=True).filter(
        Process.tenant_id == tid,
        KpiData.is_active.is_(True),
    ).order_by(KpiData.created_at.desc()).limit(8).all()

    recent_entries = [{
        "kpi_code": r.kpi_code,
        "kpi_name": r.kpi_name,
        "value": r.actual_value,
        "when": r.created_at.isoformat() if r.created_at else None,
        "who": (f"{r.first_name or ''} {r.last_name or ''}").strip() or r.email or "—",
    } for r in recent_data]

    return jsonify({
        "success": True,
        "today": today.isoformat(),
        "metrics": {
            "overdue_activities": overdue,
            "today_due_activities": today_due,
            "upcoming_7d_activities": upcoming,
            "completed_7d_activities": completed_week,
            "measurements_today": measurements_today,
            "measurements_7d": measurements_week,
        },
        "recent_entries": recent_entries,
    })


# ─── Rapor 9: Yıllar Arası Evrim Filmi ──────────────────────────────────────

@app_bp.route("/raporlar/evrim-filmi")
@login_required
def raporlar_evrim_filmi():
    return render_template("platform/raporlar/evrim_filmi.html")


@app_bp.route("/raporlar/api/evrim-filmi")
@login_required
def raporlar_api_evrim_filmi():
    """Yıllar boyunca strateji ağacının evrimi — her yıl için snapshot."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    py_list = list_plan_years(tid)
    if not py_list:
        return jsonify({"success": True, "years": [], "summary": {"total_years": 0}})

    snapshots = []
    for py in sorted(py_list, key=lambda x: x.year):
        # Bu plan yılındaki tüm strateji-alt-süreç-PG
        strategies = Strategy.query.filter_by(
            tenant_id=tid, plan_year_id=py.id, is_active=True
        ).order_by(Strategy.code).all()

        nodes = []
        for s in strategies:
            sub_list = []
            for ss in s.sub_strategies:
                if not getattr(ss, "is_active", True):
                    continue
                # Bu alt-stratejiye bağlı süreçler (process_sub_strategy_links üzerinden)
                proc_codes = []
                for pssl in ss.process_sub_strategy_links:
                    p = pssl.process
                    if p and p.is_active and (not p.plan_year_id or p.plan_year_id == py.id):
                        proc_codes.append(p.code or f"P-{p.id}")
                sub_list.append({
                    "code": ss.code or "",
                    "title": ss.title or "",
                    "process_codes": list(set(proc_codes)),
                })
            nodes.append({
                "code": s.code or "",
                "title": s.title or "",
                "id": s.id,
                "sub_strategies": sub_list,
            })

        # Yılın özet sayıları
        proc_count = Process.query.filter_by(
            tenant_id=tid, plan_year_id=py.id, is_active=True
        ).count()
        kpi_count = ProcessKpi.query.join(Process).filter(
            Process.tenant_id == tid,
            Process.plan_year_id == py.id,
            ProcessKpi.is_active.is_(True),
            Process.is_active.is_(True),
        ).count()
        meas_count = db.session.query(func.count(KpiData.id)).filter(
            KpiData.year == py.year,
            KpiData.process_kpi_id.in_(
                db.session.query(ProcessKpi.id).join(Process).filter(
                    Process.tenant_id == tid, Process.plan_year_id == py.id
                )
            ),
        ).scalar() or 0

        snapshots.append({
            "year": py.year,
            "status": py.status,
            "strategies": nodes,
            "stats": {
                "strategy_count": len(strategies),
                "sub_count": sum(len(s["sub_strategies"]) for s in nodes),
                "process_count": proc_count,
                "kpi_count": kpi_count,
                "measurement_count": meas_count,
            },
        })

    return jsonify({
        "success": True,
        "summary": {
            "total_years": len(snapshots),
            "year_range": f"{snapshots[0]['year']}–{snapshots[-1]['year']}" if snapshots else "—",
        },
        "years": snapshots,
    })


# ─── Rapor 10: AI Yıl Sonu Sunum Üretici ────────────────────────────────────

@app_bp.route("/raporlar/ai-sunum")
@login_required
def raporlar_ai_sunum():
    return render_template("platform/raporlar/ai_sunum.html")


@app_bp.route("/raporlar/api/ai-sunum/preview")
@login_required
def raporlar_api_ai_sunum_preview():
    """Sunumun veri özetini ve üretilecek slayt başlıklarını döner (preview)."""
    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    active_py = get_active_plan_year_for_user(current_user)
    tenant = db.session.get(Tenant, tid)

    strat_count = Strategy.query.filter_by(
        tenant_id=tid, is_active=True,
        plan_year_id=active_py.id if active_py else None,
    ).count() if active_py else Strategy.query.filter_by(tenant_id=tid, is_active=True).count()

    proc_count = Process.query.filter_by(
        tenant_id=tid, is_active=True,
        plan_year_id=active_py.id if active_py else None,
    ).count() if active_py else Process.query.filter_by(tenant_id=tid, is_active=True).count()

    init_count = Initiative.query.filter_by(tenant_id=tid, is_active=True).count()
    user_count = User.query.filter_by(tenant_id=tid, is_active=True).count()
    meas_count = db.session.query(func.count(KpiData.id)).join(
        ProcessKpi
    ).join(Process).filter(Process.tenant_id == tid).scalar() or 0

    slides = [
        {"no": 1, "title": "Kapak", "content": f"{tenant.name if tenant else 'Kurum'} — {active_py.year if active_py else 'Yıllık'} Stratejik Plan Sunumu"},
        {"no": 2, "title": "Yönetici Özeti", "content": "AI tarafından üretilen 3-cümle özet"},
        {"no": 3, "title": "Vizyon ve Misyon", "content": "Kurumsal kimlik alanları"},
        {"no": 4, "title": "Stratejik Yıl Özeti", "content": f"{strat_count} strateji · {proc_count} süreç · {init_count} initiative"},
        {"no": 5, "title": "Strateji Ağacı", "content": "Ana stratejiler + ağırlık dağılımı"},
        {"no": 6, "title": "Süreç Sağlık Heatmap", "content": f"{proc_count} süreç performans özeti"},
        {"no": 7, "title": "K-Vektör Analizi", "content": "Ağırlık × performans çarpıklık tespiti"},
        {"no": 8, "title": "Initiative Portföyü", "content": f"{init_count} girişimin bütçe + ilerleme durumu"},
        {"no": 9, "title": "Yıllar Arası Değişim", "content": "Önceki yılla karşılaştırma — strateji, süreç, PG farkları"},
        {"no": 10, "title": "İnsan Kaynakları", "content": f"{user_count} çalışan + departman dağılımı"},
        {"no": 11, "title": "Risk Haritası", "content": "Top 10 risk + olasılık × etki matrisi"},
        {"no": 12, "title": "ESG & Sürdürülebilirlik", "content": "E/S/G metrikleri + SDG katkı"},
        {"no": 13, "title": "Veri Kalitesi", "content": f"{meas_count:,} toplam ölçüm · doluluk oranları"},
        {"no": 14, "title": "Önümüzdeki Dönem Yol Haritası", "content": "Yeni initiative önerileri + öncelikler"},
        {"no": 15, "title": "Kapanış", "content": "Teşekkür + Q&A"},
    ]

    return jsonify({
        "success": True,
        "preview": {
            "tenant_name": tenant.name if tenant else "—",
            "plan_year": active_py.year if active_py else None,
            "data_points": {
                "strategies": strat_count,
                "processes": proc_count,
                "initiatives": init_count,
                "users": user_count,
                "kpi_measurements": meas_count,
            },
            "slides": slides,
            "estimated_size_kb": 280,
            "ai_tokens_estimate": 4500,
        },
    })


@app_bp.route("/raporlar/api/ai-sunum/generate", methods=["GET", "POST"])
@login_required
def raporlar_api_ai_sunum_generate():
    """Gerçek PowerPoint dosyasını üretir + indirme URL'i döner."""
    from flask import send_file
    import io
    import os
    from datetime import datetime as _dt

    tid = current_user.tenant_id
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return jsonify({
            "success": False,
            "message": "python-pptx kurulu değil. Kurulum: pip install python-pptx",
        }), 500

    active_py = get_active_plan_year_for_user(current_user)
    tenant = db.session.get(Tenant, tid)

    # Veri özetleri
    tenant_name = tenant.name if tenant else "Kurum"
    year_label = str(active_py.year) if active_py else _dt.utcnow().year

    strategies = []
    if active_py:
        strat_q = Strategy.query.filter_by(
            tenant_id=tid, is_active=True, plan_year_id=active_py.id,
        ).order_by(Strategy.code).all()
        weights = {w.strategy_id: w.weight_raw for w in
                   KVektorStrategyWeight.query.filter_by(tenant_id=tid).all()}
        total_w = sum(weights.values()) or 1
        for s in strat_q:
            strategies.append({
                "code": s.code, "title": s.title,
                "weight_pct": round((weights.get(s.id, 0) or 0) / total_w * 100, 1),
            })

    proc_count = Process.query.filter_by(
        tenant_id=tid, is_active=True,
        plan_year_id=active_py.id if active_py else None,
    ).count()
    kpi_count = ProcessKpi.query.join(Process).filter(
        Process.tenant_id == tid,
        Process.plan_year_id == (active_py.id if active_py else None),
        ProcessKpi.is_active.is_(True),
    ).count()
    initiatives = Initiative.query.filter_by(tenant_id=tid, is_active=True).all()
    user_count = User.query.filter_by(tenant_id=tid, is_active=True).count()
    meas_count = db.session.query(func.count(KpiData.id)).join(ProcessKpi).join(
        Process
    ).filter(Process.tenant_id == tid).scalar() or 0

    # AI özeti — LLM gateway ile (yoksa fallback şablon)
    ai_summary = (
        f"{tenant_name} {year_label} dönemini {len(strategies)} ana strateji, "
        f"{proc_count} aktif süreç ve {kpi_count} performans göstergesi ile yönetti. "
        f"Yıl içinde {len(initiatives)} stratejik girişim portföyünde takip edildi ve "
        f"{user_count} çalışan tarafından toplam {meas_count:,} ölçüm verisi sisteme girildi. "
        "Önümüzdeki dönemde stratejik öncelikler ve yapısal iyileştirmeler için yol haritası bu sunumda yer almaktadır."
    )
    try:
        from app.services.llm_gateway import call_llm  # opsiyonel
        prompt = (
            f"Türkçe, 4 cümlelik yönetici özeti yaz. Kurum: {tenant_name}, Yıl: {year_label}. "
            f"Veri: {len(strategies)} strateji, {proc_count} süreç, {kpi_count} PG, "
            f"{len(initiatives)} initiative, {user_count} çalışan, {meas_count} ölçüm. "
            "Resmî ama özlü ton kullan."
        )
        result = call_llm(prompt=prompt, tenant_id=tid, user_id=current_user.id,
                          endpoint="ai_yil_sonu_sunum", max_output_tokens=400)
        if result and result.get("text"):
            ai_summary = result["text"].strip()
    except Exception as e:
        current_app.logger.info(f"[ai-sunum] LLM yok/başarısız ({e}), fallback metin kullanıldı")

    # PowerPoint oluştur
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(0x4f, 0x46, 0xe5)
    DARK = RGBColor(0x0f, 0x17, 0x2a)
    GRAY = RGBColor(0x64, 0x74, 0x8b)
    LIGHT_BG = RGBColor(0xf8, 0xfa, 0xfc)

    def add_title_slide(prs, title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Title
        tx = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = DARK
        # Subtitle
        sub = slide.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(12), Inches(1))
        sub.text_frame.text = subtitle
        sub.text_frame.paragraphs[0].font.size = Pt(20)
        sub.text_frame.paragraphs[0].font.color.rgb = GRAY
        return slide

    def add_content_slide(prs, slide_no, title, body_text, body_bullets=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Header bar
        from pptx.shapes.autoshape import Shape  # noqa
        # Title
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BLUE
        # Body
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
        bf = body.text_frame
        bf.word_wrap = True
        if body_text:
            p0 = bf.paragraphs[0]
            p0.text = body_text
            p0.font.size = Pt(16)
            p0.font.color.rgb = DARK
        if body_bullets:
            for bul in body_bullets:
                bp = bf.add_paragraph()
                bp.text = "• " + str(bul)
                bp.font.size = Pt(14)
                bp.font.color.rgb = GRAY
                bp.space_before = Pt(6)
        # Footer
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.3))
        ft.text_frame.text = f"{tenant_name} · {year_label} · Slayt {slide_no}"
        ft.text_frame.paragraphs[0].font.size = Pt(10)
        ft.text_frame.paragraphs[0].font.color.rgb = GRAY
        return slide

    # Slaytları üret
    add_title_slide(prs,
        f"{tenant_name}",
        f"{year_label} Stratejik Plan Sunumu · {_dt.utcnow().strftime('%d %B %Y')}")

    add_content_slide(prs, 2, "Yönetici Özeti", ai_summary)

    add_content_slide(prs, 3, f"Stratejik Yıl Özeti — {year_label}",
        "Tüm sistem aşağıdaki ana boyutlarda izlendi:",
        body_bullets=[
            f"{len(strategies)} ana strateji",
            f"{proc_count} aktif süreç",
            f"{kpi_count} performans göstergesi",
            f"{len(initiatives)} stratejik girişim (initiative)",
            f"{user_count} aktif çalışan",
            f"{meas_count:,} toplam KPI ölçümü",
        ])

    add_content_slide(prs, 4, "Ana Stratejiler ve K-Vektör Ağırlıkları",
        None, body_bullets=[
            f"{s['code']} — {s['title']} (Ağırlık: %{s['weight_pct']})"
            for s in strategies[:12]
        ] or ["Aktif plan yılında strateji bulunamadı."])

    add_content_slide(prs, 5, "Initiative Portföyü",
        f"Aktif ve tamamlanmış toplam {len(initiatives)} girişim.",
        body_bullets=[
            f"{i.code or 'INI-'+str(i.id)} — {i.name} (durum: {i.status}, ilerleme: %{i.progress_pct or 0})"
            for i in initiatives[:10]
        ])

    add_content_slide(prs, 6, "Veri Akışı ve Ölçüm Hacmi",
        f"{tenant_name} sistemine girilen toplam {meas_count:,} ölçüm verisi.",
        body_bullets=[
            f"Aktif süreç: {proc_count}",
            f"Aktif PG: {kpi_count}",
            f"Ortalama ölçüm/PG: {round(meas_count / max(kpi_count, 1), 1)}",
            "Veri girişi düzenli sürdürüldü — kalite skoru raporu detay sağlar.",
        ])

    add_content_slide(prs, 7, "Önümüzdeki Dönem — Öneriler",
        "Bu sunumun arka planındaki Kokpitim verisinden çıkan öneriler:",
        body_bullets=[
            "Stratejik öncelikleri K-Vektör çarpıklık raporundan gözden geçir",
            "Yıllar arası karşılaştırma sayfasından evrim haritasını incele",
            "Veri kalitesi raporundaki eksik PG'leri planla",
            "Initiative portföyünde geri kalanlara odak",
        ])

    add_content_slide(prs, 8, "Teşekkürler",
        f"{tenant_name} ekibinin katkılarıyla hazırlandı.",
        body_bullets=["Sorular ve detay için: kokpitim.com"])

    # Bellek üzerinden gönder
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = f"{tenant_name.replace(' ', '_')}_{year_label}_yil_sonu_sunum.pptx"
    return send_file(
        buf,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=filename,
    )


# ═══════════════════════════════════════════════════════════════════════════
# FAZ 1 — 15 HIZLI KAZANC RAPORU (PLAN-TUM-RAPORLAR.md Faz 1)
# ═══════════════════════════════════════════════════════════════════════════

def _tid_or_none():
    return current_user.tenant_id if current_user.tenant_id else None


# ─── ST-01: Stratejik Hiyerarşi Sunburst ─────────────────────────────────────

@app_bp.route("/raporlar/sunburst")
@login_required
def raporlar_sunburst():
    return render_template("platform/raporlar/sunburst.html")


@app_bp.route("/raporlar/api/sunburst")
@login_required
def raporlar_api_sunburst():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    active_py = get_active_plan_year_for_user(current_user)
    py_id = active_py.id if active_py else None

    strat_q = Strategy.query.filter_by(tenant_id=tid, is_active=True)
    if py_id:
        strat_q = strat_q.filter(or_(Strategy.plan_year_id == py_id, Strategy.plan_year_id.is_(None)))
    strategies = strat_q.order_by(Strategy.code).all()

    weights = {w.strategy_id: w.weight_raw or 0 for w in
               KVektorStrategyWeight.query.filter_by(tenant_id=tid).all()}
    total_w = sum(weights.values()) or 1

    # Tree yapısı: vizyon → strateji → alt → süreç
    tree = {"name": "Vizyon", "value": 100, "children": []}
    for s in strategies:
        s_node = {
            "name": (s.code or "?") + " " + (s.title or "")[:30],
            "value": round((weights.get(s.id, 0) / total_w) * 100, 1),
            "children": [],
        }
        sub_count = max(len([ss for ss in s.sub_strategies if getattr(ss, "is_active", True)]), 1)
        for ss in s.sub_strategies:
            if not getattr(ss, "is_active", True):
                continue
            ss_node = {
                "name": (ss.code or "?") + " " + (ss.title or "")[:25],
                "value": round(s_node["value"] / sub_count, 1),
                "children": [],
            }
            for pssl in ss.process_sub_strategy_links:
                p = pssl.process
                if not p or not p.is_active:
                    continue
                if py_id and p.plan_year_id and p.plan_year_id != py_id:
                    continue
                ss_node["children"].append({
                    "name": (p.code or "?") + " " + (p.name or "")[:25],
                    "value": round((pssl.contribution_pct or 10) / 5, 1),
                })
            s_node["children"].append(ss_node)
        tree["children"].append(s_node)

    return jsonify({"success": True, "tree": tree, "summary": {
        "strategies": len(strategies),
        "plan_year": active_py.year if active_py else None,
    }})


# ─── ST-06: VRIO Portföy 4-Köşe ─────────────────────────────────────────────

@app_bp.route("/raporlar/vrio-portfoy")
@login_required
def raporlar_vrio_portfoy():
    return render_template("platform/raporlar/vrio_portfoy.html")


@app_bp.route("/raporlar/api/vrio-portfoy")
@login_required
def raporlar_api_vrio_portfoy():
    from app.models.strategy_frameworks import VRIOResource
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    resources = VRIOResource.query.filter_by(tenant_id=tid, is_active=True).order_by(VRIOResource.name).all()

    buckets = {
        "Rekabetçi Dezavantaj": [],
        "Rekabet Paritesi": [],
        "Geçici Rekabet Avantajı": [],
        "Kullanılmayan Avantaj": [],
        "Sürdürülebilir Rekabet Avantajı": [],
    }
    for r in resources:
        buckets[r.competitive_label].append({
            "id": r.id, "name": r.name, "category": r.category, "note": r.note,
            "v": r.is_valuable, "r": r.is_rare, "i": r.is_inimitable, "o": r.is_organized,
        })

    return jsonify({"success": True, "summary": {
        "total": len(resources),
        "sustainable": len(buckets["Sürdürülebilir Rekabet Avantajı"]),
        "temporary": len(buckets["Geçici Rekabet Avantajı"]),
        "parity": len(buckets["Rekabet Paritesi"]),
        "unused": len(buckets["Kullanılmayan Avantaj"]),
        "disadvantage": len(buckets["Rekabetçi Dezavantaj"]),
    }, "buckets": buckets})


# ─── ST-09: OKR Cascade Görseli ─────────────────────────────────────────────

@app_bp.route("/raporlar/okr-cascade")
@login_required
def raporlar_okr_cascade():
    return render_template("platform/raporlar/okr_cascade.html")


@app_bp.route("/raporlar/api/okr-cascade")
@login_required
def raporlar_api_okr_cascade():
    from app.models.okr import OkrObjective, OkrKeyResult
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    active_py = get_active_plan_year_for_user(current_user)
    py_id = active_py.id if active_py else None

    q = OkrObjective.query.filter_by(tenant_id=tid, is_active=True)
    if py_id:
        q = q.filter_by(plan_year_id=py_id)
    objectives = q.order_by(OkrObjective.quarter.nullsfirst(), OkrObjective.order_no).all()

    items = []
    for obj in objectives:
        krs = []
        for kr in obj.key_results:
            if not kr.is_active:
                continue
            krs.append({
                "id": kr.id, "title": kr.title, "metric": kr.metric,
                "start": kr.start_value, "target": kr.target_value, "current": kr.current_value,
                "progress_pct": kr.progress_pct,
                "linked_kpi_id": kr.linked_process_kpi_id,
            })
        strat_name = obj.linked_strategy.title if obj.linked_strategy else None
        sub_name = obj.linked_sub_strategy.title if obj.linked_sub_strategy else None
        items.append({
            "id": obj.id, "title": obj.title, "quarter": obj.quarter,
            "owner": obj.owner, "strategy": strat_name, "sub_strategy": sub_name,
            "kr_count": len(krs), "key_results": krs,
            "avg_progress": round(sum((k["progress_pct"] or 0) for k in krs) / max(len(krs), 1), 1),
        })

    return jsonify({"success": True, "summary": {
        "total_objectives": len(items),
        "total_krs": sum(o["kr_count"] for o in items),
        "avg_completion": round(sum(o["avg_progress"] for o in items) / max(len(items), 1), 1),
        "plan_year": active_py.year if active_py else None,
    }, "objectives": items})


# ─── ST-11: Initiative Roadmap Gantt ────────────────────────────────────────

@app_bp.route("/raporlar/initiative-roadmap")
@login_required
def raporlar_initiative_roadmap():
    return render_template("platform/raporlar/initiative_roadmap.html")


@app_bp.route("/raporlar/api/initiative-roadmap")
@login_required
def raporlar_api_initiative_roadmap():
    from app.models.initiative import InitiativeMilestone
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    inits = Initiative.query.filter_by(tenant_id=tid, is_active=True).order_by(Initiative.start_year).all()

    items = []
    for i in inits:
        milestones = []
        for m in i.milestones:
            milestones.append({
                "id": m.id, "name": m.name,
                "target_date": m.target_date.isoformat() if m.target_date else None,
                "completed_date": m.completed_date.isoformat() if m.completed_date else None,
                "status": m.status,
            })
        items.append({
            "id": i.id, "code": i.code, "name": i.name,
            "start_year": i.start_year, "end_year": i.end_year,
            "status": i.status, "priority": i.priority,
            "progress_pct": i.progress_pct,
            "milestones": milestones,
        })
    return jsonify({"success": True, "summary": {
        "total": len(items),
        "year_range": f"{min((i['start_year'] for i in items), default=0)}–{max((i['end_year'] for i in items), default=0)}" if items else "—",
        "total_milestones": sum(len(i["milestones"]) for i in items),
    }, "initiatives": items})


# ─── OP-04: 7 Muda Waste Analizi ───────────────────────────────────────────

@app_bp.route("/raporlar/muda-analizi")
@login_required
def raporlar_muda_analizi():
    return render_template("platform/raporlar/muda_analizi.html")


@app_bp.route("/raporlar/api/muda-analizi")
@login_required
def raporlar_api_muda_analizi():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    try:
        from app.services.muda_analyzer import MudaAnalyzerService
    except Exception as e:
        return jsonify({"success": False, "message": f"Muda servisi yüklenemedi: {e}"}), 500

    procs = Process.query.filter_by(tenant_id=tid, is_active=True).all()
    process_findings = []
    muda_counter = defaultdict(int)
    for p in procs[:50]:  # ilk 50 süreçle sınırla
        try:
            findings = MudaAnalyzerService.analyze_process_inefficiency(p.id, tid)
            if findings:
                for f in findings:
                    muda_counter[f.get("muda_type", "other")] += 1
                process_findings.append({
                    "process_id": p.id, "code": p.code, "name": p.name,
                    "findings_count": len(findings),
                    "findings": findings[:5],
                })
        except Exception:
            continue

    # 7 muda tipinin Türkçe etiketleri
    labels = {
        "overproduction": "Aşırı Üretim",
        "waiting": "Bekleme",
        "transport": "Taşıma/Nakliye",
        "overprocessing": "Aşırı İşlem",
        "inventory": "Stok",
        "motion": "Hareket",
        "defects": "Kusurlar",
    }
    muda_summary = [{"key": k, "label": labels.get(k, k), "count": v}
                    for k, v in sorted(muda_counter.items(), key=lambda x: -x[1])]

    return jsonify({"success": True, "summary": {
        "total_processes_analyzed": len(procs),
        "processes_with_findings": len(process_findings),
        "total_findings": sum(muda_counter.values()),
    }, "by_muda": muda_summary, "by_process": process_findings})


# ─── OP-15: CMMI Olgunluk Heatmap ──────────────────────────────────────────

@app_bp.route("/raporlar/cmmi-heatmap")
@login_required
def raporlar_cmmi_heatmap():
    return render_template("platform/raporlar/cmmi_heatmap.html")


@app_bp.route("/raporlar/api/cmmi-heatmap")
@login_required
def raporlar_api_cmmi_heatmap():
    from app.models.k_radar_domain import ProcessMaturity
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    # Her süreç için en son maturity kaydı
    rows = db.session.query(
        Process.id, Process.code, Process.name,
        func.max(ProcessMaturity.maturity_level).label("level"),
    ).join(ProcessMaturity, ProcessMaturity.process_id == Process.id).filter(
        Process.tenant_id == tid, Process.is_active.is_(True),
        ProcessMaturity.is_active.is_(True),
    ).group_by(Process.id, Process.code, Process.name).all()

    processes = [{"id": r[0], "code": r[1], "name": r[2], "level": int(r[3])} for r in rows]

    # Seviye dağılımı
    distribution = defaultdict(int)
    for p in processes:
        distribution[p["level"]] += 1

    avg_level = round(sum(p["level"] for p in processes) / max(len(processes), 1), 2)

    level_labels = {1: "Initial", 2: "Managed", 3: "Defined", 4: "Quantitatively Managed", 5: "Optimizing"}
    distribution_list = [{"level": k, "label": level_labels[k], "count": distribution.get(k, 0)}
                         for k in [1, 2, 3, 4, 5]]

    return jsonify({"success": True, "summary": {
        "total_processes": len(processes),
        "avg_level": avg_level,
        "level_5_count": distribution.get(5, 0),
    }, "distribution": distribution_list, "processes": processes})


# ─── OP-17: Operasyonel İstatistik Sayfası (süreç bazlı) ───────────────────

@app_bp.route("/raporlar/operasyon-istatistik")
@login_required
def raporlar_op_istatistik():
    return render_template("platform/raporlar/op_istatistik.html")


@app_bp.route("/raporlar/api/operasyon-istatistik")
@login_required
def raporlar_api_op_istatistik():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    proc_ids_subq = db.session.query(Process.id).filter(
        Process.tenant_id == tid, Process.is_active.is_(True),
    ).subquery()

    # Süreç başına: PG sayısı, faaliyet sayısı, son ölçüm tarihi
    rows = db.session.query(
        Process.id, Process.code, Process.name, Process.status,
        func.count(func.distinct(ProcessKpi.id)).label("kpi_count"),
        func.count(func.distinct(ProcessActivity.id)).label("activity_count"),
    ).outerjoin(ProcessKpi, and_(ProcessKpi.process_id == Process.id, ProcessKpi.is_active.is_(True))
    ).outerjoin(ProcessActivity, and_(ProcessActivity.process_id == Process.id, ProcessActivity.is_active.is_(True))
    ).filter(Process.tenant_id == tid, Process.is_active.is_(True)
    ).group_by(Process.id, Process.code, Process.name, Process.status).order_by(Process.code).all()

    items = [{
        "id": r[0], "code": r[1], "name": r[2], "status": r[3],
        "kpi_count": r[4], "activity_count": r[5],
    } for r in rows]

    return jsonify({"success": True, "summary": {
        "total_processes": len(items),
        "total_kpis": sum(i["kpi_count"] for i in items),
        "total_activities": sum(i["activity_count"] for i in items),
    }, "processes": items})


# ─── FN-05: ROI per Strategy ────────────────────────────────────────────────

@app_bp.route("/raporlar/roi-per-strategy")
@login_required
def raporlar_roi_strategy():
    return render_template("platform/raporlar/roi_strategy.html")


@app_bp.route("/raporlar/api/roi-per-strategy")
@login_required
def raporlar_api_roi_strategy():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    active_py = get_active_plan_year_for_user(current_user)
    py_id = active_py.id if active_py else None

    strat_q = Strategy.query.filter_by(tenant_id=tid, is_active=True)
    if py_id:
        strat_q = strat_q.filter(or_(Strategy.plan_year_id == py_id, Strategy.plan_year_id.is_(None)))
    strategies = strat_q.order_by(Strategy.code).all()

    # Skor motorundan strateji skorları
    try:
        today = _date.today()
        score_year = active_py.year if active_py else today.year
        proc_scores, _ = compute_process_scores_internal(tid, score_year, today, persist_pg_scores=False, plan_year=active_py)
    except Exception:
        proc_scores = {}

    rows = []
    for s in strategies:
        # Bu stratejiye bağlı initiative bütçesi (manual or implicit by year matching)
        budget = db.session.query(func.coalesce(func.sum(Initiative.budget_total), 0)).filter(
            Initiative.tenant_id == tid, Initiative.is_active.is_(True),
            Initiative.strategy_id == s.id,
        ).scalar() or 0
        spent = db.session.query(func.coalesce(func.sum(Initiative.budget_spent), 0)).filter(
            Initiative.tenant_id == tid, Initiative.is_active.is_(True),
            Initiative.strategy_id == s.id,
        ).scalar() or 0

        # Strateji skoru: bağlı süreçlerin ortalaması
        sub_ids = [ss.id for ss in s.sub_strategies if getattr(ss, "is_active", True)]
        related_proc_ids = set()
        if sub_ids:
            related_proc_ids = set(r[0] for r in db.session.query(ProcessKpi.process_id).filter(
                ProcessKpi.sub_strategy_id.in_(sub_ids), ProcessKpi.is_active.is_(True),
            ).all())
        scores = [proc_scores.get(pid) for pid in related_proc_ids if proc_scores.get(pid) is not None]
        avg_score = round(sum(scores) / len(scores), 1) if scores else None

        # ROI: spent başına skor (ham metric)
        roi_score = None
        if spent and avg_score is not None:
            roi_score = round(avg_score / (float(spent) / 100000), 3)  # 100K bazlı

        rows.append({
            "code": s.code, "title": s.title,
            "budget": float(budget), "spent": float(spent),
            "avg_score": avg_score,
            "roi_score": roi_score,
        })

    return jsonify({"success": True, "summary": {
        "total_strategies": len(rows),
        "total_budget": sum(r["budget"] for r in rows),
        "total_spent": sum(r["spent"] for r in rows),
        "plan_year": active_py.year if active_py else None,
    }, "strategies": rows})


# ─── HR-07: Bireysel Hedef Hizalama ────────────────────────────────────────

@app_bp.route("/raporlar/bireysel-hizalama")
@login_required
def raporlar_bireysel_hizalama():
    return render_template("platform/raporlar/bireysel_hizalama.html")


@app_bp.route("/raporlar/api/bireysel-hizalama")
@login_required
def raporlar_api_bireysel_hizalama():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    # Her aktif kullanıcı için: kaç bireysel PG'si var, kaçı bir sürece bağlı
    users = User.query.filter_by(tenant_id=tid, is_active=True).all()
    rows = []
    for u in users:
        pgs = IndividualPerformanceIndicator.query.filter_by(
            user_id=u.id, is_active=True
        ).all()
        if not pgs:
            continue
        total = len(pgs)
        aligned = sum(1 for p in pgs if p.source_process_id or p.source_process_kpi_id)
        rows.append({
            "user_id": u.id, "name": f"{u.first_name or ''} {u.last_name or ''}".strip() or u.email,
            "email": u.email, "department": u.department,
            "total_pgs": total, "aligned_pgs": aligned,
            "alignment_pct": round(aligned / total * 100, 1),
        })
    rows.sort(key=lambda x: -x["alignment_pct"])

    overall = round(sum(r["alignment_pct"] for r in rows) / max(len(rows), 1), 1)

    return jsonify({"success": True, "summary": {
        "users_with_pg": len(rows),
        "total_users": len(users),
        "overall_alignment_pct": overall,
    }, "users": rows[:50]})


# ─── RK-01: Risk Heatmap Detay ─────────────────────────────────────────────

@app_bp.route("/raporlar/risk-heatmap")
@login_required
def raporlar_risk_heatmap():
    return render_template("platform/raporlar/risk_heatmap.html")


@app_bp.route("/raporlar/api/risk-heatmap")
@login_required
def raporlar_api_risk_heatmap():
    from app.models.k_radar_domain import RiskHeatmapItem
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    risks = RiskHeatmapItem.query.filter_by(tenant_id=tid, is_active=True).order_by(
        RiskHeatmapItem.rpn.desc().nullslast()
    ).all()

    # 5×5 grid
    grid = [[[] for _ in range(5)] for _ in range(5)]  # grid[probability-1][impact-1]
    for r in risks:
        p = max(1, min(r.probability or 1, 5))
        i = max(1, min(r.impact or 1, 5))
        grid[p - 1][i - 1].append({
            "id": r.id, "title": r.title, "status": r.status or "Open",
            "rpn": r.rpn, "source": r.source_type,
        })

    summary = {
        "total": len(risks),
        "open": sum(1 for r in risks if (r.status or "Open").lower() == "open"),
        "mitigated": sum(1 for r in risks if (r.status or "").lower() == "mitigated"),
        "high_rpn": sum(1 for r in risks if (r.rpn or 0) >= 15),
    }
    return jsonify({"success": True, "summary": summary, "grid": grid,
                    "top_risks": [{"id": r.id, "title": r.title, "rpn": r.rpn,
                                   "probability": r.probability, "impact": r.impact,
                                   "status": r.status} for r in risks[:10]]})


# ─── RK-08: 2FA Kullanım Raporu ────────────────────────────────────────────

@app_bp.route("/raporlar/iki-fa")
@login_required
def raporlar_iki_fa():
    return render_template("platform/raporlar/iki_fa.html")


@app_bp.route("/raporlar/api/iki-fa")
@login_required
def raporlar_api_iki_fa():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    total = User.query.filter_by(tenant_id=tid, is_active=True).count()
    enabled = User.query.filter_by(tenant_id=tid, is_active=True, totp_enabled=True).count()
    disabled_admins = db.session.query(User.id, User.email, User.first_name, User.last_name).join(
        User.role
    ).filter(
        User.tenant_id == tid, User.is_active.is_(True),
        User.totp_enabled.is_(False),
        User.role.has(name="tenant_admin"),
    ).all()

    return jsonify({"success": True, "summary": {
        "total_users": total,
        "totp_enabled": enabled,
        "totp_disabled": total - enabled,
        "enable_pct": round(enabled / max(total, 1) * 100, 1),
        "admins_without_2fa": len(disabled_admins),
    }, "admins_without_2fa": [{
        "id": a.id, "email": a.email,
        "name": f"{a.first_name or ''} {a.last_name or ''}".strip(),
    } for a in disabled_admins]})


# ─── ES-01: Carbon Footprint Toplam Trend ──────────────────────────────────

@app_bp.route("/raporlar/carbon-trend")
@login_required
def raporlar_carbon_trend():
    return render_template("platform/raporlar/carbon_trend.html")


@app_bp.route("/raporlar/api/carbon-trend")
@login_required
def raporlar_api_carbon_trend():
    from app.models.esg import EsgMetric, EsgMetricValue
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    # Tüm scope 1+2+3 metrikleri
    metrics = EsgMetric.query.filter_by(tenant_id=tid, is_active=True).filter(
        EsgMetric.scope.in_(["scope1", "scope2", "scope3"])
    ).all()

    by_year_scope = defaultdict(lambda: {"scope1": 0, "scope2": 0, "scope3": 0})
    for m in metrics:
        for v in m.values:
            if v.value is not None:
                by_year_scope[v.year][m.scope] += v.value or 0

    years_sorted = sorted(by_year_scope.keys())
    trend = []
    for y in years_sorted:
        d = by_year_scope[y]
        trend.append({
            "year": y, "scope1": round(d["scope1"], 2), "scope2": round(d["scope2"], 2),
            "scope3": round(d["scope3"], 2),
            "total": round(d["scope1"] + d["scope2"] + d["scope3"], 2),
        })

    return jsonify({"success": True, "summary": {
        "metrics_count": len(metrics),
        "years_with_data": len(trend),
        "latest_total": trend[-1]["total"] if trend else 0,
        "first_total": trend[0]["total"] if trend else 0,
    }, "trend": trend})


# ─── AI-02: AI Strateji Danışmanı ──────────────────────────────────────────

@app_bp.route("/raporlar/ai-danisman")
@login_required
def raporlar_ai_danisman():
    return render_template("platform/raporlar/ai_danisman.html")


@app_bp.route("/raporlar/api/ai-danisman")
@login_required
def raporlar_api_ai_danisman():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    try:
        from app.services.ai_pivot_advisor_service import generate_pivot_recommendations
        result = generate_pivot_recommendations(tenant_id=tid, use_llm=True, user_id=current_user.id)
        return jsonify({"success": True, "data": result})
    except Exception as e:
        current_app.logger.warning(f"[ai-danisman] hata: {e}")
        return jsonify({"success": False, "message": f"AI danışman hatası: {e}"}), 500


# ─── AI-04: AI Coach ───────────────────────────────────────────────────────

@app_bp.route("/raporlar/ai-coach")
@login_required
def raporlar_ai_coach():
    return render_template("platform/raporlar/ai_coach.html")


@app_bp.route("/raporlar/api/ai-coach")
@login_required
def raporlar_api_ai_coach():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    # En düşük skorlu 3 stratejiyi tespit et + AI ile öneri al
    active_py = get_active_plan_year_for_user(current_user)
    try:
        today = _date.today()
        score_year = active_py.year if active_py else today.year
        proc_scores, _ = compute_process_scores_internal(tid, score_year, today, persist_pg_scores=False, plan_year=active_py)
    except Exception:
        proc_scores = {}

    strat_q = Strategy.query.filter_by(tenant_id=tid, is_active=True)
    if active_py:
        strat_q = strat_q.filter(or_(Strategy.plan_year_id == active_py.id, Strategy.plan_year_id.is_(None)))
    strategies = strat_q.all()

    strat_scores = []
    for s in strategies:
        sub_ids = [ss.id for ss in s.sub_strategies if getattr(ss, "is_active", True)]
        if not sub_ids:
            continue
        related = set(r[0] for r in db.session.query(ProcessKpi.process_id).filter(
            ProcessKpi.sub_strategy_id.in_(sub_ids)
        ).all())
        sc = [proc_scores.get(pid) for pid in related if proc_scores.get(pid) is not None]
        if sc:
            strat_scores.append({
                "code": s.code, "title": s.title,
                "score": round(sum(sc) / len(sc), 1),
            })

    strat_scores.sort(key=lambda x: x["score"])
    bottom3 = strat_scores[:3]

    # AI coach servisini çağır
    try:
        from services.ai_coach_service import analyze_strategic_performance
        coach_text = analyze_strategic_performance(
            data={
                "tenant_id": tid,
                "bottom_strategies": bottom3,
                "all_strategies": strat_scores,
            },
            tenant_id=tid, user_id=current_user.id,
        )
    except Exception as e:
        current_app.logger.info(f"[ai-coach] fallback ({e})")
        coach_text = (
            f"Sistemde {len(strat_scores)} ana strateji izlenmektedir. "
            f"En düşük performanslı 3 strateji: " +
            ", ".join(f"{s['code']} ({s['score']})" for s in bottom3) + ". "
            "Bu stratejilere bağlı PG'lerin gözden geçirilmesi, hedef revizyonu ve "
            "ek initiative tahsisi önerilmektedir."
        )

    return jsonify({"success": True, "data": {
        "strategies_analyzed": len(strat_scores),
        "bottom3": bottom3,
        "coach_advice": coach_text,
    }})


# ─── AI-05: AI Early Warning Dashboard ─────────────────────────────────────

@app_bp.route("/raporlar/early-warning")
@login_required
def raporlar_early_warning():
    return render_template("platform/raporlar/early_warning.html")


@app_bp.route("/raporlar/api/early-warning")
@login_required
def raporlar_api_early_warning():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    # Trend analizi: son 3 dönem hedef altında olan PG'leri tespit
    active_py = get_active_plan_year_for_user(current_user)
    py_id = active_py.id if active_py else None

    kpi_q = ProcessKpi.query.join(Process).filter(
        Process.tenant_id == tid, Process.is_active.is_(True),
        ProcessKpi.is_active.is_(True),
    )
    if py_id:
        kpi_q = kpi_q.filter(Process.plan_year_id == py_id)
    all_kpis = kpi_q.all()

    warnings = []
    for k in all_kpis[:100]:
        # Son 6 ölçüm
        recent = KpiData.query.filter_by(
            process_kpi_id=k.id, is_active=True,
        ).order_by(KpiData.data_date.desc()).limit(6).all()
        if len(recent) < 3:
            continue
        # Status_percentage düşüş trendi
        scores = [r.status_percentage for r in recent if r.status_percentage is not None]
        if len(scores) >= 3:
            last3 = scores[:3]
            if all(s < 70 for s in last3):
                warnings.append({
                    "kpi_id": k.id, "kpi_code": k.code or f"PG-{k.id}",
                    "kpi_name": k.name,
                    "process_name": k.process.name if k.process else "?",
                    "last_scores": last3,
                    "severity": "high" if all(s < 50 for s in last3) else "medium",
                })

    return jsonify({"success": True, "summary": {
        "total_kpis_checked": len(all_kpis),
        "warnings_count": len(warnings),
        "high_severity": sum(1 for w in warnings if w["severity"] == "high"),
    }, "warnings": warnings[:25]})


# ═══════════════════════════════════════════════════════════════════════════
# FAZ 2 — PERSONA DASHBOARDS + AI ÜRÜNLERİ
# ═══════════════════════════════════════════════════════════════════════════

# ─── PD-01: CFO Dashboard ──────────────────────────────────────────────────

@app_bp.route("/raporlar/cfo-dashboard")
@login_required
def raporlar_cfo_dashboard():
    return render_template("platform/raporlar/cfo_dashboard.html")


@app_bp.route("/raporlar/api/cfo-dashboard")
@login_required
def raporlar_api_cfo_dashboard():
    from app.models.llm_usage import LLMUsageLog
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    # Initiative bütçe özeti
    inits = Initiative.query.filter_by(tenant_id=tid, is_active=True).all()
    total_budget = sum(float(i.budget_total or 0) for i in inits)
    total_spent = sum(float(i.budget_spent or 0) for i in inits)
    over_budget = sum(1 for i in inits if (i.budget_spent or 0) > (i.budget_total or 0))

    # En büyük 5 initiative
    top5_inits = sorted(inits, key=lambda x: float(x.budget_total or 0), reverse=True)[:5]
    top5 = [{
        "code": i.code, "name": i.name, "status": i.status,
        "budget": float(i.budget_total or 0),
        "spent": float(i.budget_spent or 0),
        "progress": i.progress_pct or 0,
        "usage_pct": round(float(i.budget_spent or 0) / float(i.budget_total or 1) * 100, 1) if i.budget_total else 0,
    } for i in top5_inits]

    # LLM kullanım maliyeti (son 30 gün)
    last_30 = datetime.utcnow() - timedelta(days=30)
    llm_cost = db.session.query(func.coalesce(func.sum(LLMUsageLog.cost_usd), 0)).filter(
        LLMUsageLog.tenant_id == tid,
        LLMUsageLog.created_at >= last_30,
    ).scalar() or 0
    llm_calls = LLMUsageLog.query.filter(
        LLMUsageLog.tenant_id == tid,
        LLMUsageLog.created_at >= last_30,
    ).count()

    # Recurring task tahmini (yıllık)
    try:
        from app.models.portfolio_project import RecurringTask, Project
        recurring = db.session.query(func.count(RecurringTask.id)).join(
            Project, Project.id == RecurringTask.project_id,
        ).filter(Project.tenant_id == tid, RecurringTask.is_active.is_(True)).scalar() or 0
    except Exception:
        recurring = 0

    # Status dağılımı
    by_status = defaultdict(int)
    for i in inits:
        by_status[i.status] += 1

    # Strateji bazlı bütçe atıf
    by_strategy = defaultdict(lambda: {"budget": 0, "spent": 0, "count": 0})
    for i in inits:
        key = i.strategy_id or 0
        by_strategy[key]["budget"] += float(i.budget_total or 0)
        by_strategy[key]["spent"] += float(i.budget_spent or 0)
        by_strategy[key]["count"] += 1

    strat_budgets = []
    for sid, d in by_strategy.items():
        if sid:
            s = db.session.get(Strategy, sid)
            label = f"{s.code} {s.title}" if s else f"#{sid}"
        else:
            label = "(stratejisiz)"
        strat_budgets.append({
            "label": label, "budget": d["budget"], "spent": d["spent"],
            "count": d["count"],
            "usage_pct": round(d["spent"] / d["budget"] * 100, 1) if d["budget"] else 0,
        })
    strat_budgets.sort(key=lambda x: -x["budget"])

    return jsonify({"success": True, "metrics": {
        "total_budget": total_budget,
        "total_spent": total_spent,
        "remaining": total_budget - total_spent,
        "usage_pct": round(total_spent / total_budget * 100, 1) if total_budget else 0,
        "initiative_count": len(inits),
        "over_budget_count": over_budget,
        "llm_cost_30d_usd": float(llm_cost),
        "llm_calls_30d": llm_calls,
        "recurring_count": recurring,
    }, "by_status": dict(by_status),
       "top_initiatives": top5,
       "by_strategy": strat_budgets[:10]})


# ─── PD-02: COO Dashboard ──────────────────────────────────────────────────

@app_bp.route("/raporlar/coo-dashboard")
@login_required
def raporlar_coo_dashboard():
    return render_template("platform/raporlar/coo_dashboard.html")


@app_bp.route("/raporlar/api/coo-dashboard")
@login_required
def raporlar_api_coo_dashboard():
    from app.models.k_radar_domain import BottleneckLog, ProcessMaturity
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    active_py = get_active_plan_year_for_user(current_user)
    py_id = active_py.id if active_py else None

    # Süreç skorları
    try:
        today = _date.today()
        proc_scores, _ = compute_process_scores_internal(
            tid, active_py.year if active_py else today.year, today,
            persist_pg_scores=False, plan_year=active_py,
        )
    except Exception:
        proc_scores = {}

    procs = Process.query.filter_by(
        tenant_id=tid, is_active=True,
        plan_year_id=py_id,
    ).order_by(Process.code).all() if py_id else Process.query.filter_by(tenant_id=tid, is_active=True).all()

    proc_health = []
    for p in procs:
        score = proc_scores.get(p.id)
        proc_health.append({
            "id": p.id, "code": p.code, "name": p.name,
            "score": round(score, 1) if score is not None else None,
            "weight": p.weight,
        })

    # Skor dağılımı
    scored = [p["score"] for p in proc_health if p["score"] is not None]
    distribution = {
        "good": sum(1 for s in scored if s >= 70),
        "medium": sum(1 for s in scored if 50 <= s < 70),
        "low": sum(1 for s in scored if s < 50),
        "no_data": len(proc_health) - len(scored),
    }

    # Darboğaz
    bottlenecks = BottleneckLog.query.filter_by(tenant_id=tid, is_active=True).order_by(
        BottleneckLog.triggered_at.desc().nullslast()
    ).limit(10).all()
    bn_list = [{
        "id": b.id, "process_id": b.process_id, "severity": b.severity, "note": b.note,
        "triggered_at": b.triggered_at.isoformat() if b.triggered_at else None,
        "resolved": b.resolved_at is not None,
    } for b in bottlenecks]

    # Geciken faaliyet sayımı
    overdue_count = ProcessActivity.query.join(Process).filter(
        Process.tenant_id == tid,
        ProcessActivity.is_active.is_(True),
        ProcessActivity.end_at < datetime.utcnow(),
        ProcessActivity.status != "Tamamlandı",
    ).count()

    # CMMI ortalama
    maturity_rows = db.session.query(func.avg(ProcessMaturity.maturity_level)).filter(
        ProcessMaturity.tenant_id == tid, ProcessMaturity.is_active.is_(True),
    ).scalar()
    avg_cmmi = round(float(maturity_rows), 2) if maturity_rows else None

    return jsonify({"success": True, "metrics": {
        "total_processes": len(proc_health),
        "avg_score": round(sum(scored) / len(scored), 1) if scored else None,
        "overdue_activities": overdue_count,
        "active_bottlenecks": sum(1 for b in bottlenecks if not b.resolved_at),
        "avg_cmmi_level": avg_cmmi,
    }, "distribution": distribution,
       "processes": proc_health,
       "bottlenecks": bn_list})


# ─── PD-03: CHRO Dashboard ──────────────────────────────────────────────────

@app_bp.route("/raporlar/chro-dashboard")
@login_required
def raporlar_chro_dashboard():
    return render_template("platform/raporlar/chro_dashboard.html")


@app_bp.route("/raporlar/api/chro-dashboard")
@login_required
def raporlar_api_chro_dashboard():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    # Çalışan + departman özeti
    total_users = User.query.filter_by(tenant_id=tid, is_active=True).count()
    totp_enabled = User.query.filter_by(tenant_id=tid, is_active=True, totp_enabled=True).count()

    dept_rows = db.session.query(
        User.department, func.count(User.id),
    ).filter(User.tenant_id == tid, User.is_active.is_(True)).group_by(User.department).all()
    departments = [{"dept": d or "(belirsiz)", "count": c} for d, c in dept_rows]
    departments.sort(key=lambda x: -x["count"])

    # Bireysel PG özeti
    pg_users = db.session.query(func.count(func.distinct(IndividualPerformanceIndicator.user_id))).join(
        User, User.id == IndividualPerformanceIndicator.user_id
    ).filter(
        User.tenant_id == tid, User.is_active.is_(True),
        IndividualPerformanceIndicator.is_active.is_(True),
    ).scalar() or 0
    total_pgs = db.session.query(func.count(IndividualPerformanceIndicator.id)).join(
        User, User.id == IndividualPerformanceIndicator.user_id
    ).filter(
        User.tenant_id == tid, User.is_active.is_(True),
        IndividualPerformanceIndicator.is_active.is_(True),
    ).scalar() or 0

    # Rol dağılımı
    from app.models.core import Role
    role_rows = db.session.query(Role.name, func.count(User.id)).join(
        User, User.role_id == Role.id
    ).filter(User.tenant_id == tid, User.is_active.is_(True)).group_by(Role.name).all()
    roles = [{"role": r or "(yok)", "count": c} for r, c in role_rows]

    # En çok bireysel PG sahibi top 10 user
    top_user_rows = db.session.query(
        User.id, User.first_name, User.last_name, User.email, User.department,
        func.count(IndividualPerformanceIndicator.id).label("pg_count"),
    ).join(IndividualPerformanceIndicator, IndividualPerformanceIndicator.user_id == User.id).filter(
        User.tenant_id == tid, User.is_active.is_(True),
        IndividualPerformanceIndicator.is_active.is_(True),
    ).group_by(User.id, User.first_name, User.last_name, User.email, User.department).order_by(
        func.count(IndividualPerformanceIndicator.id).desc()
    ).limit(10).all()
    top_users = [{
        "id": r[0], "name": f"{r[1] or ''} {r[2] or ''}".strip() or r[3],
        "email": r[3], "department": r[4], "pg_count": r[5],
    } for r in top_user_rows]

    # Süreç-üye M:N (her kullanıcı kaç süreçte üye)
    from app.models.process import process_members
    member_rows = db.session.query(
        process_members.c.user_id, func.count(process_members.c.process_id),
    ).join(Process, Process.id == process_members.c.process_id).filter(
        Process.tenant_id == tid, Process.is_active.is_(True),
    ).group_by(process_members.c.user_id).all()
    process_membership = {uid: cnt for uid, cnt in member_rows}
    most_loaded = sorted(process_membership.items(), key=lambda x: -x[1])[:5]
    most_loaded_users = []
    for uid, cnt in most_loaded:
        u = db.session.get(User, uid)
        if u:
            most_loaded_users.append({
                "id": uid, "name": f"{u.first_name or ''} {u.last_name or ''}".strip() or u.email,
                "process_count": cnt,
            })

    return jsonify({"success": True, "metrics": {
        "total_users": total_users,
        "total_departments": len(departments),
        "totp_pct": round(totp_enabled / max(total_users, 1) * 100, 1),
        "users_with_pg": pg_users,
        "total_pgs": total_pgs,
        "avg_pg_per_user": round(total_pgs / max(pg_users, 1), 1),
    }, "departments": departments[:15],
       "roles": roles,
       "top_users_by_pg": top_users,
       "most_loaded_in_processes": most_loaded_users})


# ─── AI-11: AI Quarterly Review Hazırlayıcısı ──────────────────────────────

@app_bp.route("/raporlar/quarterly-review")
@login_required
def raporlar_quarterly_review():
    return render_template("platform/raporlar/quarterly_review.html")


@app_bp.route("/raporlar/api/quarterly-review")
@login_required
def raporlar_api_quarterly_review():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    today = _date.today()
    quarter = (today.month - 1) // 3 + 1
    active_py = get_active_plan_year_for_user(current_user)

    # Son çeyreğin metrikleri
    q_start_month = (quarter - 1) * 3 + 1
    q_start = _date(today.year, q_start_month, 1)
    if quarter < 4:
        q_end = _date(today.year, q_start_month + 3, 1) - timedelta(days=1)
    else:
        q_end = _date(today.year, 12, 31)

    # Çeyrek içinde girilen KPI ölçüm sayısı
    proc_ids = db.session.query(Process.id).filter(
        Process.tenant_id == tid, Process.is_active.is_(True),
    ).subquery()
    kpi_ids = db.session.query(ProcessKpi.id).filter(
        ProcessKpi.process_id.in_(proc_ids), ProcessKpi.is_active.is_(True),
    ).subquery()

    measurements_q = KpiData.query.filter(
        KpiData.process_kpi_id.in_(kpi_ids),
        KpiData.data_date.between(q_start, q_end),
        KpiData.is_active.is_(True),
    ).count()

    # Yeni initiative + biten initiative bu çeyrek
    new_inits = Initiative.query.filter(
        Initiative.tenant_id == tid,
        Initiative.created_at >= q_start,
    ).count()
    completed_inits = Initiative.query.filter(
        Initiative.tenant_id == tid,
        Initiative.status == "completed",
        Initiative.updated_at >= q_start,
    ).count()

    # AI ile çeyrek özeti üret
    summary_text = (
        f"{today.year} Q{quarter} ({q_start.strftime('%d %b')} - {q_end.strftime('%d %b')}) "
        f"döneminde sistem {measurements_q:,} KPI ölçümü topladı. "
        f"{new_inits} yeni stratejik girişim başlatıldı, {completed_inits} girişim tamamlandı. "
        "Çeyrek değerlendirme toplantısı için aşağıdaki gündem önerilmektedir."
    )
    try:
        from app.services.llm_gateway import call_llm
        prompt = (
            f"Türkçe, 3-4 cümlelik bir çeyrek değerlendirme özeti yaz. "
            f"Çeyrek: {today.year} Q{quarter}. "
            f"Veri: {measurements_q} ölçüm, {new_inits} yeni girişim, {completed_inits} tamamlanmış girişim. "
            "Resmî, yöneticiye uygun ton."
        )
        r = call_llm(prompt=prompt, tenant_id=tid, user_id=current_user.id,
                     endpoint="ai_quarterly_review", max_output_tokens=300)
        if r and r.get("text"):
            summary_text = r["text"].strip()
    except Exception as e:
        current_app.logger.info(f"[quarterly-review] LLM fallback: {e}")

    # Önerilen ajanda
    agenda = [
        f"1. Yönetici özeti — {today.year} Q{quarter} performans snapshot (10 dk)",
        "2. Stratejik hedef başarı durumu (Vizyon skor + 6 ana strateji) (15 dk)",
        "3. Initiative portföyü gözden geçirme + bütçe durumu (20 dk)",
        "4. OKR check-in — KR ilerlemeleri (15 dk)",
        "5. Risk register güncelleme — yeni risk + mitigation (10 dk)",
        "6. Önümüzdeki çeyrek için odak alanları ve replan kararları (15 dk)",
        "7. Açık tartışma ve aksiyon maddeleri (15 dk)",
    ]

    # Ön çalışma soruları
    prep_questions = [
        "Bu çeyrekte hedeflerinin altında kalan PG'ler için sebep analizi hazırlandı mı?",
        "Bütçe sapması olan initiative'ler için açıklama metni var mı?",
        "Yeni risk olarak değerlendirilmesi gereken durumlar belirlendi mi?",
        "Önümüzdeki çeyrek için yeni initiative önerisi var mı?",
        "Geçen çeyreğin aksiyon maddelerinin durumu güncellendi mi?",
    ]

    return jsonify({"success": True, "data": {
        "year": today.year, "quarter": quarter,
        "period_start": q_start.isoformat(), "period_end": q_end.isoformat(),
        "plan_year": active_py.year if active_py else None,
        "metrics": {
            "measurements_q": measurements_q,
            "new_initiatives": new_inits,
            "completed_initiatives": completed_inits,
        },
        "ai_summary": summary_text,
        "agenda": agenda,
        "prep_questions": prep_questions,
    }})


# ─── AI-12: AI Strateji Hikayeleştirici ────────────────────────────────────

@app_bp.route("/raporlar/strateji-hikayesi")
@login_required
def raporlar_strateji_hikayesi():
    return render_template("platform/raporlar/strateji_hikayesi.html")


@app_bp.route("/raporlar/api/strateji-hikayesi")
@login_required
def raporlar_api_strateji_hikayesi():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    tenant = db.session.get(Tenant, tid)
    py_list = list_plan_years(tid)
    if not py_list:
        return jsonify({"success": True, "data": {"narrative": "Henüz plan yılı yok.", "highlights": []}})

    py_sorted = sorted(py_list, key=lambda x: x.year)
    # Her yıl için snapshot özet
    snapshots = []
    for py in py_sorted:
        s_count = Strategy.query.filter_by(tenant_id=tid, plan_year_id=py.id, is_active=True).count()
        p_count = Process.query.filter_by(tenant_id=tid, plan_year_id=py.id, is_active=True).count()
        k_count = ProcessKpi.query.join(Process).filter(
            Process.tenant_id == tid, Process.plan_year_id == py.id,
            ProcessKpi.is_active.is_(True),
        ).count()
        m_count = db.session.query(func.count(KpiData.id)).filter(
            KpiData.year == py.year,
        ).scalar() or 0
        snapshots.append({
            "year": py.year, "status": py.status,
            "strategy_count": s_count, "process_count": p_count,
            "kpi_count": k_count, "measurement_count": m_count,
        })

    # Kırılım noktaları (yıl yıl yapısal değişim)
    highlights = []
    for i in range(1, len(snapshots)):
        prev, curr = snapshots[i - 1], snapshots[i]
        delta_s = curr["strategy_count"] - prev["strategy_count"]
        delta_p = curr["process_count"] - prev["process_count"]
        if abs(delta_s) >= 1 or abs(delta_p) >= 2:
            highlights.append({
                "year": curr["year"],
                "change": f"Strateji {prev['strategy_count']}→{curr['strategy_count']}, Süreç {prev['process_count']}→{curr['process_count']}",
            })

    # AI ile hikaye yaz
    fallback_narrative = (
        f"{tenant.name if tenant else 'Kurum'} {snapshots[0]['year']} yılında "
        f"{snapshots[0]['strategy_count']} stratejik direkle başlayıp "
        f"{snapshots[-1]['year']} yılına {snapshots[-1]['strategy_count']} strateji + "
        f"{snapshots[-1]['process_count']} süreçle ulaştı. "
        f"Toplam {sum(s['measurement_count'] for s in snapshots):,} ölçüm verisi "
        "stratejik kararların temelini oluşturdu."
    )
    narrative = fallback_narrative
    try:
        from app.services.llm_gateway import call_llm
        data_str = "\n".join([
            f"  {s['year']}: {s['strategy_count']} strateji, {s['process_count']} süreç, {s['kpi_count']} PG, {s['measurement_count']:,} ölçüm"
            for s in snapshots
        ])
        prompt = (
            f"Sen kurumsal stratejik danışmansın. {tenant.name if tenant else 'Bir kurum'} için "
            f"{len(snapshots)} yıllık stratejik plan evrimini anlatan 4-6 cümlelik bir hikaye yaz. "
            f"Veri:\n{data_str}\n\n"
            f"Önemli değişim noktaları: {[h['year'] for h in highlights]}.\n"
            "Resmî ama akıcı bir dil kullan. Sayılarla başla, evrime ve dönüşüme vurgu yap."
        )
        r = call_llm(prompt=prompt, tenant_id=tid, user_id=current_user.id,
                     endpoint="ai_strateji_hikaye", max_output_tokens=600)
        if r and r.get("text"):
            narrative = r["text"].strip()
    except Exception as e:
        current_app.logger.info(f"[strateji-hikaye] LLM fallback: {e}")

    return jsonify({"success": True, "data": {
        "tenant_name": tenant.name if tenant else "—",
        "year_range": f"{snapshots[0]['year']}–{snapshots[-1]['year']}",
        "narrative": narrative,
        "snapshots": snapshots,
        "highlights": highlights,
    }})


# ═══════════════════════════════════════════════════════════════════════════
# FAZ 3 — PREMIUM ÜRÜNLER (indirilebilir PDF/PPTX/Excel/ZIP)
# ═══════════════════════════════════════════════════════════════════════════

def _pdf_helpers():
    """Reportlab PDF üretici yardımcılar."""
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Image,
    )
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    return {
        "SimpleDocTemplate": SimpleDocTemplate, "Paragraph": Paragraph,
        "Spacer": Spacer, "PageBreak": PageBreak, "Table": Table,
        "TableStyle": TableStyle, "A4": A4, "cm": cm, "colors": colors,
        "styles": getSampleStyleSheet(),
        "ParagraphStyle": ParagraphStyle,
        "TA_CENTER": TA_CENTER, "TA_JUSTIFY": TA_JUSTIFY,
    }


def _ai_text(prompt: str, fallback: str, tid: int, endpoint: str, max_tokens: int = 400) -> str:
    """LLM çağrısı yap, başarısızsa fallback metin dön."""
    try:
        from app.services.llm_gateway import call_llm
        r = call_llm(prompt=prompt, tenant_id=tid, user_id=current_user.id,
                     endpoint=endpoint, max_output_tokens=max_tokens)
        if r and r.get("text"):
            return r["text"].strip()
    except Exception as e:
        current_app.logger.info(f"[{endpoint}] LLM fallback: {e}")
    return fallback


# ─── ST-16: Stratejik Yıllık Kitap PDF ─────────────────────────────────────

@app_bp.route("/raporlar/stratejik-yillik")
@login_required
def raporlar_stratejik_yillik():
    return render_template("platform/raporlar/stratejik_yillik.html")


@app_bp.route("/raporlar/api/stratejik-yillik/preview")
@login_required
def raporlar_api_stratejik_yillik_preview():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    tenant = db.session.get(Tenant, tid)
    active_py = get_active_plan_year_for_user(current_user)

    sections = [
        {"no": 1, "title": "Kapak", "content": "Kurum logo + yıl + başlık"},
        {"no": 2, "title": "Önsöz", "content": "AI özet — 1 sayfa giriş"},
        {"no": 3, "title": "Yönetici Özeti", "content": "Yıl özetinde 3 cümle + 5 kritik metrik"},
        {"no": 4, "title": "Kurumsal Kimlik", "content": "Vizyon, misyon, değerler, etik"},
        {"no": 5, "title": "Stratejik Direkler", "content": "Ana stratejiler + K-Vektör ağırlık"},
        {"no": 6, "title": "Süreç Mükemmelliği", "content": "Süreç sağlığı + CMMI olgunluk"},
        {"no": 7, "title": "Initiative Portföyü", "content": "Yıl içi girişimler + bütçe"},
        {"no": 8, "title": "Performans Göstergeleri", "content": "PG hedef-gerçek özet"},
        {"no": 9, "title": "Risk Yönetimi", "content": "Risk heatmap + top 10"},
        {"no": 10, "title": "Sürdürülebilirlik (ESG)", "content": "Carbon + sosyal + yönetişim"},
        {"no": 11, "title": "İK ve Yetenek", "content": "Çalışan + departman + bireysel başarı"},
        {"no": 12, "title": "Yıllar Arası Karşılaştırma", "content": "Önceki yıllarla evrim"},
        {"no": 13, "title": "Önümüzdeki Dönem", "content": "Yol haritası + AI önerileri"},
        {"no": 14, "title": "Ekler", "content": "Veri kaynağı, metodoloji, terimler"},
    ]
    return jsonify({"success": True, "preview": {
        "tenant_name": tenant.name if tenant else "—",
        "plan_year": active_py.year if active_py else None,
        "sections": sections,
        "estimated_pages": 35,
        "format": "PDF (reportlab)",
    }})


@app_bp.route("/raporlar/api/stratejik-yillik/generate", methods=["GET", "POST"])
@login_required
def raporlar_api_stratejik_yillik_generate():
    from flask import send_file
    import io
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400

    h = _pdf_helpers()
    tenant = db.session.get(Tenant, tid)
    active_py = get_active_plan_year_for_user(current_user)
    year_label = active_py.year if active_py else _date.today().year
    tname = tenant.name if tenant else "Kurum"

    # Veri toplama
    strat_q = Strategy.query.filter_by(tenant_id=tid, is_active=True)
    if active_py:
        strat_q = strat_q.filter_by(plan_year_id=active_py.id)
    strategies = strat_q.order_by(Strategy.code).all()
    weights = {w.strategy_id: w.weight_raw or 0 for w in
               KVektorStrategyWeight.query.filter_by(tenant_id=tid).all()}
    total_w = sum(weights.values()) or 1

    proc_count = Process.query.filter_by(tenant_id=tid, is_active=True,
                                          plan_year_id=active_py.id if active_py else None).count()
    kpi_count = ProcessKpi.query.join(Process).filter(
        Process.tenant_id == tid, Process.plan_year_id == (active_py.id if active_py else None),
        ProcessKpi.is_active.is_(True),
    ).count()
    initiatives = Initiative.query.filter_by(tenant_id=tid, is_active=True).all()
    user_count = User.query.filter_by(tenant_id=tid, is_active=True).count()
    meas_count = db.session.query(func.count(KpiData.id)).join(ProcessKpi).join(
        Process).filter(Process.tenant_id == tid).scalar() or 0

    # AI önsöz
    onsoz = _ai_text(
        prompt=f"{tname} {year_label} Stratejik Yıllık Kitabı için 4-5 cümlelik bir önsöz yaz. "
               f"Veri: {len(strategies)} strateji, {proc_count} süreç, {kpi_count} PG, "
               f"{len(initiatives)} initiative, {user_count} çalışan, {meas_count} ölçüm. "
               "Resmî, ilham verici ton kullan.",
        fallback=(f"{tname} olarak {year_label} yılında stratejik planlama yolculuğumuzda "
                  f"{len(strategies)} ana direk üzerinden {proc_count} süreç ve {kpi_count} "
                  f"performans göstergesi ile yönettik. {user_count} çalışanımızın katkısıyla "
                  f"{meas_count:,} ölçüm verisi sisteme aktarıldı. Bu yıllık, "
                  "yıl içinde elde ettiğimiz başarıları, karşılaştığımız sorunları ve "
                  "önümüzdeki dönem için yol haritamızı bir arada sunmaktadır."),
        tid=tid, endpoint="ai_stratejik_yillik_onsoz", max_tokens=350,
    )

    # PDF üret
    buf = io.BytesIO()
    doc = h["SimpleDocTemplate"](
        buf, pagesize=h["A4"],
        leftMargin=2*h["cm"], rightMargin=2*h["cm"],
        topMargin=2.5*h["cm"], bottomMargin=2.5*h["cm"],
        title=f"{tname} {year_label} Stratejik Yıllık", author="Kokpitim",
    )
    styles = h["styles"]
    title_style = h["ParagraphStyle"](
        "TitleBig", parent=styles["Title"], fontSize=32, leading=40,
        textColor=h["colors"].HexColor("#0f172a"), spaceAfter=20, alignment=h["TA_CENTER"],
    )
    h1 = h["ParagraphStyle"]("H1", parent=styles["Heading1"], fontSize=20,
        textColor=h["colors"].HexColor("#4f46e5"), spaceAfter=16, spaceBefore=12)
    h2 = h["ParagraphStyle"]("H2", parent=styles["Heading2"], fontSize=14,
        textColor=h["colors"].HexColor("#1e293b"), spaceAfter=10, spaceBefore=10)
    body = h["ParagraphStyle"]("Body", parent=styles["BodyText"], fontSize=11,
        leading=16, textColor=h["colors"].HexColor("#0f172a"),
        spaceAfter=8, alignment=h["TA_JUSTIFY"])
    small = h["ParagraphStyle"]("Small", parent=body, fontSize=9,
        textColor=h["colors"].HexColor("#64748b"))

    P = h["Paragraph"]
    elems = []

    # Kapak
    elems.append(h["Spacer"](1, 6*h["cm"]))
    elems.append(P(tname, title_style))
    elems.append(P(f"{year_label}", title_style))
    elems.append(h["Spacer"](1, 1*h["cm"]))
    elems.append(P("STRATEJİK YILLIK", h["ParagraphStyle"](
        "Sub", parent=body, fontSize=16, alignment=h["TA_CENTER"],
        textColor=h["colors"].HexColor("#64748b"), spaceAfter=8)))
    elems.append(h["Spacer"](1, 3*h["cm"]))
    elems.append(P(f"Hazırlık tarihi: {_date.today().strftime('%d %B %Y')}",
                    h["ParagraphStyle"]("date", parent=small, alignment=h["TA_CENTER"])))
    elems.append(h["PageBreak"]())

    # Önsöz
    elems.append(P("Önsöz", h1))
    elems.append(P(onsoz, body))
    elems.append(h["PageBreak"]())

    # Yönetici özeti
    elems.append(P("Yönetici Özeti", h1))
    elems.append(P(f"<b>{year_label}</b> yılı için kurum geneli görünümü:", body))
    summary_rows = [
        ["Strateji sayısı", str(len(strategies))],
        ["Aktif süreç", str(proc_count)],
        ["Performans göstergesi", str(kpi_count)],
        ["Stratejik girişim", str(len(initiatives))],
        ["Aktif çalışan", str(user_count)],
        ["Toplam ölçüm", f"{meas_count:,}"],
    ]
    tbl = h["Table"](summary_rows, colWidths=[8*h["cm"], 5*h["cm"]])
    tbl.setStyle(h["TableStyle"]([
        ("BACKGROUND", (0, 0), (0, -1), h["colors"].HexColor("#f8fafc")),
        ("TEXTCOLOR", (1, 0), (1, -1), h["colors"].HexColor("#4f46e5")),
        ("FONTNAME", (1, 0), (1, -1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 11),
        ("BOX", (0, 0), (-1, -1), 1, h["colors"].HexColor("#e2e8f0")),
        ("INNERGRID", (0, 0), (-1, -1), 0.5, h["colors"].HexColor("#e2e8f0")),
        ("PADDING", (0, 0), (-1, -1), 10),
    ]))
    elems.append(tbl)
    elems.append(h["PageBreak"]())

    # Kurumsal Kimlik
    elems.append(P("Kurumsal Kimlik", h1))
    if tenant:
        if tenant.purpose:
            elems.append(P("Amaç", h2)); elems.append(P(tenant.purpose, body))
        if tenant.vision:
            elems.append(P("Vizyon", h2)); elems.append(P(tenant.vision, body))
        if tenant.core_values:
            elems.append(P("Değerler", h2)); elems.append(P(tenant.core_values, body))
        if tenant.code_of_ethics:
            elems.append(P("Etik Kurallar", h2)); elems.append(P(tenant.code_of_ethics, body))
    else:
        elems.append(P("Kurumsal kimlik bilgisi tanımlanmamış.", small))
    elems.append(h["PageBreak"]())

    # Stratejik Direkler
    elems.append(P(f"Stratejik Direkler — {len(strategies)} Ana Strateji", h1))
    if strategies:
        rows = [["Kod", "Ad", "K-Vektör %"]]
        for s in strategies:
            w_pct = round(weights.get(s.id, 0) / total_w * 100, 1)
            rows.append([s.code or "—", (s.title or "")[:60], f"%{w_pct}"])
        tbl = h["Table"](rows, colWidths=[2*h["cm"], 11*h["cm"], 3*h["cm"]])
        tbl.setStyle(h["TableStyle"]([
            ("BACKGROUND", (0, 0), (-1, 0), h["colors"].HexColor("#4f46e5")),
            ("TEXTCOLOR", (0, 0), (-1, 0), h["colors"].white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 10),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1),
             [h["colors"].white, h["colors"].HexColor("#f8fafc")]),
            ("BOX", (0, 0), (-1, -1), 1, h["colors"].HexColor("#e2e8f0")),
            ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#e2e8f0")),
            ("PADDING", (0, 0), (-1, -1), 8),
        ]))
        elems.append(tbl)
    else:
        elems.append(P("Strateji tanımlanmamış.", small))
    elems.append(h["PageBreak"]())

    # Initiative Portföyü
    elems.append(P(f"Initiative Portföyü — {len(initiatives)} Girişim", h1))
    if initiatives:
        rows = [["Kod", "Ad", "Durum", "İlerleme", "Bütçe (₺)"]]
        for i in initiatives[:30]:
            rows.append([
                i.code or "—", (i.name or "")[:40],
                i.status or "—", f"%{i.progress_pct or 0}",
                f"{float(i.budget_total or 0):,.0f}",
            ])
        tbl = h["Table"](rows, colWidths=[3*h["cm"], 6*h["cm"], 2.5*h["cm"], 2*h["cm"], 3*h["cm"]])
        tbl.setStyle(h["TableStyle"]([
            ("BACKGROUND", (0, 0), (-1, 0), h["colors"].HexColor("#16a34a")),
            ("TEXTCOLOR", (0, 0), (-1, 0), h["colors"].white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1),
             [h["colors"].white, h["colors"].HexColor("#f8fafc")]),
            ("BOX", (0, 0), (-1, -1), 1, h["colors"].HexColor("#e2e8f0")),
            ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#e2e8f0")),
            ("PADDING", (0, 0), (-1, -1), 6),
        ]))
        elems.append(tbl)
    else:
        elems.append(P("Initiative tanımlanmamış.", small))
    elems.append(h["PageBreak"]())

    # Önümüzdeki dönem
    elems.append(P("Önümüzdeki Dönem", h1))
    next_year_text = _ai_text(
        prompt=f"{tname} için {year_label+1} yılı önerilen 3-4 stratejik öncelik yaz. "
               "Mevcut yıldaki verilere bakarak (özellikle düşük performanslı alanlar).",
        fallback=(f"{year_label+1} yılında stratejik plan döngüsünün doğal devamı olarak: "
                  "(1) mevcut stratejilerin değerlendirilmesi, "
                  "(2) yeni initiative tahsisleri, "
                  "(3) operasyonel mükemmellik odakları, "
                  "(4) ESG metriklerinin yatırımcı standartlarına çekilmesi öngörülmektedir."),
        tid=tid, endpoint="ai_stratejik_yillik_next", max_tokens=350,
    )
    elems.append(P(next_year_text, body))
    elems.append(h["PageBreak"]())

    # Kapanış
    elems.append(P("Teşekkürler", h1))
    elems.append(P(f"Bu yıllık, {tname} ekibinin yıl boyunca verdiği emeğin yansımasıdır. "
                   "Kokpitim platformu üzerinde sistemli olarak toplanan veriler ile "
                   "üretilmiştir.", body))
    elems.append(h["Spacer"](1, 3*h["cm"]))
    elems.append(P(f"<b>{tname}</b> · {year_label} Stratejik Yıllık",
                    h["ParagraphStyle"]("FootCenter", parent=body, alignment=h["TA_CENTER"])))
    elems.append(P("Kokpitim — Kurumsal Performans Yönetim Platformu",
                    h["ParagraphStyle"]("FootSmall", parent=small, alignment=h["TA_CENTER"])))

    doc.build(elems)
    buf.seek(0)

    filename = f"{tname.replace(' ', '_')}_{year_label}_stratejik_yillik.pdf"
    return send_file(buf, mimetype="application/pdf",
                     as_attachment=True, download_name=filename)


# ─── AI-17: Yatırımcı Sunum PPTX ───────────────────────────────────────────

@app_bp.route("/raporlar/yatirimci-sunum")
@login_required
def raporlar_yatirimci_sunum():
    return render_template("platform/raporlar/yatirimci_sunum.html")


@app_bp.route("/raporlar/api/yatirimci-sunum/preview")
@login_required
def raporlar_api_yatirimci_sunum_preview():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    tenant = db.session.get(Tenant, tid)
    active_py = get_active_plan_year_for_user(current_user)

    slides = [
        {"no": 1, "title": "Kapak", "content": f"{tenant.name if tenant else 'Kurum'} — Yatırımcı Sunumu"},
        {"no": 2, "title": "Yönetici Özeti", "content": "AI 4-cümle özet"},
        {"no": 3, "title": "Vizyon & Misyon", "content": "Stratejik kimlik"},
        {"no": 4, "title": "Pazardaki Yerimiz", "content": "Sektör + rakip benchmark"},
        {"no": 5, "title": "İş Modeli", "content": "Gelir kalemleri + büyüme stratejisi"},
        {"no": 6, "title": "Stratejik Direkler", "content": f"{Strategy.query.filter_by(tenant_id=tid, is_active=True).count()} ana strateji"},
        {"no": 7, "title": "Operasyonel Mükemmellik", "content": "Süreç olgunluk + verimlilik"},
        {"no": 8, "title": "Finansal Performans", "content": "Bütçe + initiative ROI"},
        {"no": 9, "title": "İnsan Kaynağı", "content": "Çalışan + yetenek geliştirme"},
        {"no": 10, "title": "ESG & Sürdürülebilirlik", "content": "Çevre + sosyal + yönetişim"},
        {"no": 11, "title": "Risk Yönetimi", "content": "Risk register + mitigation"},
        {"no": 12, "title": "Initiative Portföyü", "content": f"{Initiative.query.filter_by(tenant_id=tid, is_active=True).count()} girişim"},
        {"no": 13, "title": "Geçmiş Performans (7 yıl)", "content": "Yıllar arası trend"},
        {"no": 14, "title": "3 Yıllık Projeksiyon", "content": "Senaryo bazlı tahmin"},
        {"no": 15, "title": "Yatırım Talebi", "content": "Bütçe + use of funds"},
        {"no": 16, "title": "Yönetim Ekibi", "content": "Lider ekip"},
        {"no": 17, "title": "Soru-Cevap", "content": "Hazır cevaplar"},
        {"no": 18, "title": "Teşekkürler", "content": "İletişim"},
    ]
    return jsonify({"success": True, "preview": {
        "tenant_name": tenant.name if tenant else "—",
        "plan_year": active_py.year if active_py else None,
        "slides": slides,
        "format": "PPTX (python-pptx, 16:9)",
        "slide_count": len(slides),
    }})


@app_bp.route("/raporlar/api/yatirimci-sunum/generate", methods=["GET", "POST"])
@login_required
def raporlar_api_yatirimci_sunum_generate():
    from flask import send_file
    import io
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return jsonify({"success": False, "message": "python-pptx yok"}), 500

    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    tenant = db.session.get(Tenant, tid)
    active_py = get_active_plan_year_for_user(current_user)
    year_label = str(active_py.year) if active_py else str(_date.today().year)
    tname = tenant.name if tenant else "Kurum"

    strategies = Strategy.query.filter_by(tenant_id=tid, is_active=True,
        plan_year_id=active_py.id if active_py else None).order_by(Strategy.code).all()
    inits = Initiative.query.filter_by(tenant_id=tid, is_active=True).all()
    proc_count = Process.query.filter_by(tenant_id=tid, is_active=True).count()
    user_count = User.query.filter_by(tenant_id=tid, is_active=True).count()
    meas_count = db.session.query(func.count(KpiData.id)).join(ProcessKpi).join(
        Process).filter(Process.tenant_id == tid).scalar() or 0
    total_budget = sum(float(i.budget_total or 0) for i in inits)

    # AI özet
    investor_summary = _ai_text(
        prompt=f"{tname} için yatırımcılara yönelik 4-cümle özet yaz. "
               f"Strateji: {len(strategies)}, Süreç: {proc_count}, "
               f"Çalışan: {user_count}, Initiative bütçe: {total_budget:,.0f}₺. "
               "Güven verici, sayısal, profesyonel ton.",
        fallback=(f"{tname}, {len(strategies)} stratejik direk üzerinde "
                  f"{proc_count} süreç ve {user_count} çalışanın katkısıyla "
                  f"sürdürülebilir büyüme yolculuğundadır. {len(inits)} aktif stratejik girişim "
                  f"ile {total_budget:,.0f}₺'lik yatırım portföyü yönetilmekte, "
                  f"{meas_count:,} veri noktası ile karar verme süreçleri desteklenmektedir. "
                  "Bu sunum, mevcut performansımızı ve yatırımcı önerimizi sunmaktadır."),
        tid=tid, endpoint="ai_yatirimci_sunum", max_tokens=400,
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(0x4f, 0x46, 0xe5)
    DARK = RGBColor(0x0f, 0x17, 0x2a)
    GRAY = RGBColor(0x64, 0x74, 0x8b)
    GREEN = RGBColor(0x10, 0xb9, 0x81)

    def add_slide_title(slide, txt, color):
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        p = tx.text_frame.paragraphs[0]
        p.text = txt
        p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = color

    def add_body(slide, text, bullets=None, font=14):
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
        bf = body.text_frame; bf.word_wrap = True
        if text:
            p = bf.paragraphs[0]; p.text = text
            p.font.size = Pt(font); p.font.color.rgb = DARK
        if bullets:
            for b in bullets:
                bp = bf.add_paragraph(); bp.text = "• " + str(b)
                bp.font.size = Pt(font - 2); bp.font.color.rgb = GRAY
                bp.space_before = Pt(4)

    def add_footer(slide, no):
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.3))
        ft.text_frame.text = f"{tname} · {year_label} · Slayt {no}"
        ft.text_frame.paragraphs[0].font.size = Pt(10)
        ft.text_frame.paragraphs[0].font.color.rgb = GRAY

    # Slaytları üret
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, tname, DARK)
    add_body(s, f"\n\n{year_label} Yatırımcı Sunumu\n\n{_date.today().strftime('%d %B %Y')}", font=20)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, "Yönetici Özeti", BLUE)
    add_body(s, investor_summary); add_footer(s, 2)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, "Vizyon", BLUE)
    add_body(s, (tenant.vision if tenant else None) or "Vizyon tanımlanmamış."); add_footer(s, 3)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, "Misyon ve Değerler", BLUE)
    bullets = []
    if tenant and tenant.purpose: bullets.append("Amaç: " + (tenant.purpose[:200]))
    if tenant and tenant.core_values: bullets.append("Değerler: " + (tenant.core_values[:200]))
    add_body(s, None, bullets); add_footer(s, 4)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, "Stratejik Direkler", BLUE)
    add_body(s, f"{len(strategies)} ana stratejik direk:",
        bullets=[f"{st.code} — {st.title}" for st in strategies[:8]])
    add_footer(s, 5)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, "Operasyonel Performans", BLUE)
    add_body(s, "Yıllık operasyonel göstergeler:",
        bullets=[f"{proc_count} aktif süreç", f"{meas_count:,} KPI ölçümü",
                 f"{user_count} aktif çalışan",
                 "Süreçler CMMI olgunluk seviyesinde izlenmektedir"])
    add_footer(s, 6)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, "Finansal Görünüm", BLUE)
    add_body(s, "Stratejik girişim portföyü:",
        bullets=[f"Toplam initiative bütçesi: ₺{total_budget:,.0f}",
                 f"Aktif initiative: {len(inits)}",
                 f"Ortalama bütçe/initiative: ₺{total_budget/max(len(inits),1):,.0f}",
                 "EVM (PV/EV/AC) modeliyle takip"])
    add_footer(s, 7)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, "Initiative Portföyü", BLUE)
    add_body(s, f"En öne çıkan {min(8, len(inits))} stratejik girişim:",
        bullets=[f"{i.code or 'INI'} — {i.name} (%{i.progress_pct or 0} ilerleme, durum: {i.status})"
                 for i in sorted(inits, key=lambda x: -(x.progress_pct or 0))[:8]])
    add_footer(s, 8)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, "İnsan Sermayesi", BLUE)
    add_body(s, f"{user_count} aktif çalışan ile yetenek odaklı büyüme.",
        bullets=["Bireysel performans takibi", "Süreç sahipliği ve liderlik geliştirme",
                 "2FA güvenlik standartları"])
    add_footer(s, 9)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, "ESG ve Sürdürülebilirlik", GREEN)
    add_body(s, "Çevre, sosyal ve yönetişim taahhütleri:",
        bullets=["Scope 1+2+3 emisyon takibi", "SDG katkı haritası",
                 "İş güvenliği ve çeşitlilik metrikleri", "Bağımsız yatırımcı raporlaması"])
    add_footer(s, 10)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, "Risk Yönetimi", BLUE)
    add_body(s, "Stratejik risk register ile proaktif yönetim:",
        bullets=["5×5 risk heatmap (olasılık × etki)",
                 "Sahip atamalı mitigation planları",
                 "Çeyreklik risk değerlendirme", "AI Erken Uyarı sistemi"])
    add_footer(s, 11)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, "Yatırım Talebi", BLUE)
    add_body(s, "Önerilen yatırım kullanım planı:",
        bullets=["Ürün geliştirme ve Ar-Ge (~%40)",
                 "Operasyonel kapasite (~%30)",
                 "Pazarlama ve büyüme (~%20)",
                 "Yedek (~%10)"])
    add_footer(s, 12)

    s = prs.slides.add_slide(prs.slide_layouts[6]); add_slide_title(s, "Teşekkürler", DARK)
    add_body(s, f"\n\n{tname}\n\nİletişim: investor@example.com\n\nDetaylı analiz: Kokpitim platformu")
    add_footer(s, 13)

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    filename = f"{tname.replace(' ', '_')}_{year_label}_yatirimci_sunum.pptx"
    return send_file(buf,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True, download_name=filename)


# ─── ES-08: GRI/CDP/TCFD ESG Rapor PDF ─────────────────────────────────────

@app_bp.route("/raporlar/esg-rapor")
@login_required
def raporlar_esg_rapor():
    return render_template("platform/raporlar/esg_rapor.html")


@app_bp.route("/raporlar/api/esg-rapor/generate", methods=["GET", "POST"])
@login_required
def raporlar_api_esg_rapor_generate():
    from flask import send_file
    from app.models.esg import EsgMetric, EsgMetricValue
    import io

    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    tenant = db.session.get(Tenant, tid)
    tname = tenant.name if tenant else "Kurum"
    year = _date.today().year

    h = _pdf_helpers()
    buf = io.BytesIO()
    doc = h["SimpleDocTemplate"](buf, pagesize=h["A4"],
        leftMargin=2*h["cm"], rightMargin=2*h["cm"], topMargin=2*h["cm"], bottomMargin=2*h["cm"],
        title=f"{tname} ESG Yıllık Raporu {year}")
    styles = h["styles"]
    h1 = h["ParagraphStyle"]("H1", parent=styles["Heading1"], fontSize=18,
        textColor=h["colors"].HexColor("#16a34a"), spaceAfter=14)
    h2 = h["ParagraphStyle"]("H2", parent=styles["Heading2"], fontSize=13,
        textColor=h["colors"].HexColor("#0f172a"), spaceAfter=8)
    body = h["ParagraphStyle"]("Body", parent=styles["BodyText"], fontSize=10.5,
        leading=15, textColor=h["colors"].HexColor("#0f172a"),
        spaceAfter=6, alignment=h["TA_JUSTIFY"])

    P = h["Paragraph"]; elems = []

    # Kapak
    elems.append(P(f"<font color='#16a34a'>{tname}</font>",
                    h["ParagraphStyle"]("Cover", parent=body, fontSize=24, alignment=h["TA_CENTER"])))
    elems.append(h["Spacer"](1, 0.3*h["cm"]))
    elems.append(P(f"ESG YILLIK RAPOR — {year}",
                    h["ParagraphStyle"]("Cover2", parent=body, fontSize=18, alignment=h["TA_CENTER"],
                    textColor=h["colors"].HexColor("#475569"))))
    elems.append(h["Spacer"](1, 0.5*h["cm"]))
    elems.append(P("GRI · CDP · TCFD Uyumlu",
                    h["ParagraphStyle"]("Cover3", parent=body, fontSize=11, alignment=h["TA_CENTER"],
                    textColor=h["colors"].HexColor("#94a3b8"))))
    elems.append(h["PageBreak"]())

    # GRI 102 — Genel Bilgiler
    elems.append(P("GRI 102 — Organizasyonel Profil", h1))
    elems.append(P(f"<b>Kurum:</b> {tname}", body))
    if tenant:
        if tenant.sector: elems.append(P(f"<b>Sektör:</b> {tenant.sector}", body))
        if tenant.employee_count: elems.append(P(f"<b>Çalışan sayısı:</b> {tenant.employee_count}", body))
        if tenant.vision: elems.append(P(f"<b>Vizyon:</b> {tenant.vision[:300]}", body))
    elems.append(h["Spacer"](1, 0.5*h["cm"]))

    # GRI 305 — Emisyonlar
    elems.append(P("GRI 305 — Emisyonlar", h1))
    e_metrics = EsgMetric.query.filter_by(tenant_id=tid, is_active=True).filter(
        EsgMetric.category == "E").all()
    if e_metrics:
        rows = [["Kod", "Metrik", "Scope", "Birim"]]
        for m in e_metrics:
            rows.append([m.code or "—", m.name or "—", m.scope or "—", m.unit or "—"])
        tbl = h["Table"](rows, colWidths=[2.5*h["cm"], 7*h["cm"], 3*h["cm"], 3*h["cm"]])
        tbl.setStyle(h["TableStyle"]([
            ("BACKGROUND", (0, 0), (-1, 0), h["colors"].HexColor("#16a34a")),
            ("TEXTCOLOR", (0, 0), (-1, 0), h["colors"].white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1),
             [h["colors"].white, h["colors"].HexColor("#f0fdf4")]),
            ("BOX", (0, 0), (-1, -1), 1, h["colors"].HexColor("#e2e8f0")),
            ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#e2e8f0")),
            ("PADDING", (0, 0), (-1, -1), 6),
        ]))
        elems.append(tbl)

        # En son yıl toplam
        by_year_scope = defaultdict(lambda: defaultdict(float))
        for m in e_metrics:
            for v in m.values:
                if v.value is not None and m.scope:
                    by_year_scope[v.year][m.scope] += v.value
        years_sorted = sorted(by_year_scope.keys())
        if years_sorted:
            latest = years_sorted[-1]
            elems.append(h["Spacer"](1, 0.5*h["cm"]))
            elems.append(P(f"<b>{latest} Yılı Toplam Emisyonlar:</b>", body))
            for scope in ("scope1", "scope2", "scope3"):
                v = by_year_scope[latest].get(scope, 0)
                if v > 0:
                    elems.append(P(f"  {scope.upper()}: {v:,.2f} tCO₂e", body))
    else:
        elems.append(P("Henüz çevre metriği tanımlanmamış.", body))

    elems.append(h["PageBreak"]())

    # GRI 403 — Sosyal
    elems.append(P("GRI 403 — Sosyal Göstergeler (LTIFR, vb.)", h1))
    s_metrics = EsgMetric.query.filter_by(tenant_id=tid, is_active=True, category="S").all()
    if s_metrics:
        for m in s_metrics:
            elems.append(P(f"<b>{m.name}</b>", body))
            elems.append(P(f"Kod: {m.code or '—'} · Birim: {m.unit or '—'}", body))
            if m.description: elems.append(P(m.description, body))
    else:
        elems.append(P("Henüz sosyal metriği tanımlanmamış.", body))

    elems.append(h["Spacer"](1, 1*h["cm"]))

    # GRI 102-22 — Yönetişim
    elems.append(P("GRI 102-22 — Yönetişim", h1))
    g_metrics = EsgMetric.query.filter_by(tenant_id=tid, is_active=True, category="G").all()
    if g_metrics:
        for m in g_metrics:
            elems.append(P(f"<b>{m.name}:</b> {m.description or '—'}", body))
    else:
        elems.append(P("Henüz yönetişim metriği tanımlanmamış.", body))

    elems.append(h["PageBreak"]())

    # TCFD — İklim Risk Açıklaması
    elems.append(P("TCFD — İklim Risk Açıklaması", h1))
    elems.append(P("Yönetişim, Strateji, Risk Yönetimi, Metrik ve Hedefler dört "
                   "boyutta iklim ile ilgili açıklamalar:", body))
    elems.append(P("<b>1. Yönetişim:</b> İklim risklerinin yönetim kurulu seviyesinde "
                   "değerlendirilmesi.", body))
    elems.append(P("<b>2. Strateji:</b> İklim ile ilgili kısa/orta/uzun vadeli stratejik "
                   "etki analizi.", body))
    elems.append(P("<b>3. Risk yönetimi:</b> Risk register'a iklim risklerinin entegrasyonu.", body))
    elems.append(P("<b>4. Metrik ve hedefler:</b> Scope 1/2/3 emisyon ölçümü ve Net Zero "
                   "hedefi (varsa).", body))

    # Kapanış
    elems.append(h["Spacer"](1, 1*h["cm"]))
    elems.append(P("Bu rapor Kokpitim platformu üzerinden, GRI/CDP/TCFD standart "
                   "yapısına uygun olarak otomatik üretilmiştir. Detay veri ve "
                   "yöntem için ek ekleri talep edebilirsiniz.",
                   h["ParagraphStyle"]("Foot", parent=body, fontSize=9,
                       textColor=h["colors"].HexColor("#64748b"))))

    doc.build(elems); buf.seek(0)
    filename = f"{tname.replace(' ', '_')}_ESG_Raporu_{year}.pdf"
    return send_file(buf, mimetype="application/pdf",
                     as_attachment=True, download_name=filename)


# ─── RK-14: Audit Çıktı Paketi PDF ─────────────────────────────────────────

@app_bp.route("/raporlar/audit-paketi")
@login_required
def raporlar_audit_paketi():
    return render_template("platform/raporlar/audit_paketi.html")


@app_bp.route("/raporlar/api/audit-paketi/generate", methods=["GET", "POST"])
@login_required
def raporlar_api_audit_paketi_generate():
    from flask import send_file
    from app.models.audit import AuditLog
    import io

    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    tenant = db.session.get(Tenant, tid)
    tname = tenant.name if tenant else "Kurum"
    year = _date.today().year

    h = _pdf_helpers()
    buf = io.BytesIO()
    doc = h["SimpleDocTemplate"](buf, pagesize=h["A4"],
        leftMargin=2*h["cm"], rightMargin=2*h["cm"], topMargin=2*h["cm"], bottomMargin=2*h["cm"],
        title=f"{tname} Audit Çıktı Paketi {year}")
    styles = h["styles"]
    h1 = h["ParagraphStyle"]("H1", parent=styles["Heading1"], fontSize=18,
        textColor=h["colors"].HexColor("#dc2626"), spaceAfter=14)
    body = h["ParagraphStyle"]("Body", parent=styles["BodyText"], fontSize=10.5,
        leading=14, textColor=h["colors"].HexColor("#0f172a"),
        spaceAfter=6, alignment=h["TA_JUSTIFY"])
    small = h["ParagraphStyle"]("Sm", parent=body, fontSize=9,
        textColor=h["colors"].HexColor("#64748b"))

    P = h["Paragraph"]; elems = []

    # Kapak
    elems.append(P(f"<b>{tname}</b>",
                    h["ParagraphStyle"]("C1", parent=body, fontSize=24, alignment=h["TA_CENTER"])))
    elems.append(P(f"AUDIT ÇIKTI PAKETİ — {year}",
                    h["ParagraphStyle"]("C2", parent=body, fontSize=16, alignment=h["TA_CENTER"],
                    textColor=h["colors"].HexColor("#dc2626"))))
    elems.append(h["Spacer"](1, 1*h["cm"]))
    elems.append(P("Bu doküman, üçüncü parti denetçilere kurum içi stratejik plan, "
                   "performans, risk ve kullanıcı aktivite verisinin özet sunumudur.",
                   h["ParagraphStyle"]("CI", parent=body, alignment=h["TA_CENTER"])))
    elems.append(h["PageBreak"]())

    # 1. Tenant özeti
    elems.append(P("1. Kurum Bilgileri", h1))
    info_rows = [["Kurum", tname]]
    if tenant:
        if tenant.short_name: info_rows.append(["Kısa Ad", tenant.short_name])
        if tenant.sector: info_rows.append(["Sektör", tenant.sector])
        if tenant.employee_count: info_rows.append(["Çalışan", str(tenant.employee_count)])
        info_rows.append(["Plan Year Aktif", "Evet" if tenant.plan_year_enabled else "Hayır"])
        info_rows.append(["Tenant ID", str(tid)])
    tbl = h["Table"](info_rows, colWidths=[5*h["cm"], 11*h["cm"]])
    tbl.setStyle(h["TableStyle"]([
        ("BACKGROUND", (0, 0), (0, -1), h["colors"].HexColor("#f8fafc")),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("BOX", (0, 0), (-1, -1), 0.5, h["colors"].HexColor("#cbd5e1")),
        ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#e2e8f0")),
        ("PADDING", (0, 0), (-1, -1), 8),
    ]))
    elems.append(tbl)
    elems.append(h["PageBreak"]())

    # 2. Stratejik plan özet
    elems.append(P("2. Stratejik Plan Yapısı", h1))
    strat = Strategy.query.filter_by(tenant_id=tid, is_active=True).count()
    sub = db.session.query(func.count(SubStrategy.id)).join(Strategy).filter(
        Strategy.tenant_id == tid, SubStrategy.is_active.is_(True)).scalar() or 0
    proc = Process.query.filter_by(tenant_id=tid, is_active=True).count()
    kpi = ProcessKpi.query.join(Process).filter(
        Process.tenant_id == tid, ProcessKpi.is_active.is_(True)).count()
    sp_rows = [["Ana Strateji", str(strat)], ["Alt Strateji", str(sub)],
               ["Süreç", str(proc)], ["Performans Göstergesi", str(kpi)]]
    tbl = h["Table"](sp_rows, colWidths=[8*h["cm"], 4*h["cm"]])
    tbl.setStyle(h["TableStyle"]([
        ("BACKGROUND", (0, 0), (0, -1), h["colors"].HexColor("#f8fafc")),
        ("FONTSIZE", (0, 0), (-1, -1), 11),
        ("BOX", (0, 0), (-1, -1), 0.5, h["colors"].HexColor("#cbd5e1")),
        ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#e2e8f0")),
        ("PADDING", (0, 0), (-1, -1), 8),
    ]))
    elems.append(tbl)
    elems.append(h["PageBreak"]())

    # 3. Audit log son 90 gün
    elems.append(P("3. Audit Log — Son 90 Gün", h1))
    last_90 = datetime.utcnow() - timedelta(days=90)
    audit_logs = AuditLog.query.filter(
        AuditLog.tenant_id == tid, AuditLog.created_at >= last_90,
    ).order_by(AuditLog.created_at.desc()).limit(50).all()

    # Action dağılımı
    action_count = defaultdict(int)
    for a in audit_logs:
        action_count[a.action] += 1
    if action_count:
        elems.append(P("<b>Aksiyon Dağılımı (top 50):</b>", body))
        for act, cnt in sorted(action_count.items(), key=lambda x: -x[1]):
            elems.append(P(f"  {act}: {cnt}", small))
        elems.append(h["Spacer"](1, 0.5*h["cm"]))

    # En son 20 audit log
    if audit_logs:
        elems.append(P("<b>Son 20 İşlem:</b>", body))
        rows = [["Tarih", "Kullanıcı", "Aksiyon", "Kaynak"]]
        for a in audit_logs[:20]:
            rows.append([
                a.created_at.strftime("%d.%m %H:%M") if a.created_at else "—",
                (a.username or "—")[:25],
                a.action or "—",
                (a.resource_type or "—")[:20],
            ])
        tbl = h["Table"](rows, colWidths=[3*h["cm"], 4.5*h["cm"], 4*h["cm"], 4*h["cm"]])
        tbl.setStyle(h["TableStyle"]([
            ("BACKGROUND", (0, 0), (-1, 0), h["colors"].HexColor("#dc2626")),
            ("TEXTCOLOR", (0, 0), (-1, 0), h["colors"].white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1),
             [h["colors"].white, h["colors"].HexColor("#fef2f2")]),
            ("PADDING", (0, 0), (-1, -1), 5),
            ("BOX", (0, 0), (-1, -1), 0.5, h["colors"].HexColor("#fecaca")),
            ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#fecaca")),
        ]))
        elems.append(tbl)
    else:
        elems.append(P("Son 90 günde audit log kaydı yok.", small))
    elems.append(h["PageBreak"]())

    # 4. Kullanıcı + 2FA
    elems.append(P("4. Kullanıcı Güvenlik Özeti", h1))
    total_u = User.query.filter_by(tenant_id=tid, is_active=True).count()
    totp_u = User.query.filter_by(tenant_id=tid, is_active=True, totp_enabled=True).count()
    elems.append(P(f"<b>Toplam aktif kullanıcı:</b> {total_u}", body))
    elems.append(P(f"<b>2FA etkin:</b> {totp_u} (%{round(totp_u/max(total_u,1)*100,1)})", body))
    elems.append(P(f"<b>2FA pasif:</b> {total_u - totp_u}", body))

    # Footer
    elems.append(h["Spacer"](1, 2*h["cm"]))
    elems.append(P(f"Bu rapor {_date.today().strftime('%d.%m.%Y')} tarihinde "
                   f"Kokpitim platformundan otomatik üretilmiştir.",
                   h["ParagraphStyle"]("Foot", parent=small, alignment=h["TA_CENTER"])))

    doc.build(elems); buf.seek(0)
    filename = f"{tname.replace(' ', '_')}_Audit_Paketi_{year}.pdf"
    return send_file(buf, mimetype="application/pdf",
                     as_attachment=True, download_name=filename)


# ─── HR-01: Bireysel Karne PDF Batch (ZIP) ─────────────────────────────────

@app_bp.route("/raporlar/bireysel-karne-batch")
@login_required
def raporlar_bireysel_karne_batch():
    return render_template("platform/raporlar/bireysel_karne_batch.html")


@app_bp.route("/raporlar/api/bireysel-karne-batch/preview")
@login_required
def raporlar_api_bireysel_karne_batch_preview():
    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    # PG'si olan kullanıcı sayısı
    users_with_pg = db.session.query(func.count(func.distinct(IndividualPerformanceIndicator.user_id))).join(
        User, User.id == IndividualPerformanceIndicator.user_id
    ).filter(
        User.tenant_id == tid, User.is_active.is_(True),
        IndividualPerformanceIndicator.is_active.is_(True),
    ).scalar() or 0
    total = User.query.filter_by(tenant_id=tid, is_active=True).count()
    return jsonify({"success": True, "preview": {
        "total_users": total, "users_with_pg": users_with_pg,
        "format": "ZIP (PDF her kullanıcı için)",
        "estimated_size_mb": round(users_with_pg * 0.05, 2),
    }})


@app_bp.route("/raporlar/api/bireysel-karne-batch/generate", methods=["GET", "POST"])
@login_required
def raporlar_api_bireysel_karne_batch_generate():
    from flask import send_file
    import io, zipfile

    tid = _tid_or_none()
    if not tid:
        return jsonify({"success": False, "message": "Tenant yok"}), 400
    tenant = db.session.get(Tenant, tid)
    tname = tenant.name if tenant else "Kurum"

    # PG'si olan kullanıcılar
    user_ids_with_pg = [r[0] for r in db.session.query(
        func.distinct(IndividualPerformanceIndicator.user_id)
    ).join(User, User.id == IndividualPerformanceIndicator.user_id).filter(
        User.tenant_id == tid, User.is_active.is_(True),
        IndividualPerformanceIndicator.is_active.is_(True),
    ).all()]

    if not user_ids_with_pg:
        return jsonify({"success": False, "message": "Bireysel PG'si olan kullanıcı yok."}), 400

    h = _pdf_helpers()
    P = h["Paragraph"]
    styles = h["styles"]
    h1 = h["ParagraphStyle"]("H1", parent=styles["Heading1"], fontSize=18,
        textColor=h["colors"].HexColor("#4f46e5"), spaceAfter=12)
    body = h["ParagraphStyle"]("Body", parent=styles["BodyText"], fontSize=10.5,
        leading=14, textColor=h["colors"].HexColor("#0f172a"), spaceAfter=4)

    # ZIP container
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for uid in user_ids_with_pg[:100]:  # max 100 user (güvenlik)
            u = db.session.get(User, uid)
            if not u:
                continue
            uname = f"{u.first_name or ''} {u.last_name or ''}".strip() or u.email
            pgs = IndividualPerformanceIndicator.query.filter_by(user_id=uid, is_active=True).all()
            aligned = sum(1 for p in pgs if p.source_process_id or p.source_process_kpi_id)
            alignment_pct = round(aligned / max(len(pgs), 1) * 100, 1)

            # Tek PDF üret
            pdf_buf = io.BytesIO()
            doc = h["SimpleDocTemplate"](pdf_buf, pagesize=h["A4"],
                leftMargin=2*h["cm"], rightMargin=2*h["cm"], topMargin=2*h["cm"], bottomMargin=2*h["cm"],
                title=f"{uname} Bireysel Karne")
            elems = []
            elems.append(P(f"<b>{tname}</b>",
                h["ParagraphStyle"]("Top", parent=body, fontSize=10,
                    textColor=h["colors"].HexColor("#64748b"))))
            elems.append(P(f"BİREYSEL PERFORMANS KARNESİ", h1))
            elems.append(P(f"<b>{uname}</b>", body))
            elems.append(P(f"E-posta: {u.email}", body))
            if u.department: elems.append(P(f"Departman: {u.department}", body))
            elems.append(h["Spacer"](1, 0.5*h["cm"]))
            elems.append(P(f"<b>Özet:</b> {len(pgs)} bireysel PG · "
                f"{aligned} kurum stratejisine bağlı (%{alignment_pct} hizalama)", body))
            elems.append(h["Spacer"](1, 0.5*h["cm"]))

            if pgs:
                rows = [["Kod", "PG Adı", "Hedef", "Birim", "Tip"]]
                for p in pgs[:30]:
                    rows.append([
                        (p.code or "—")[:15],
                        (p.name or "")[:40],
                        (p.target_value or "—")[:10],
                        (p.unit or "—")[:10],
                        "Kurumdan" if p.source_process_id else "Bireysel",
                    ])
                tbl = h["Table"](rows, colWidths=[2*h["cm"], 7*h["cm"], 2.5*h["cm"], 2*h["cm"], 3*h["cm"]])
                tbl.setStyle(h["TableStyle"]([
                    ("BACKGROUND", (0, 0), (-1, 0), h["colors"].HexColor("#4f46e5")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), h["colors"].white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1),
                     [h["colors"].white, h["colors"].HexColor("#f8fafc")]),
                    ("PADDING", (0, 0), (-1, -1), 5),
                    ("BOX", (0, 0), (-1, -1), 0.5, h["colors"].HexColor("#e2e8f0")),
                    ("INNERGRID", (0, 0), (-1, -1), 0.3, h["colors"].HexColor("#e2e8f0")),
                ]))
                elems.append(tbl)
            elems.append(h["Spacer"](1, 1.5*h["cm"]))
            elems.append(P("Bu karne Kokpitim platformu üzerinden otomatik üretilmiştir.",
                h["ParagraphStyle"]("F", parent=body, fontSize=9,
                    textColor=h["colors"].HexColor("#94a3b8"))))

            doc.build(elems)
            pdf_buf.seek(0)
            safe_name = uname.replace(" ", "_").replace("/", "_")[:40]
            zf.writestr(f"{safe_name}_karne.pdf", pdf_buf.read())

    zip_buf.seek(0)
    filename = f"{tname.replace(' ', '_')}_bireysel_karne_batch.zip"
    return send_file(zip_buf, mimetype="application/zip",
                     as_attachment=True, download_name=filename)
