from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# Renk paleti
DARK_BG     = RGBColor(0x0F, 0x17, 0x2A)   # koyu lacivert arka plan
ACCENT      = RGBColor(0x38, 0xBD, 0xF8)   # cyan/sky mavi
ACCENT2     = RGBColor(0x34, 0xD3, 0x99)   # yeşil
ACCENT3     = RGBColor(0xFB, 0xBF, 0x24)   # sarı/turuncu
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCB, 0xD5, 0xE1)
CARD_BG     = RGBColor(0x1E, 0x29, 0x3B)
RED_ACCENT  = RGBColor(0xF8, 0x71, 0x71)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # boş layout

def add_bg(slide, color=DARK_BG):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

def box(slide, x, y, w, h, fill_color, radius=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def line(slide, x, y, w, color=ACCENT, thickness=2):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ─────────────────────────────────────────────
# SLIDE 1 — KAPAK
# ─────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
add_bg(s1)

# Sol dekoratif dikey çubuk
box(s1, 0, 0, 0.08, 7.5, ACCENT)

# Sağ alt dekoratif blok
box(s1, 10.5, 5.8, 2.83, 1.7, CARD_BG)
box(s1, 10.5, 5.8, 0.06, 1.7, ACCENT2)

# Ana başlık
txt(s1, "KOKPİTİM", 1.0, 1.8, 8, 1.4, size=64, bold=True, color=WHITE)
line(s1, 1.0, 3.3, 5.5, color=ACCENT)

# Alt başlık
txt(s1, "Kurumsal Performans Yönetim Platformu", 1.0, 3.5, 9, 0.8, size=24, color=ACCENT)

# Açıklama
txt(s1, "Strateji → Süreç → KPI → Bireysel Performans\nTek Entegre Platformda",
    1.0, 4.3, 9, 1.0, size=16, color=LIGHT_GRAY)

# Versiyon/tarih
txt(s1, "Versiyon 2.0  |  2026", 1.0, 6.5, 5, 0.5, size=13, color=LIGHT_GRAY)

# Durum badge
box(s1, 10.6, 5.95, 2.5, 0.55, ACCENT2)
txt(s1, "✓  %98 Tamamlandı", 10.7, 5.95, 2.4, 0.55, size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 2 — PROJE TANIMI
# ─────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
add_bg(s2)
box(s2, 0, 0, 0.08, 7.5, ACCENT)

txt(s2, "Proje Tanımı", 0.5, 0.3, 9, 0.8, size=32, bold=True, color=WHITE)
line(s2, 0.5, 1.15, 12.3, color=ACCENT)

# 3 kart
cards = [
    (ACCENT,  "HEDEF",      "Türkiye pazarına özel çok kiracılı\n(multi-tenant) SaaS platform"),
    (ACCENT2, "FARK",       "Strateji→Süreç→KPI→Bireysel zinciri\nentegre — rakipler ayrı modülle çözüyor"),
    (ACCENT3, "DURUM",      "900 / 920 saat efor\n%98 tamamlandı — Production Ready"),
]
cx = 0.5
for color, title, body in cards:
    box(s2, cx, 1.4, 3.9, 2.6, CARD_BG)
    box(s2, cx, 1.4, 3.9, 0.06, color)
    txt(s2, title, cx+0.15, 1.5, 3.6, 0.5, size=15, bold=True, color=color)
    txt(s2, body,  cx+0.15, 2.05, 3.6, 1.8, size=13, color=LIGHT_GRAY)
    cx += 4.15

# Alt satır — kullanıcı rolleri
txt(s2, "Kullanıcı Rolleri:", 0.5, 4.2, 3, 0.4, size=14, bold=True, color=ACCENT)
roles = ["Admin", "Tenant Admin", "Executive Manager", "User"]
rx = 0.5
for r in roles:
    box(s2, rx, 4.65, 2.8, 0.5, CARD_BG)
    box(s2, rx, 4.65, 0.06, 0.5, ACCENT)
    txt(s2, r, rx+0.15, 4.65, 2.6, 0.5, size=13, color=WHITE)
    rx += 3.05

# Port notu
txt(s2, "Çalışma Portu: 5001  |  Multi-tenant mimari  |  Python 3.11 / Flask",
    0.5, 6.7, 12, 0.4, size=12, color=LIGHT_GRAY)


# ─────────────────────────────────────────────
# SLIDE 3 — TEKNOLOJİ STACK
# ─────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
add_bg(s3)
box(s3, 0, 0, 0.08, 7.5, ACCENT2)

txt(s3, "Teknoloji Stack", 0.5, 0.3, 9, 0.8, size=32, bold=True, color=WHITE)
line(s3, 0.5, 1.15, 12.3, color=ACCENT2)

cols = [
    ("Backend", ACCENT, [
        "Python 3.11 + Flask 3.0",
        "SQLAlchemy 2.0 (ORM)",
        "Flask-Migrate / Alembic",
        "Flask-Login (Auth)",
        "Flask-SocketIO 5.3 + Eventlet",
        "Celery 5 + Redis",
        "Marshmallow 3 (Serialization)",
        "PyJWT 2.8 (JWT)",
    ]),
    ("Frontend", ACCENT2, [
        "Vue.js 3",
        "Alpine.js",
        "Chart.js 4.4",
        "SweetAlert2 11",
        "Tailwind CSS (CDN)",
        "Service Worker (PWA)",
        "Web Push API",
        "IndexedDB",
    ]),
    ("Altyapı & DevOps", ACCENT3, [
        "PostgreSQL 15+ (üretim)",
        "SQLite (geliştirme)",
        "Redis 7 (cache/broker)",
        "Sentry (hata takip)",
        "pytest + coverage",
        "Flask-Limiter",
        "Flask-Caching",
        "scikit-learn (ML)",
    ]),
]

cx = 0.5
for (title, color, items) in cols:
    box(s3, cx, 1.4, 4.0, 5.6, CARD_BG)
    box(s3, cx, 1.4, 4.0, 0.06, color)
    txt(s3, title, cx+0.15, 1.5, 3.7, 0.5, size=15, bold=True, color=color)
    iy = 2.1
    for item in items:
        txt(s3, f"• {item}", cx+0.2, iy, 3.6, 0.38, size=12, color=LIGHT_GRAY)
        iy += 0.38
    cx += 4.25


# ─────────────────────────────────────────────
# SLIDE 4 — MİMARİ YAPI
# ─────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
add_bg(s4)
box(s4, 0, 0, 0.08, 7.5, ACCENT3)

txt(s4, "Mimari Yapı", 0.5, 0.3, 9, 0.8, size=32, bold=True, color=WHITE)
line(s4, 0.5, 1.15, 12.3, color=ACCENT3)

# Blueprint akışı
txt(s4, "Uygulama Akışı", 0.5, 1.35, 5, 0.4, size=14, bold=True, color=ACCENT3)
box(s4, 0.5, 1.8, 12.3, 0.7, CARD_BG)
txt(s4, "app.py  →  create_app()  →  config.py  →  extensions init  →  blueprints register  →  app.run(5001)",
    0.65, 1.85, 12.0, 0.6, size=13, color=ACCENT, bold=True)

# Modül tablosu
txt(s4, "Micro Modüller  (~113 route)", 0.5, 2.65, 8, 0.4, size=14, bold=True, color=ACCENT3)

modules = [
    ("admin",          "22", "Kullanıcı, tenant, paket yönetimi"),
    ("surec",          "21", "Süreç yönetimi, KPI takibi"),
    ("api",            "15", "REST API (Swagger UI)"),
    ("sp",             "12", "Stratejik Planlama, SWOT"),
    ("bireysel",       "11", "Bireysel performans, karne"),
    ("proje",          "10", "Proje & görev yönetimi"),
    ("analiz",          "7", "Kurum analiz merkezi"),
    ("shared/auth",     "5", "Profil, login"),
    ("kurum",           "8", "Kurum yönetimi, strateji"),
    ("k_radar",         "—", "K-Radar performans analizi (YENİ)"),
]

col1_x, col2_x, col3_x = 0.5, 2.3, 4.5
ry = 3.1
header_color = CARD_BG
box(s4, 0.5, 3.05, 12.3, 0.4, ACCENT3)
txt(s4, "Modül", 0.55, 3.08, 1.7, 0.35, size=12, bold=True, color=DARK_BG)
txt(s4, "Route", 2.3,  3.08, 0.8, 0.35, size=12, bold=True, color=DARK_BG)
txt(s4, "Açıklama", 3.3, 3.08, 8, 0.35, size=12, bold=True, color=DARK_BG)

for i, (mod, cnt, desc) in enumerate(modules):
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x16, 0x21, 0x32)
    row_y = 3.45 + i * 0.37
    box(s4, 0.5, row_y, 12.3, 0.37, bg)
    txt(s4, mod,  0.6,  row_y, 1.7, 0.37, size=11, color=ACCENT)
    txt(s4, cnt,  2.3,  row_y, 0.8, 0.37, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s4, desc, 3.3,  row_y, 9.0, 0.37, size=11, color=LIGHT_GRAY)


# ─────────────────────────────────────────────
# SLIDE 5 — PERFORMANS İYİLEŞTİRMELERİ
# ─────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
add_bg(s5)
box(s5, 0, 0, 0.08, 7.5, RED_ACCENT)

txt(s5, "Performans İyileştirmeleri", 0.5, 0.3, 10, 0.8, size=32, bold=True, color=WHITE)
line(s5, 0.5, 1.15, 12.3, color=RED_ACCENT)

metrics = [
    ("Güvenlik Skoru",    "C",       "A+",    "+3 Seviye",  ACCENT3),
    ("Yanıt Süresi",      "~500ms",  "~100ms","%80 ↓",      ACCENT2),
    ("DB Sorgu Sayısı",   "301",     "4",     "%98.7 ↓",    ACCENT),
    ("Mobil Skoru",       "40/100",  "95/100","+137%",       ACCENT3),
    ("PWA Skoru",         "0/100",   "90/100","+90 puan",   ACCENT2),
    ("Test Coverage",     "0%",      "%50+",  "+50 puan",   ACCENT),
    ("API Endpoint",      "0",       "25+",   "—",          ACCENT3),
    ("ML Özellikleri",    "❌",      "✅",    "—",          RED_ACCENT),
]

# Tablo başlığı
box(s5, 0.5, 1.35, 12.3, 0.45, ACCENT)
for lbl, x in [("Metrik",0.6),("Öncesi",4.8),("Sonrası",7.2),("İyileşme",10.2)]:
    txt(s5, lbl, x, 1.38, 2.5, 0.4, size=13, bold=True, color=DARK_BG)

for i, (metric, before, after, delta, color) in enumerate(metrics):
    row_y = 1.8 + i * 0.56
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x16, 0x21, 0x32)
    box(s5, 0.5, row_y, 12.3, 0.56, bg)
    box(s5, 0.5, row_y, 0.08, 0.56, color)
    txt(s5, metric, 0.65, row_y, 4.1, 0.56, size=13, color=WHITE)
    txt(s5, before, 4.8,  row_y, 2.3, 0.56, size=13, color=RED_ACCENT)
    txt(s5, after,  7.2,  row_y, 2.8, 0.56, size=13, color=ACCENT2, bold=True)
    txt(s5, delta,  10.2, row_y, 2.5, 0.56, size=13, color=ACCENT3, bold=True)


# ─────────────────────────────────────────────
# SLIDE 6 — TAMAMLANAN ÖZELLİKLER
# ─────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
add_bg(s6)
box(s6, 0, 0, 0.08, 7.5, ACCENT2)

txt(s6, "Tamamlanan Özellikler", 0.5, 0.3, 10, 0.8, size=32, bold=True, color=WHITE)
line(s6, 0.5, 1.15, 12.3, color=ACCENT2)

feature_groups = [
    ("Güvenlik & Stabilite", ACCENT, [
        "Security headers (HSTS, CSP)",
        "Rate limiting (Flask-Limiter)",
        "Error tracking (Sentry)",
        "Input validation (Marshmallow)",
        "Audit logging",
        "Unit test %50+ coverage",
    ]),
    ("Analitik & Raporlama", ACCENT2, [
        "Trend analizi",
        "Sağlık skoru hesaplama",
        "Karşılaştırma analizi",
        "Anomali tespiti",
        "Tahminleme (ML tabanlı)",
        "Dashboard builder + Excel export",
    ]),
    ("API & Entegrasyon", ACCENT3, [
        "RESTful API v1 (25+ endpoint)",
        "Swagger / OpenAPI dokümantasyonu",
        "OAuth2 kimlik doğrulama",
        "JWT token yönetimi",
        "API key sistemi",
        "Webhook sistemi",
    ]),
    ("Mobil & PWA", RED_ACCENT, [
        "Progressive Web App (PWA)",
        "Service Worker + Offline mod",
        "Web Push bildirimleri",
        "Mobil-first responsive UI",
        "Background sync",
        "Install prompt",
    ]),
]

cx = 0.5
for (title, color, items) in feature_groups:
    box(s6, cx, 1.4, 3.0, 5.6, CARD_BG)
    box(s6, cx, 1.4, 3.0, 0.06, color)
    txt(s6, title, cx+0.12, 1.5, 2.8, 0.45, size=13, bold=True, color=color)
    iy = 2.0
    for item in items:
        txt(s6, f"✓  {item}", cx+0.12, iy, 2.8, 0.4, size=11, color=LIGHT_GRAY)
        iy += 0.4
    cx += 3.2


# ─────────────────────────────────────────────
# SLIDE 7 — K-RADAR YENİ ÖZELLİK
# ─────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
add_bg(s7)
box(s7, 0, 0, 0.08, 7.5, ACCENT3)

txt(s7, "K-Radar  —  Yeni Analiz Motoru", 0.5, 0.3, 10, 0.8, size=32, bold=True, color=WHITE)
line(s7, 0.5, 1.15, 12.3, color=ACCENT3)

# Sol — açıklama
box(s7, 0.5, 1.4, 5.8, 5.6, CARD_BG)
box(s7, 0.5, 1.4, 5.8, 0.06, ACCENT3)
txt(s7, "K-Radar Nedir?", 0.65, 1.5, 5.5, 0.5, size=16, bold=True, color=ACCENT3)
txt(s7, (
    "K-Radar, Kokpitim'in kurumsal performans\n"
    "değerlendirme motoru olarak tasarlanmış,\n"
    "çok boyutlu analiz altyapısıdır.\n\n"
    "Domain tabanlı ağırlıklı skorlama ile\n"
    "kurum sağlığını tek ekranda gösterir.\n\n"
    "ML tabanlı öneri motoru ve K-Vektör\n"
    "entegrasyonu ile karar destek sağlar."
), 0.65, 2.1, 5.5, 4.5, size=13, color=LIGHT_GRAY)

# Sağ — özellikler
box(s7, 6.6, 1.4, 6.2, 5.6, CARD_BG)
box(s7, 6.6, 1.4, 6.2, 0.06, ACCENT2)
txt(s7, "Temel Yetenekler", 6.75, 1.5, 5.9, 0.5, size=16, bold=True, color=ACCENT2)

kradar_features = [
    (ACCENT,  "Domain Tabanlı Skor",   "Ağırlıklı puan hesabı, kurum odaklı"),
    (ACCENT2, "K-Vektör Entegrasyonu", "Çok boyutlu kurumsal vektör analizi"),
    (ACCENT3, "Öneri Aksiyonları",     "ML bazlı aksiyon önerileri"),
    (RED_ACCENT,"Radar Görsel",        "Radar chart + ısı haritası"),
    (ACCENT,  "Scheduler Servisi",     "Otomatik periyodik hesaplama"),
    (ACCENT2, "PDF / Excel Export",    "Raporlama ve dışa aktarım"),
]

fy = 2.1
for (color, title, desc) in kradar_features:
    box(s7, 6.75, fy, 5.9, 0.68, RGBColor(0x16, 0x21, 0x32))
    box(s7, 6.75, fy, 0.07, 0.68, color)
    txt(s7, title, 6.9, fy+0.03, 2.4, 0.35, size=12, bold=True, color=color)
    txt(s7, desc,  6.9, fy+0.35, 5.6, 0.32, size=11, color=LIGHT_GRAY)
    fy += 0.75


# ─────────────────────────────────────────────
# SLIDE 8 — SPRINT ÖZETI / YOL HARİTASI
# ─────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
add_bg(s8)
box(s8, 0, 0, 0.08, 7.5, ACCENT)

txt(s8, "Sprint Özeti & Yol Haritası", 0.5, 0.3, 10, 0.8, size=32, bold=True, color=WHITE)
line(s8, 0.5, 1.15, 12.3, color=ACCENT)

sprints = [
    ("Sprint 0",     "Quick Wins",               "20 saat",  "✅"),
    ("Sprint 1-2",   "Frontend Modernizasyon",    "120 saat", "✅"),
    ("Sprint 3-4",   "Performans Optimizasyonu",  "80 saat",  "✅"),
    ("Sprint 5-6",   "Güvenlik & Stabilite",      "100 saat", "✅"),
    ("Sprint 7-9",   "Real-Time & Bildirimler",   "150 saat", "✅"),
    ("Sprint 10-12", "Analitik & Raporlama",      "180 saat", "✅"),
    ("Sprint 13-15", "API & Entegrasyonlar",      "160 saat", "✅"),
    ("Sprint 16-18", "AI & Otomasyon",            "200 saat", "✅"),
    ("Sprint 19-21", "Mobil & PWA",               "180 saat", "✅"),
]

box(s8, 0.5, 1.35, 9.3, 0.42, ACCENT)
for lbl, x in [("Sprint",0.6),("Konu",2.3),("Efor",7.0)]:
    txt(s8, lbl, x, 1.38, 2.5, 0.38, size=12, bold=True, color=DARK_BG)

for i, (sprint, topic, effort, status) in enumerate(sprints):
    row_y = 1.77 + i * 0.5
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x16, 0x21, 0x32)
    box(s8, 0.5, row_y, 9.3, 0.5, bg)
    txt(s8, sprint,  0.6, row_y, 1.6, 0.5, size=12, color=ACCENT,  bold=True)
    txt(s8, topic,   2.3, row_y, 4.6, 0.5, size=12, color=WHITE)
    txt(s8, effort,  7.0, row_y, 1.5, 0.5, size=12, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)
    txt(s8, status,  8.6, row_y, 1.0, 0.5, size=13, color=ACCENT2,  align=PP_ALIGN.CENTER)

# Sağ — kalan görevler
box(s8, 10.1, 1.35, 3.1, 5.8, CARD_BG)
box(s8, 10.1, 1.35, 3.1, 0.06, ACCENT3)
txt(s8, "Kalan Görevler", 10.2, 1.45, 2.9, 0.4, size=13, bold=True, color=ACCENT3)
remaining = [
    "React Native App (opsiyonel)",
    "HTTPS sertifikası (prod)",
    "Redis Cluster kurulumu",
    "Load balancer yapılandırma",
    "CI/CD pipeline",
    "Monitoring & alerting",
    "Backup stratejisi",
    "PDF export gerçek impl.",
]
ry2 = 1.95
for r in remaining:
    txt(s8, f"◦  {r}", 10.2, ry2, 2.9, 0.42, size=11, color=LIGHT_GRAY)
    ry2 += 0.42

txt(s8, "Toplam Efor:  900 / 920 saat", 0.5, 6.65, 9.3, 0.45,
    size=13, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────
# SLIDE 9 — ÖZET / KAPANIŞ
# ─────────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
add_bg(s9)
box(s9, 0, 0, 0.08, 7.5, ACCENT2)

# Büyük yeşil daire (dekoratif)
circ = s9.shapes.add_shape(9, Inches(8.8), Inches(1.5), Inches(4.2), Inches(4.2))
circ.fill.solid()
circ.fill.fore_color.rgb = RGBColor(0x1A, 0x2E, 0x26)
circ.line.color.rgb = ACCENT2
circ.line.width = util.Pt(2)

txt(s9, "Özet", 0.5, 0.3, 5, 0.8, size=32, bold=True, color=WHITE)
line(s9, 0.5, 1.15, 7.8, color=ACCENT2)

summaries = [
    (ACCENT,  "Platform",     "Çok kiracılı SaaS, entegre strateji-süreç-KPI-bireysel zinciri"),
    (ACCENT2, "Teknik",       "Python/Flask + Vue.js + PostgreSQL + Redis + Celery"),
    (ACCENT3, "Performans",   "%80 yanıt süresi iyileşmesi, %98.7 daha az DB sorgusu"),
    (RED_ACCENT,"Güvenlik",   "Security headers, rate limiting, audit logging, Sentry"),
    (ACCENT,  "AI & ML",      "Anomali tespiti, tahminleme, akıllı öneri motoru"),
    (ACCENT2, "Yeni Özellik", "K-Radar: çok boyutlu kurumsal performans analiz motoru"),
]

sy = 1.35
for (color, title, desc) in summaries:
    box(s9, 0.5, sy, 7.9, 0.75, CARD_BG)
    box(s9, 0.5, sy, 0.07, 0.75, color)
    txt(s9, title, 0.65, sy+0.05, 1.6, 0.35, size=12, bold=True, color=color)
    txt(s9, desc,  2.35, sy+0.05, 6.0, 0.65, size=12, color=LIGHT_GRAY)
    sy += 0.82

# Dairenin içine büyük %98
txt(s9, "%98", 9.1, 2.6, 3.5, 1.4, size=64, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
txt(s9, "Tamamlandı", 9.1, 3.95, 3.5, 0.6, size=18, color=WHITE, align=PP_ALIGN.CENTER)
txt(s9, "Production Ready", 9.1, 4.55, 3.5, 0.5, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Alt not
box(s9, 0.5, 6.8, 12.3, 0.5, CARD_BG)
txt(s9, "Kokpitim  v2.0  |  2026  |  Flask + Vue.js + PostgreSQL + Redis  |  Production Ready",
    0.5, 6.83, 12.3, 0.44, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ─── Kaydet ───
out = "c:/kokpitim/Kokpitim_Sunum.pptx"
prs.save(out)
print(f"Sunum oluşturuldu: {out}")
